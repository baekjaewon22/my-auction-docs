from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

BASE = Path("outputs/019ed30d-da34-7c91-b393-4f6b0561de79/presentations/employee-request-purpose-guide")
OUT = BASE / "output"
PREVIEW = BASE / "preview_pdf"
OUT.mkdir(parents=True, exist_ok=True)
PREVIEW.mkdir(parents=True, exist_ok=True)

PPTX = OUT / "사원용_신청_목적_작성_가이드.pptx"
PDF = OUT / "사원용_신청_목적_작성_가이드.pdf"

W, H = 1280, 720
SLIDE_W, SLIDE_H = Inches(13.333333), Inches(7.5)
SCALE_X, SCALE_Y = 13.333333 / W, 7.5 / H

C = {
    "bg": (247, 244, 236),
    "ink": (31, 41, 51),
    "muted": (104, 113, 125),
    "line": (216, 208, 194),
    "blue": (26, 115, 232),
    "orange": (245, 124, 0),
    "green": (24, 128, 56),
    "red": (217, 48, 37),
    "dark": (24, 32, 42),
    "white": (255, 255, 255),
    "pale_green": (244, 251, 246),
    "pale_red": (255, 247, 245),
}

FONT_REG = r"C:\Windows\Fonts\malgun.ttf"
FONT_BOLD = r"C:\Windows\Fonts\malgunbd.ttf"


def rgb(c):
    return RGBColor(*c)


def px(v, axis="x"):
    return Inches(v * (SCALE_X if axis == "x" else SCALE_Y))


def add_textbox(slide, text, x, y, w, h, size=20, color="ink", bold=False, font="맑은 고딕", align="left"):
    box = slide.shapes.add_textbox(px(x), px(y, "y"), px(w), px(h, "y"))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(C[color] if isinstance(color, str) else color)
    return box


def add_rect(slide, x, y, w, h, fill="white", line="line", radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        px(x), px(y, "y"), px(w), px(h, "y")
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(C[fill] if isinstance(fill, str) else fill)
    shape.line.color.rgb = rgb(C[line] if isinstance(line, str) else line)
    shape.line.width = Pt(1)
    return shape


def add_line(slide, x, y, w, color="line", width=1):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y, "y"), px(w), px(width, "y"))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(C[color])
    line.line.fill.background()
    return line


def base_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(C["bg"])
    add_textbox(slide, "사원용 신청 목적 작성 가이드", 54, 34, 420, 22, 10, "muted", True)
    add_textbox(slide, f"{page:02d}", 1184, 34, 42, 22, 10, "muted", True, align="right")
    add_line(slide, 54, 66, 1172)
    return slide


def kicker(slide, label):
    add_rect(slide, 54, 96, 28, 3, "orange", "orange", False)
    add_textbox(slide, label, 96, 90, 360, 20, 10, "muted", True)


def title(slide, txt, y=126, size=36):
    add_textbox(slide, txt, 54, y, 800, 104, size, "ink", True, "Georgia")


def pill(slide, txt, x, y, w, color):
    add_rect(slide, x, y, w, 34, "white", color)
    add_textbox(slide, txt, x + 14, y + 8, w - 28, 16, 10.5, color, True, align="center")


def bullet(slide, txt, x, y, w, color="ink"):
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y + 9, "y"), px(6), px(6, "y"))
    oval.fill.solid()
    oval.fill.fore_color.rgb = rgb(C["orange"])
    oval.line.fill.background()
    add_textbox(slide, txt, x + 18, y, w, 40, 15, color)


def draw_wrapped(draw, xy, text, font, fill, max_width, line_gap=8):
    x, y = xy
    lines = []
    for para in text.split("\n"):
        words = para.split(" ")
        line = ""
        for word in words:
            test = (line + " " + word).strip()
            if draw.textbbox((0, 0), test, font=font)[2] <= max_width or not line:
                line = test
            else:
                lines.append(line)
                line = word
        lines.append(line)
    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap


def image_base(page):
    img = Image.new("RGB", (W, H), C["bg"])
    d = ImageDraw.Draw(img)
    f10 = ImageFont.truetype(FONT_BOLD, 18)
    d.text((54, 34), "사원용 신청 목적 작성 가이드", font=f10, fill=C["muted"])
    d.text((1184, 34), f"{page:02d}", font=f10, fill=C["muted"])
    d.rectangle((54, 66, 1226, 67), fill=C["line"])
    return img, d


def font(size, bold=False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REG, size)


slides = []
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# 1
s = base_slide(prs, 1)
add_textbox(s, "TEAM GUIDE", 54, 108, 240, 22, 11, "orange", True)
add_textbox(s, "사원용 신청 목적은\n결재자가 바로 판단할 수 있게\n‘왜 필요한지’를 쓰는 칸입니다.", 54, 166, 760, 190, 36, "ink", True, "Georgia")
add_rect(s, 870, 132, 300, 382, "dark", "dark")
add_textbox(s, "목적 문장은\n짧을수록 좋지만\n근거는 빠지면 안 됩니다.", 908, 180, 224, 120, 25, "white", True, "Georgia")
add_textbox(s, "이 가이드는 팀원이 신청서를 작성할 때 목적란을 어떻게 채워야 하는지 설명합니다.", 908, 330, 220, 90, 14, (201, 210, 220))
pill(s, "대상: 팀원", 54, 520, 120, "blue")
pill(s, "용도: 작성 가이드", 188, 520, 168, "green")
pill(s, "핵심: 판단 가능성", 370, 520, 180, "orange")
add_textbox(s, "배포용 PDF와 편집 가능한 PPT 원본으로 함께 제공", 54, 618, 520, 26, 11, "muted")
img, d = image_base(1)
d.text((54, 108), "TEAM GUIDE", font=font(20, True), fill=C["orange"])
draw_wrapped(d, (54, 166), "사원용 신청 목적은\n결재자가 바로 판단할 수 있게\n‘왜 필요한지’를 쓰는 칸입니다.", font(47, True), C["ink"], 760, 8)
d.rounded_rectangle((870, 132, 1170, 514), radius=18, fill=C["dark"])
draw_wrapped(d, (908, 180), "목적 문장은\n짧을수록 좋지만\n근거는 빠지면 안 됩니다.", font(32, True), C["white"], 230)
draw_wrapped(d, (908, 330), "이 가이드는 팀원이 신청서를 작성할 때 목적란을 어떻게 채워야 하는지 설명합니다.", font(20), (201, 210, 220), 220)
for txt, x, w, col in [("대상: 팀원", 54, 120, "blue"), ("용도: 작성 가이드", 188, 168, "green"), ("핵심: 판단 가능성", 370, 180, "orange")]:
    d.rounded_rectangle((x, 520, x + w, 554), radius=14, outline=C[col], fill=C["white"], width=2)
    d.text((x + 16, 528), txt, font=font(16, True), fill=C[col])
d.text((54, 618), "배포용 PDF와 편집 가능한 PPT 원본으로 함께 제공", font=font(17), fill=C["muted"])
slides.append(img)

# 2
s = base_slide(prs, 2); kicker(s, "WHY IT MATTERS"); title(s, "신청 목적은 ‘무엇을 해주세요’가 아니라 ‘왜 해야 하는지’를 설명합니다.", size=32)
boxes = [
    ("결재자가 보는 것", "업무상 필요한 신청인지\n비용·시간·권한이 타당한지\n승인 후 처리 방향이 명확한지", 80, "blue"),
    ("작성자가 해야 할 것", "신청 배경을 한 문장으로 쓰고\n구체적 대상·기간·결과를 적고\n불필요한 감정 표현은 줄입니다", 476, "orange"),
    ("좋은 목적의 효과", "반려가 줄고\n추가 질문이 줄고\n처리 속도가 빨라집니다", 872, "green"),
]
for label, val, x, col in boxes:
    add_rect(s, x, 276, 330, 210)
    add_rect(s, x, 276, 5, 210, col, col, False)
    add_textbox(s, label, x + 22, 294, 280, 20, 11, col, True)
    add_textbox(s, val, x + 22, 324, 280, 120, 16, "ink", True, "Georgia")
bullet(s, "결재자는 신청 목적만 보고도 승인 판단의 70%를 끝낼 수 있어야 합니다.", 96, 560, 940)
img, d = image_base(2); d.rectangle((54,100,82,103), fill=C["orange"]); d.text((96,90),"WHY IT MATTERS",font=font(18,True),fill=C["muted"])
draw_wrapped(d,(54,126),"신청 목적은 ‘무엇을 해주세요’가 아니라 ‘왜 해야 하는지’를 설명합니다.",font(39,True),C["ink"],850)
for label, val, x, col in boxes:
    d.rounded_rectangle((x,276,x+330,486), radius=16, fill=C["white"], outline=C["line"])
    d.rectangle((x,276,x+5,486), fill=C[col])
    d.text((x+22,294), label, font=font(18,True), fill=C[col])
    draw_wrapped(d,(x+22,324), val, font(22,True), C["ink"], 280)
d.ellipse((96,570,102,576), fill=C["orange"]); d.text((114,560), "결재자는 신청 목적만 보고도 승인 판단의 70%를 끝낼 수 있어야 합니다.", font=font(22), fill=C["ink"])
slides.append(img)

# 3
s = base_slide(prs, 3); kicker(s, "WRITING FORMULA"); title(s, "목적 문장은 네 조각으로 쓰면 대부분 충분합니다.")
parts = [("상황","현재 어떤 업무 상황인지"),("필요","왜 신청이 필요한지"),("사용","무엇에 사용할 것인지"),("기대 결과","처리 후 어떤 결과가 나는지")]
cols = ["blue","orange","green","red"]
for i,(lab,val) in enumerate(parts):
    x=78+i*274; add_rect(s,x,278,220,170); add_textbox(s,str(i+1),x+18,296,36,32,28,cols[i],True,"Georgia"); add_textbox(s,lab,x+60,304,120,22,15,"ink",True); add_textbox(s,val,x+24,354,170,48,14,"muted")
add_textbox(s, "작성 공식", 88, 512, 120, 22, 12, "orange", True)
add_textbox(s, "[업무 상황] 때문에 [필요 사항]을 [사용 목적]으로 신청하며,\n[기대 결과]를 달성하고자 합니다.", 88, 548, 1040, 58, 23, "ink", True, "Georgia")
img,d=image_base(3); d.rectangle((54,100,82,103),fill=C["orange"]); d.text((96,90),"WRITING FORMULA",font=font(18,True),fill=C["muted"])
draw_wrapped(d,(54,126),"목적 문장은 네 조각으로 쓰면 대부분 충분합니다.",font(44,True),C["ink"],800)
for i,(lab,val) in enumerate(parts):
    x=78+i*274; d.rounded_rectangle((x,278,x+220,448),radius=16,fill=C["white"],outline=C["line"]); d.text((x+18,296),str(i+1),font=font(36,True),fill=C[cols[i]]); d.text((x+60,304),lab,font=font(22,True),fill=C["ink"]); draw_wrapped(d,(x+24,354),val,font(20),C["muted"],170)
d.text((88,512),"작성 공식",font=font(18,True),fill=C["orange"]); draw_wrapped(d,(88,548),"[업무 상황] 때문에 [필요 사항]을 [사용 목적]으로 신청하며,\n[기대 결과]를 달성하고자 합니다.",font(30,True),C["ink"],1040)
slides.append(img)

# 4
s=base_slide(prs,4); kicker(s,"GOOD EXAMPLES"); title(s,"유형별 목적은 ‘업무 대상’과 ‘결과’를 함께 적습니다.")
rows=[("비품·장비","고객 상담 녹취 품질 개선을 위해 마이크를 구매하여 상담 기록 정확도를 높이고자 합니다."),("교육·세미나","경매 실무 상담 역량 강화를 위해 관련 교육을 수강하고, 팀 내 사례 공유 자료로 활용하고자 합니다."),("외근·출장","의뢰 물건 현장 확인과 주변 시세 조사를 위해 외근을 신청하며, 조사 결과를 물건 검토에 반영하고자 합니다."),("자료·권한","담당 사건 진행 현황 확인을 위해 시스템 접근 권한을 신청하며, 업무 처리 지연을 방지하고자 합니다.")]
for i,(lab,val) in enumerate(rows):
    y=230+i*84; add_textbox(s,lab,76,y+8,150,24,15,cols[i],True); add_line(s,236,y+18,1, "line", 44); add_textbox(s,val,268,y,840,54,18,"ink",False,"Georgia")
add_textbox(s,"문장 끝은 ‘하고자 합니다’, ‘활용하고자 합니다’, ‘방지하고자 합니다’처럼 처리 목적이 보이게 마무리합니다.",76,610,960,36,11,"muted")
img,d=image_base(4); d.rectangle((54,100,82,103),fill=C["orange"]); d.text((96,90),"GOOD EXAMPLES",font=font(18,True),fill=C["muted"]); draw_wrapped(d,(54,126),"유형별 목적은 ‘업무 대상’과 ‘결과’를 함께 적습니다.",font(42,True),C["ink"],850)
for i,(lab,val) in enumerate(rows):
    y=230+i*84; d.text((76,y+8),lab,font=font(22,True),fill=C[cols[i]]); d.rectangle((236,y+18,237,y+62),fill=C["line"]); draw_wrapped(d,(268,y),val,font(25),C["ink"],840)
d.text((76,610),"문장 끝은 ‘하고자 합니다’, ‘활용하고자 합니다’, ‘방지하고자 합니다’처럼 처리 목적이 보이게 마무리합니다.",font=font(16),fill=C["muted"])
slides.append(img)

# 5
s=base_slide(prs,5); kicker(s,"BAD TO BETTER"); title(s,"반려되는 목적은 대개 짧아서가 아니라 판단 정보가 부족합니다.")
pairs=[("필요해서 신청합니다.","고객 상담 일정 증가로 상담 내용을 안정적으로 기록하기 위해 녹취 장비를 신청합니다."),("업무에 사용하려고 합니다.","담당 사건 자료 정리와 팀 공유를 위해 클라우드 저장공간을 신청합니다."),("외근 다녀오겠습니다.","2026타경1234 물건 현장 확인 및 주변 시세 조사를 위해 외근을 신청합니다.")]
for i,(bad,good) in enumerate(pairs):
    y=224+i*118; add_rect(s,80,y,430,76,(255,247,245),(240,179,168)); add_textbox(s,"부족",104,y+14,60,18,11,"red",True); add_textbox(s,bad,104,y+38,360,24,18,"ink",False,"Georgia"); add_textbox(s,"→",548,y+24,40,26,24,"muted",False,align="center"); add_rect(s,620,y,560,76,(244,251,246),(169,216,181)); add_textbox(s,"개선",644,y+14,60,18,11,"green",True); add_textbox(s,good,644,y+36,490,28,16,"ink",False,"Georgia")
img,d=image_base(5); d.rectangle((54,100,82,103),fill=C["orange"]); d.text((96,90),"BAD TO BETTER",font=font(18,True),fill=C["muted"]); draw_wrapped(d,(54,126),"반려되는 목적은 대개 짧아서가 아니라 판단 정보가 부족합니다.",font(39,True),C["ink"],900)
for i,(bad,good) in enumerate(pairs):
    y=224+i*118; d.rounded_rectangle((80,y,510,y+76),radius=16,fill=C["pale_red"],outline=(240,179,168)); d.text((104,y+14),"부족",font=font(16,True),fill=C["red"]); d.text((104,y+38),bad,font=font(24),fill=C["ink"]); d.text((548,y+24),"→",font=font(32),fill=C["muted"]); d.rounded_rectangle((620,y,1180,y+76),radius=16,fill=C["pale_green"],outline=(169,216,181)); d.text((644,y+14),"개선",font=font(16,True),fill=C["green"]); draw_wrapped(d,(644,y+36),good,font(21),C["ink"],490)
slides.append(img)

# 6
s=base_slide(prs,6); kicker(s,"DO / DON'T"); title(s,"목적란에는 업무 판단에 필요한 정보만 남깁니다.")
add_rect(s,90,236,500,300,(244,251,246),(169,216,181)); add_textbox(s,"DO",124,268,100,34,30,"green",True,"Georgia")
for j,t in enumerate(["대상: 고객명, 사건번호, 물건, 업무명을 적기","이유: 일정 증가, 처리 지연 방지, 품질 개선 등","결과: 보고서 작성, 상담 대응, 자료 공유처럼 끝 상태 적기"]): bullet(s,t,126,332+j*54,390)
add_rect(s,690,236,500,300,(255,247,245),(240,179,168)); add_textbox(s,"DON'T",724,268,140,34,30,"red",True,"Georgia")
for j,t in enumerate(["‘필요해서’, ‘개인 사정상’처럼 이유가 비어 있는 표현","너무 긴 배경 설명이나 감정적인 표현","승인 후 무엇이 달라지는지 알 수 없는 문장"]): bullet(s,t,726,332+j*54,390)
img,d=image_base(6); d.rectangle((54,100,82,103),fill=C["orange"]); d.text((96,90),"DO / DON'T",font=font(18,True),fill=C["muted"]); draw_wrapped(d,(54,126),"목적란에는 업무 판단에 필요한 정보만 남깁니다.",font(42,True),C["ink"],850)
for x,head,col,items,bg in [(90,"DO","green",["대상: 고객명, 사건번호, 물건, 업무명을 적기","이유: 일정 증가, 처리 지연 방지, 품질 개선 등","결과: 보고서 작성, 상담 대응, 자료 공유처럼 끝 상태 적기"],C["pale_green"]),(690,"DON'T","red",["‘필요해서’, ‘개인 사정상’처럼 이유가 비어 있는 표현","너무 긴 배경 설명이나 감정적인 표현","승인 후 무엇이 달라지는지 알 수 없는 문장"],C["pale_red"])]:
    d.rounded_rectangle((x,236,x+500,536),radius=18,fill=bg,outline=(169,216,181) if col=="green" else (240,179,168)); d.text((x+34,268),head,font=font(38,True),fill=C[col])
    for j,t in enumerate(items):
        y=332+j*54; d.ellipse((x+36,y+9,x+42,y+15),fill=C["orange"]); draw_wrapped(d,(x+56,y),t,font(20),C["ink"],390)
slides.append(img)

# 7
s=base_slide(prs,7); kicker(s,"CHECKLIST"); title(s,"제출 전 30초만 확인하면 반려를 크게 줄일 수 있습니다.")
checks=[("업무 관련성","개인 편의가 아니라 업무 필요성이 보이는가?"),("구체성","대상·기간·장소·사건번호 중 필요한 정보가 들어갔는가?"),("사용 목적","신청한 것을 어디에 사용할지 적었는가?"),("기대 결과","승인 후 어떤 업무 결과가 나오는지 보이는가?"),("문장 길이","한두 문장으로 읽히는가?")]
for i,(lab,val) in enumerate(checks):
    y=208+i*72; circ=s.shapes.add_shape(MSO_SHAPE.OVAL,px(90),px(y+8,"y"),px(24),px(24,"y")); circ.fill.solid(); circ.fill.fore_color.rgb=rgb(C["blue"] if i<4 else C["orange"]); circ.line.fill.background(); add_textbox(s,"✓",95,y+8,16,18,13,"white",True,align="center"); add_textbox(s,lab,136,y,180,24,17,"ink",True); add_textbox(s,val,336,y,720,30,18,"muted",False,"Georgia"); add_line(s,136,y+46,910)
img,d=image_base(7); d.rectangle((54,100,82,103),fill=C["orange"]); d.text((96,90),"CHECKLIST",font=font(18,True),fill=C["muted"]); draw_wrapped(d,(54,126),"제출 전 30초만 확인하면 반려를 크게 줄일 수 있습니다.",font(42,True),C["ink"],850)
for i,(lab,val) in enumerate(checks):
    y=208+i*72; d.ellipse((90,y+8,114,y+32),fill=C["blue"] if i<4 else C["orange"]); d.text((95,y+4),"✓",font=font(18,True),fill=C["white"]); d.text((136,y),lab,font=font(24,True),fill=C["ink"]); d.text((336,y),val,font=font(25),fill=C["muted"]); d.rectangle((136,y+46,1046,y+47),fill=C["line"])
slides.append(img)

# 8
s=base_slide(prs,8); kicker(s,"FINAL TEMPLATE"); title(s,"복붙용 기본 문장으로 시작하고, 빈칸만 업무에 맞게 바꾸세요.")
add_rect(s,96,238,1088,170); add_textbox(s,"[업무 상황]으로 인해 [신청 대상]이 필요하여 신청합니다.\n해당 사항은 [구체적 사용 목적]에 활용하고, [기대 결과]를 달성하기 위한 목적입니다.",136,278,1008,90,24,"ink",True,"Georgia")
add_rect(s,96,462,1088,96,"dark","dark"); add_textbox(s,"한 줄 원칙",136,486,130,24,13,"orange",True); add_textbox(s,"신청 목적은 ‘상황 + 필요 + 사용 + 결과’가 보이면 충분합니다.",288,482,760,34,24,"white",True,"Georgia"); add_textbox(s,"팀장 검토 전에는 목적 문장만 따로 읽어 보고, 신청서 제목 없이도 이해되는지 확인합니다.",136,610,900,32,11,"muted")
img,d=image_base(8); d.rectangle((54,100,82,103),fill=C["orange"]); d.text((96,90),"FINAL TEMPLATE",font=font(18,True),fill=C["muted"]); draw_wrapped(d,(54,126),"복붙용 기본 문장으로 시작하고, 빈칸만 업무에 맞게 바꾸세요.",font(40,True),C["ink"],900)
d.rounded_rectangle((96,238,1184,408),radius=18,fill=C["white"],outline=C["line"]); draw_wrapped(d,(136,278),"[업무 상황]으로 인해 [신청 대상]이 필요하여 신청합니다.\n해당 사항은 [구체적 사용 목적]에 활용하고, [기대 결과]를 달성하기 위한 목적입니다.",font(29,True),C["ink"],1008)
d.rounded_rectangle((96,462,1184,558),radius=18,fill=C["dark"]); d.text((136,486),"한 줄 원칙",font=font(19,True),fill=C["orange"]); d.text((288,482),"신청 목적은 ‘상황 + 필요 + 사용 + 결과’가 보이면 충분합니다.",font=font(30,True),fill=C["white"]); d.text((136,610),"팀장 검토 전에는 목적 문장만 따로 읽어 보고, 신청서 제목 없이도 이해되는지 확인합니다.",font=font(18),fill=C["muted"])
slides.append(img)

prs.save(PPTX)

image_paths = []
for i, img in enumerate(slides, start=1):
    p = PREVIEW / f"slide-{i:02d}.png"
    img.save(p)
    image_paths.append(p)

pdf_images = [Image.open(p).convert("RGB") for p in image_paths]
pdf_images[0].save(PDF, save_all=True, append_images=pdf_images[1:], resolution=150.0)

print(PPTX)
print(PDF)
