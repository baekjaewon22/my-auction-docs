# -*- coding: utf-8 -*-
"""
오케스트레이터: 전체 보고서 생성 파이프라인
- final.py의 main() 함수를 모듈화한 것
- WebSocket으로 진행상황 전송
"""

import os
import re
import time
import logging
import base64
import json
from urllib.parse import urljoin
from typing import Optional, Callable

from pptx import Presentation
from pptx.util import Inches, Pt
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

from ..core.config import (
    settings, CAPTURE_DIR, OUTPUT_DIR, SELENIUM_PROFILE_DIR,
    ensure_dirs, load_config, save_config,
)
from ..core.utils import normalize_myauction_detail_url, track_file, cleanup_generated_files
from ..models.schemas import ReportRequest, ProgressUpdate

from . import crawler
from . import capturer
from . import pdf_processor
from . import ppt_builder
from . import forced_execution_estimator
from . import briefing_opinion
from . import briefing_rights
from .selenium_driver import (
    create_driver, login_myauction, click_tab_safe,
    switch_to_new_window, wait_document_ready, safe_click, navigate_with_retry,
    account_profile_dir, log_detail_page_diagnostics, _dismiss_alert,
)

logger = logging.getLogger(__name__)
logger.info(f"Orchestrator module loaded: {__file__}")

# 이미지 패턴
IMG_PATTERN = str(CAPTURE_DIR / "building_register_{page}.png")
SALE_IMG_PATTERN = str(CAPTURE_DIR / "sale_spec_{page}.png")
STATUS_IMG_PATTERN = str(CAPTURE_DIR / "status_report_{page}.png")
REGISTRY_IMG_PATTERN = str(CAPTURE_DIR / "registry_summary_{page}.png")

# 캡처 파일 경로
COURT_GUIDE_PNG = str(CAPTURE_DIR / "court_guide_capture.png")
KAKAO_MAP_PNG = str(CAPTURE_DIR / "kakao_map.png")
KAKAO_SAT_PNG = str(CAPTURE_DIR / "kakao_satellite.png")
LAND_USE_PLAN_PNG = str(CAPTURE_DIR / "land_use_plan.png")
APPRAISAL_PREFIX = str(CAPTURE_DIR / "appraisal_location_part")
EVICTION_COST_BASIS_PNG = str(CAPTURE_DIR / "eviction_cost_basis.png")

REGISTRY_NEEDLE = "주요 등기사항 요약"
BUILDING_OVERVIEW_PNG = str(CAPTURE_DIR / "building_overview.png")

# 전체 단계 수
TOTAL_STEPS = 6


def _safe_planner_name(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9_-]+", "_", str(value or "planner")).strip("_")[:80] or "planner"


def _planner_snapshot_image(snapshot: dict, index: int) -> str:
    data_url = (
        snapshot.get("image_data_url")
        or snapshot.get("imageDataUrl")
        or (snapshot.get("message") or {}).get("image_data_url")
        or (snapshot.get("message") or {}).get("imageDataUrl")
    )
    if not isinstance(data_url, str) or not data_url.startswith("data:image/"):
        return ""
    try:
        header, encoded = data_url.split(",", 1)
        ext = "jpg" if "jpeg" in header.lower() else "png"
        path = str(CAPTURE_DIR / f"auction_planner_{index}_{_safe_planner_name(snapshot.get('calculator', 'planner'))}.{ext}")
        with open(path, "wb") as file:
            file.write(base64.b64decode(encoded))
        track_file(path)
        return path
    except Exception as exc:
        logger.warning(f"옥션플래너 이미지 저장 실패: {exc}")
        return ""


def _add_planner_appendix_slide(prs: Presentation, title: str, image_path: str) -> bool:
    if not image_path or not os.path.exists(image_path):
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    margin = Inches(0.45)
    title_box = slide.shapes.add_textbox(margin, Inches(0.25), prs.slide_width - margin * 2, Inches(0.35))
    tf = title_box.text_frame
    tf.text = title
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].font.size = Pt(18)
        tf.paragraphs[0].runs[0].font.bold = True
    top = Inches(0.75)
    max_w = prs.slide_width - margin * 2
    max_h = prs.slide_height - top - Inches(0.35)
    try:
        from PIL import Image as PILImage
        with PILImage.open(image_path) as img:
            ratio = img.width / img.height
        width = max_w
        height = int(width / ratio)
        if height > max_h:
            height = max_h
            width = int(height * ratio)
        left = int((prs.slide_width - width) / 2)
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    except Exception:
        slide.shapes.add_picture(image_path, margin, top, width=max_w)
    return True


def _insert_planner_snapshots(prs: Presentation, snapshots: list[dict]) -> list[str]:
    inserted: list[str] = []
    handled_calculators: set[str] = set()
    for index, snapshot in enumerate(snapshots or [], start=1):
        if snapshot.get("include") is False:
            continue
        calculator = str(snapshot.get("calculator") or "")
        if calculator in handled_calculators:
            continue
        handled_calculators.add(calculator)
        label = str(snapshot.get("label") or calculator or "옥션플래너")
        keywords = _planner_snapshot_note_keywords(calculator)
        image_path = _planner_snapshot_image(snapshot, index)
        if image_path:
            ok = ppt_builder.insert_single_image_by_note_keywords(prs, keywords, image_path)
            if not ok:
                ok = _add_planner_appendix_slide(prs, f"옥션플래너 - {label}", image_path)
        else:
            ok = ppt_builder.insert_key_value_table_by_note_keywords(
                prs,
                keywords,
                f"옥션플래너 - {label}",
                _planner_snapshot_rows(snapshot),
            )
        if ok:
            inserted.append(calculator)
    return inserted


_PLANNER_LABELS = {
    "appraisal_price": "감정가",
    "minimum_price": "최저매각가격",
    "bid_price": "예상 입찰가",
    "loan_amount": "대출금액",
    "loan_rate": "대출비율",
    "acquisition_tax": "취득세",
    "local_education_tax": "지방교육세",
    "rural_tax": "농어촌특별세",
    "total_tax": "취득세 합계",
    "purchase_price": "취득가격",
    "total_cost": "총 취득비용",
}


def _planner_snapshot_rows(snapshot: dict) -> list[tuple[str, str]]:
    message = snapshot.get("message") or {}
    preferred = message
    if isinstance(message, dict):
        preferred = (
            message.get("result")
            or message.get("results")
            or message.get("data")
            or message.get("input")
            or message.get("inputs")
            or message.get("payload")
            or message
        )
    rows: list[tuple[str, str]] = []
    ignored = {"source", "type", "calculator", "timestamp", "image_data_url", "imageDataUrl"}

    def append_value(path: str, value, depth: int = 0):
        if len(rows) >= 16 or depth > 3 or value in (None, ""):
            return
        if isinstance(value, dict):
            for key, item in value.items():
                if key in ignored or re.search(r"image|screenshot|capture|thumbnail", str(key), re.I):
                    continue
                append_value(f"{path}.{key}" if path else str(key), item, depth + 1)
            return
        if isinstance(value, list):
            if all(not isinstance(item, (dict, list)) for item in value):
                value = ", ".join(str(item) for item in value)
            else:
                for item_index, item in enumerate(value[:8], start=1):
                    append_value(f"{path} {item_index}", item, depth + 1)
                return
        key = path.split(".")[-1]
        label = _PLANNER_LABELS.get(key, re.sub(r"[_-]+", " ", key).strip())
        if isinstance(value, bool):
            display = "예" if value else "아니오"
        elif isinstance(value, int):
            display = f"{value:,}"
        elif isinstance(value, float):
            display = f"{value:,.2f}".rstrip("0").rstrip(".")
        elif isinstance(value, str):
            display = value[:180]
        else:
            display = json.dumps(value, ensure_ascii=False)[:180]
        rows.append((label or "항목", display))

    append_value("", preferred)
    return rows


def _planner_snapshot_note_keywords(calculator: str) -> list[str]:
    key = str(calculator or "").strip()
    if key == "acquisition-tax":
        return [
            "SLIDE_KEY=OPINION_ACQUISITION_TAX_LOAN",
            "취득세 및 대출",
            "취득세",
        ]
    if key == "loan-bid-estimator":
        return [
            "SLIDE_KEY=OPINION_BID_PRICE_TABLE",
            "EXCEL_RANGE:(1) 예상 입찰가 금액분석표",
            "예상 입찰가 금액분석표",
            "입찰가 산정표",
        ]
    if key == "acquisition-cost-sheet":
        return [
            "SLIDE_KEY=OPINION_ACQUISITION_COST_TABLE",
            "EXCEL_RANGE:(2) 취득시 비용계산표",
            "취득시 비용계산표",
            "취득비용 계산표",
            "비용계산표",
        ]
    return ["SLIDE_KEY=AUCTION_PLANNER", "SLIDE_KEY=CALC_COST"]


def _normalize_match_text(value: str) -> str:
    text = str(value or "").lower()
    text = re.sub(r"\([^)]*\)|\[[^\]]*\]", " ", text)
    text = re.sub(r"[^0-9a-z가-힣]+", "", text)
    return text


def _clean_reference_content(value: str) -> str:
    text = str(value or "").replace("\r", "\n")
    text = re.sub(r"\s*`?\[[A-Z]{2,4}-\d{2}\]`?", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


_PROPERTY_TYPE_ALIASES: list[tuple[str, tuple[str, ...]]] = [
    ("아파트형공장", ("아파트형공장", "지식산업센터")),
    ("아파트상가", ("아파트상가", "단지내상가")),
    ("다세대(빌라)", ("다세대", "빌라", "연립")),
    ("다가구주택", ("다가구",)),
    ("근린주택", ("근린주택", "상가주택")),
    ("도시형생활주택", ("도시형생활주택",)),
    ("오피스텔", ("오피스텔",)),
    ("아파트", ("아파트",)),
    ("주택", ("단독주택", "전원주택", "주택")),
    ("자동차관련시설", ("자동차관련시설", "자동차정비", "세차장")),
    ("장례관련시설", ("장례식장", "봉안", "장례관련시설")),
    ("콘도(호텔)", ("콘도", "호텔")),
    ("펜션(캠핑장)", ("펜션", "캠핑장", "야영장")),
    ("근린상가", ("근린상가",)),
    ("근린시설", ("근린생활시설", "근린시설")),
    ("숙박시설", ("숙박시설", "모텔", "여관")),
    ("목욕시설", ("목욕시설", "목욕장", "사우나")),
    ("운동시설", ("운동시설", "체육시설")),
    ("휴게시설", ("휴게시설", "휴게음식점")),
    ("노유자시설", ("노유자시설", "노인복지", "아동복지")),
    ("교육시설", ("교육시설", "학원")),
    ("주유소", ("주유소",)),
    ("병원", ("병원", "의료시설")),
    ("창고", ("창고시설", "물류센터", "창고")),
    ("공장", ("공장", "제조시설")),
    ("상가", ("상가", "판매시설")),
    ("공장용지", ("공장용지",)),
    ("창고용지", ("창고용지",)),
    ("목장용지", ("목장용지",)),
    ("기타토지", ("기타토지",)),
    ("주차장", ("주차장",)),
    ("과수원", ("과수원",)),
    ("잡종지", ("잡종지",)),
    ("임야", ("임야", "산지")),
    ("대지", ("대지",)),
    ("도로", ("도로", "사도", "현황도로")),
    ("유지", ("유지", "저수지")),
    ("하천", ("하천",)),
    ("구거", ("구거",)),
    ("묘지", ("묘지", "분묘")),
    ("전", ("전",)),
    ("답", ("답",)),
    ("축사(농가시설)", ("축사", "농가시설")),
    ("광업권", ("광업권",)),
    ("어업권", ("어업권",)),
    ("양어장", ("양어장", "양식장")),
    ("종교시설", ("종교시설", "교회", "사찰", "성당")),
    ("중장비", ("중장비", "건설기계")),
    ("선박", ("선박",)),
    ("차량", ("차량", "자동차")),
    ("학교", ("학교",)),
]


def _canonical_property_type(value: str) -> str:
    compact = _normalize_match_text(value)
    if not compact:
        return ""
    for canonical, aliases in _PROPERTY_TYPE_ALIASES:
        if any(_normalize_match_text(alias) in compact for alias in aliases):
            return canonical
    return "기타" if "기타" in compact else ""


def _find_applicable_checklist_references(data: dict, references: dict) -> list[dict]:
    checklist = references.get("checklist") if isinstance(references, dict) else []
    if not isinstance(checklist, list):
        return []
    property_text = " ".join(
        str(value or "") for value in (data.get("item_type"), data.get("item_category"))
    )
    canonical = _canonical_property_type(property_text)
    common: list[dict] = []
    specific: list[dict] = []
    seen: set[str] = set()
    for item in checklist:
        if not isinstance(item, dict):
            continue
        item_id = str(item.get("id") or f"{item.get('category')}:{item.get('title')}")
        if item_id in seen:
            continue
        title = str(item.get("title") or "")
        category = str(item.get("category") or "")
        if "공통" in category or "공통" in title:
            common.append(item)
            seen.add(item_id)
            continue
        if canonical and _canonical_property_type(title) == canonical:
            specific.append(item)
            seen.add(item_id)
    if not specific:
        fallback = (
            _find_checklist_reference(str(data.get("item_category") or ""), references)
            or _find_checklist_reference(str(data.get("item_type") or ""), references)
        )
        if fallback:
            specific.append(fallback)
    return [*common, *specific]


def _find_checklist_reference(item_category: str, references: dict) -> dict:
    category = _normalize_match_text(item_category)
    if not category:
        return {}
    checklist = references.get("checklist") if isinstance(references, dict) else []
    if not isinstance(checklist, list):
        return {}

    best: tuple[int, dict] = (0, {})
    for item in checklist:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or "")
        group = str(item.get("category") or "")
        haystacks = [_normalize_match_text(title), _normalize_match_text(group)]
        score = 0
        for haystack in haystacks:
            if not haystack:
                continue
            if haystack == category:
                score = max(score, 100)
            elif category in haystack or haystack in category:
                score = max(score, 80)
            else:
                overlap = len(set(category) & set(haystack))
                if overlap >= max(2, min(len(category), len(haystack)) // 2):
                    score = max(score, overlap)
        if score > best[0]:
            best = (score, item)
    return best[1] if best[0] >= 2 else {}


def _append_briefing_special_references(special_opinion: str, data: dict, request) -> str:
    additions: list[str] = []
    item_category = str(data.get("item_category") or "").strip()
    item_type = str(data.get("item_type") or "").strip()
    references = getattr(request, "auction_references", {}) or {}
    matched_items = _find_applicable_checklist_references(data, references)
    for matched in matched_items:
        title = str(matched.get("title") or item_type or item_category).strip()
        category_label = "공통" if "공통" in title else (item_type or item_category or title)
        content = _clean_reference_content(matched.get("content") or "")
        checklist_lines: list[str] = []
        for raw_line in content.splitlines():
            line = re.sub(r"^[\s\-•ㆍ*□☐✅✔]+", "", raw_line).strip()
            line = re.sub(r"^\d+[.)]\s*", "", line).strip()
            if len(line) < 3 or line in checklist_lines:
                continue
            checklist_lines.append(line[:110])
            if len(checklist_lines) >= 16:
                break
        header = f"물건별 체크리스트 점검: [{category_label}] {title}"
        if checklist_lines:
            additions.append("\n".join([header, *(f"  · {line}" for line in checklist_lines)]))
        else:
            additions.append(f"{header}\n  · 등기·공부·현장 확인항목을 담당자가 최종 확인해야 합니다.")

    if not additions:
        return special_opinion or ""

    sections: list[str] = []
    if str(special_opinion or "").strip():
        sections.append(str(special_opinion).strip())
    sections.extend(f"- {block}" for block in additions)
    return "\n\n".join(sections).strip()


def _short_selenium_message(exc: Exception, fallback: str) -> str:
    text = str(exc or "").strip()
    compact = re.sub(r"\s+", " ", text)
    if not text or compact.startswith("Message: Stacktrace:") or text.startswith("Message: \nStacktrace:"):
        return fallback
    compact = re.sub(r"Stacktrace:.*$", "", compact).strip()
    return compact[:500] or fallback


def _find_public_data_link(driver, timeout: int = 20):
    end = time.time() + timeout
    xpath_candidates = [
        "//div[@id='dtlw_link']//a[normalize-space(.)='부동산표시']",
        "//div[@id='dtlw_link']//a[contains(normalize-space(.), '부동산표시')]",
        "//div[@id='dtlw_link']//a[contains(normalize-space(.), '공시자료')]",
        "//a[contains(normalize-space(.), '부동산표시')]",
        "//a[contains(normalize-space(.), '공시자료')]",
        "//a[contains(@href, 'aceeair') or contains(@onclick, 'aceeair')]",
        "//*[self::a or self::button][.//img[contains(@alt, '공시') or contains(@alt, '부동산')]]",
    ]
    reject_words = ("매각물건명세서", "물건명세서", "현황조사서", "등기부", "감정평가서", "물건사진")
    accept_words = ("부동산표시", "공시자료")
    last_err = None
    while time.time() < end:
        for xpath in xpath_candidates:
            try:
                elements = driver.find_elements(By.XPATH, xpath)
                if elements:
                    logger.info(f"공시자료 링크 후보 발견: xpath={xpath}, count={len(elements)}")
                for el in elements:
                    try:
                        text = (el.text or el.get_attribute("title") or el.get_attribute("href") or el.get_attribute("onclick") or "").strip()
                        text_clean = re.sub(r"\s+", " ", text)[:160]
                        logger.info(
                            f"공시자료 링크 후보 상태: displayed={el.is_displayed()}, enabled={el.is_enabled()}, "
                            f"text={text_clean}"
                        )
                    except Exception:
                        pass
                    try:
                        probe_text = driver.execute_script("""
                            const el = arguments[0];
                            return [
                              el.innerText,
                              el.textContent,
                              el.getAttribute('title'),
                              el.getAttribute('alt'),
                              el.getAttribute('href'),
                              el.getAttribute('onclick')
                            ].filter(Boolean).join(' ');
                        """, el) or ""
                    except Exception:
                        probe_text = ""
                    if any(word in probe_text for word in reject_words) and not any(word in probe_text for word in accept_words):
                        probe_preview = re.sub(r"\s+", " ", probe_text)[:160]
                        logger.info(f"공시자료 후보 제외: text={probe_preview}")
                        continue
                    if el.is_displayed() and el.is_enabled():
                        return el
            except Exception as e:
                last_err = e
        try:
            el = driver.execute_script("""
                const needles = ['부동산표시', '공시자료'];
                const rejects = ['매각물건명세서', '물건명세서', '현황조사서', '등기부', '감정평가서', '물건사진'];
                const nodes = Array.from(document.querySelectorAll('a, button, area, input[type=button], input[type=image], img'));
                for (const node of nodes) {
                  const holder = node.closest('a, button') || node;
                  const text = [
                    node.innerText,
                    node.textContent,
                    node.getAttribute('title'),
                    node.getAttribute('alt'),
                    node.getAttribute('value'),
                    node.getAttribute('href'),
                    node.getAttribute('onclick'),
                    holder.getAttribute('title'),
                    holder.getAttribute('href'),
                    holder.getAttribute('onclick')
                  ].filter(Boolean).join(' ');
                  if (!needles.some((needle) => text.includes(needle))) continue;
                  if (rejects.some((word) => text.includes(word)) && !needles.some((needle) => text.includes(needle))) continue;
                  const rect = holder.getBoundingClientRect();
                  const style = window.getComputedStyle(holder);
                  if (rect.width > 0 && rect.height > 0 && style.visibility !== 'hidden' && style.display !== 'none') {
                    return holder;
                  }
                }
                return null;
            """)
            if el:
                return el
        except Exception as e:
            last_err = e
        time.sleep(0.4)
    raise RuntimeError("공시자료 팝업 링크를 찾지 못했습니다. 마이옥션 화면 구조가 변경되었거나 해당 사건에 공시자료 버튼이 없습니다.") from last_err


def _open_public_data_page(driver, link_el, timeout: int = 15) -> str:
    info = driver.execute_script("""
        const el = arguments[0];
        return {
          text: `${el.innerText || el.textContent || ''}`.trim(),
          href: el.getAttribute('href') || '',
          onclick: el.getAttribute('onclick') || ''
        };
    """, link_el) or {}
    logger.info(
        "공시자료 선택 링크: "
        f"text={info.get('text', '')}, href={info.get('href', '')}, onclick={info.get('onclick', '')}"
    )

    href = (info.get("href") or "").strip()
    onclick = (info.get("onclick") or "").strip()
    popup_path = driver.execute_script("""
        const href = arguments[0] || '';
        const onclick = arguments[1] || '';
        const text = `${onclick} ${href}`;
        const match = text.match(/windowOpen\\(['"]([^'"]+)['"]/) ||
                      text.match(/open\\(['"]([^'"]+)['"]/) ||
                      text.match(/location\\.href\\s*=\\s*['"]([^'"]+)['"]/);
        return match ? match[1] : '';
    """, href, onclick)

    direct_url = popup_path or href
    if direct_url and not direct_url.lower().startswith("javascript:") and direct_url != "#none":
        if direct_url.startswith("../"):
            direct_url = "/" + direct_url[3:]
        direct_url = urljoin(driver.current_url, direct_url)
        logger.info(f"공시자료 URL 직접 이동: {direct_url}")
        driver.get(direct_url)
        wait_document_ready(driver, timeout=timeout)
        return driver.current_window_handle

    before_handles = list(driver.window_handles)
    safe_click(driver, link_el)
    before_set = set(before_handles)
    end = time.time() + min(4, timeout)
    while time.time() < end:
        new_handles = list(set(driver.window_handles) - before_set)
        if new_handles:
            driver.switch_to.window(new_handles[0])
            wait_document_ready(driver, timeout=timeout)
            logger.info(f"공시자료 새 창 열림: {driver.current_url}")
            return new_handles[0]
        time.sleep(0.2)

    wait_document_ready(driver, timeout=timeout)
    logger.info(f"공시자료 현재 창 열림: {driver.current_url}")
    return driver.current_window_handle


def _safe_filename_part(value: str) -> str:
    text = re.sub(r"\s+", "", str(value or "")).strip()
    text = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", text)
    text = text.strip(" ._-")
    return text[:80]


def _briefing_output_file(data: dict, task_id: Optional[str] = None) -> str:
    case_number = _safe_filename_part(data.get("case_number") or "")
    if not case_number:
        case_number = _safe_filename_part(task_id or "") or time.strftime("%Y%m%d_%H%M%S")

    output_dir = os.path.dirname(settings.output_file) or str(OUTPUT_DIR)
    output_ext = os.path.splitext(settings.output_file)[1] or ".pptm"
    return os.path.join(output_dir, f"브리핑자료_{case_number}{output_ext}")


def _apply_author_fields(data: dict, request: ReportRequest) -> None:
    author_name = str(getattr(request, "author_name", "") or "").strip()
    author_title = str(getattr(request, "author_title", "") or "").strip()
    author_phone = str(getattr(request, "author_phone", "") or "").strip()
    author_name_title = " ".join(part for part in (author_name, author_title) if part).strip()

    data["authorName"] = author_name
    data["authorTitle"] = author_title
    data["authorPhone"] = author_phone
    data["가입자 성명"] = author_name
    data["가입자 직책"] = author_title
    data["가입자 성명 직책"] = author_name_title
    data["가입자 전화번호"] = author_phone


async def generate_report(
    request: ReportRequest,
    progress_callback: Optional[Callable] = None,
    task_id: Optional[str] = None,
) -> dict:
    logger.info(f"generate_report using orchestrator: {__file__}")
    """
    전체 보고서 생성 파이프라인
    progress_callback(ProgressUpdate) 로 진행상황 전달
    """
    ensure_dirs()

    diagnostics: list[dict] = []

    def add_diagnostic(key: str, label: str, status: str, message: str):
        diagnostics.append({
            "key": key,
            "label": label,
            "status": status,
            "message": str(message or "").strip(),
        })

    def emit(step, title, message, status="running", percent=0.0):
        if progress_callback:
            try:
                update = ProgressUpdate(
                    step=step, total_steps=TOTAL_STEPS,
                    title=title, message=message,
                    status=status, percent=percent,
                )
                # 동기/비동기 콜백 모두 지원
                import asyncio
                if asyncio.iscoroutinefunction(progress_callback):
                    try:
                        loop = asyncio.get_event_loop()
                        if loop.is_running():
                            loop.create_task(progress_callback(update))
                        else:
                            asyncio.run(progress_callback(update))
                    except RuntimeError:
                        pass
                else:
                    progress_callback(update)
            except Exception:
                pass
        logger.info(f"[{step}/{TOTAL_STEPS}] {title}: {message}")

    input_url = (request.url or "").strip()
    url = normalize_myauction_detail_url(input_url, request.myauction_id)
    logger.info(f"입력 URL: {input_url}")
    logger.info(f"URL 정규화 결과: {url}")
    logger.info(
        "URL 정규화 상세: "
        f"view_to_view3={'/view/' in input_url and '/view3/' not in input_url}, "
        f"myauction_id_appended={bool(request.myauction_id and request.myauction_id.strip() and request.myauction_id.strip() in url)}, "
        f"final_url={url}"
    )

    # 초기화
    data = {}
    prs = None
    LAND_MODE = False
    total_building = total_sale = total_registry = total_status = 0
    total_registry_land = 0
    tenant_imgs = building_registry_imgs = land_registry_imgs = []
    land_use_plan_img = kakao_map_img = kakao_sat_img = ""
    loc_left_img = loc_right_img = ""
    court_start_time = court_end_time = court_capture_png = ""
    building_overview_img = ""
    eviction_cost_basis_img = ""
    registry_land_img_pattern = ""
    checklist_match_count = 0
    checklist_applied = False
    planner_inserted_calculators: list[str] = []
    diagnostic_reasons: dict[str, str] = {}

    # ===== STEP 0: Selenium 준비 =====
    emit(0, "브라우저 준비", "Chrome 시작 중...")

    profile_dir = account_profile_dir(request.myauction_id) if request.remember_login else ""
    logger.info(
        f"Chrome profile 선택: remember_login={request.remember_login}, "
        f"profile_dir={profile_dir or '(임시/비저장 프로필)'}"
    )

    driver = create_driver(profile_dir=profile_dir, headless=True)

    try:
        # ===== STEP 1: 로그인 + 파싱 =====
        emit(1, "사이트 파싱", "마이옥션 로그인 중...", percent=5)
        login_myauction(driver, request.myauction_id, request.myauction_pw)

        emit(1, "사이트 파싱", "상세 페이지 접속 중...", percent=10)
        navigate_with_retry(driver, url, retries=3)
        wait_document_ready(driver, timeout=30)
        time.sleep(5)
        current_url = driver.current_url or ""
        page_source = driver.page_source or ""
        logger.info(f"MyAuction detail loaded: url={current_url}, title={driver.title or ''}, html_length={len(page_source)}")
        log_detail_page_diagnostics(driver, prefix="브리핑자료 상세 진입 직후")
        if "member/login.php" in current_url or ("id=\"id\"" in page_source and "passwd" in page_source):
            raise RuntimeError("마이옥션 로그인이 유지되지 않아 사건 상세 페이지 대신 로그인 화면이 열렸습니다. 저장된 마이옥션 ID/PW를 다시 확인해 주세요.")

        emit(1, "사이트 파싱", "데이터 추출 중...", percent=15)
        soup = crawler.fetch_soup_from_driver(driver)
        data = crawler.parse_myauction_detail(soup, url, driver=driver)
        _apply_author_fields(data, request)
        LAND_MODE = bool(data.get("LAND_MODE", False))

        # 토지이용계획 텍스트
        try:
            landplan_url = (data.get("landplan_url") or "").strip()
            if landplan_url:
                refined = crawler.fetch_land_zoning_from_plan(driver, landplan_url)
                if refined:
                    data["land_zoning"] = refined
        except Exception as e:
            logger.warning(f"토지이용계획 추출 실패: {e}")

        rights_analysis_opinion = ""
        special_opinion = ""
        try:
            emit(1, "사이트 파싱", "권리분석 정보 확인 중...", percent=18)
            rights_context = briefing_rights.extract_context(soup, driver=driver, task_id=task_id)
            if rights_context:
                data.update(rights_context)
            briefing_rights_data = briefing_rights.build_opinion_data(data)
            rights_analysis_opinion = briefing_opinion.build_rights_analysis_opinion(briefing_rights_data)
            special_opinion = briefing_opinion.build_special_opinion(briefing_rights_data)
        except Exception as e:
            logger.warning(f"담당자 종합의견 (2) 권리분석 문안 구성 실패: {e}")
            try:
                briefing_rights_data = briefing_rights.build_opinion_data(data)
                rights_analysis_opinion = briefing_opinion.build_rights_analysis_opinion(briefing_rights_data)
                special_opinion = briefing_opinion.build_special_opinion(briefing_rights_data)
            except Exception as fallback_error:
                logger.warning(f"담당자 종합의견 (2) 권리분석 기본 문안 구성 실패: {fallback_error}")

        # 권리분석 OCR 일부가 실패하더라도 선택 물건 체크리스트는 항상 반영한다.
        try:
            checklist_match_count = len(_find_applicable_checklist_references(
                data,
                getattr(request, "auction_references", {}) or {},
            ))
            special_opinion = _append_briefing_special_references(special_opinion, data, request)
        except Exception as e:
            diagnostic_reasons["checklist"] = _short_selenium_message(e, "체크리스트 구성 중 오류 발생")
            logger.warning(f"물건별 체크리스트 반영 실패: {e}")
        data["rights_analysis_opinion"] = rights_analysis_opinion
        data["special_opinion"] = special_opinion

        try:
            data["property_status_opinion"] = briefing_opinion.build_property_status_opinion(data)
        except Exception as e:
            logger.warning(f"담당자 종합의견 (1) 물건현황 문안 구성 실패: {e}")

        try:
            eviction_values = forced_execution_estimator.build_eviction_cost_values(data)
            for key, value in eviction_values.items():
                if key != "cost":
                    data[key] = value
        except Exception as e:
            logger.warning(f"명도비 템플릿 변수 구성 실패: {e}")

        emit(1, "사이트 파싱", "파싱 완료", percent=20)

        # ===== STEP 2: PPT 기본 채우기 =====
        emit(2, "PPT 기본값", "템플릿 로드 중...", percent=20)
        prs = Presentation(settings.pptm_template)
        logger.info(f"PPT template loaded, continue_on_capture_error=True, prs_ready={prs is not None}")
        ppt_builder.fill_slide_with_data(prs, data)
        try:
            property_status_opinion = data.get("property_status_opinion") or briefing_opinion.build_property_status_opinion(data)
            data["property_status_opinion"] = property_status_opinion
            if ppt_builder.apply_property_status_opinion(prs, property_status_opinion):
                logger.info("담당자 종합의견 (1) 물건현황 자동 작성 완료")
        except Exception as e:
            logger.warning(f"담당자 종합의견 (1) 물건현황 작성 실패: {e}")
        try:
            if ppt_builder.apply_rights_analysis_opinion(prs, rights_analysis_opinion):
                logger.info("담당자 종합의견 (2) 권리분석 자동 작성 완료")
        except Exception as e:
            logger.warning(f"담당자 종합의견 (2) 권리분석 작성 실패: {e}")
        try:
            special_opinion = data.get("special_opinion") or special_opinion
            if ppt_builder.apply_special_opinion(prs, special_opinion):
                checklist_applied = checklist_match_count > 0
                logger.info("담당자 종합의견 (3) 특이사항 자동 작성 완료")
        except Exception as e:
            diagnostic_reasons["checklist"] = _short_selenium_message(e, "특이사항 슬라이드 반영 중 오류 발생")
            logger.warning(f"담당자 종합의견 (3) 특이사항 작성 실패: {e}")
        ppt_builder.insert_main_photo(prs, data.get("photo_url", ""))
        emit(2, "PPT 기본값", "기본값 채우기 완료", percent=25)

        # ===== STEP 3: 캡처 =====
        emit(3, "문서 캡처", "관할법원안내 처리 중...", percent=25)
        base_handle = driver.current_window_handle

        # 관할법원안내
        try:
            court_start_time, court_end_time, court_popup = capturer.open_court_guide_popup(driver)
            court_capture_png = capturer.capture_court_popup(driver, COURT_GUIDE_PNG)
            driver.close()
            driver.switch_to.window(base_handle)
            wait_document_ready(driver)
        except Exception as e:
            logger.warning(f"관할법원안내 실패: {e}")
            try:
                if driver.current_window_handle != base_handle:
                    driver.switch_to.window(base_handle)
            except Exception:
                pass

        # 토지이용계획 캡처
        emit(3, "문서 캡처", "토지이용계획 캡처 중...", percent=30)
        try:
            land_use_plan_img = capturer.capture_land_use_plan(driver, LAND_USE_PLAN_PNG)
        except Exception as e:
            logger.warning(f"토지이용계획 캡처 실패: {e}")

        # 임차인/등기부 캡처
        emit(3, "문서 캡처", "임차인/등기부현황 캡처 중...", percent=35)
        try:
            tenant_imgs, building_registry_imgs, land_registry_imgs, table_capture_reasons = capturer.capture_tenant_and_registry(driver)
            diagnostic_reasons.update(table_capture_reasons)
            logger.info(
                "임차인/등기부현황 캡처 파일 확인: "
                f"tenant={[(p, os.path.exists(p), os.path.getsize(p) if os.path.exists(p) else 0) for p in tenant_imgs]}, "
                f"building={[(p, os.path.exists(p), os.path.getsize(p) if os.path.exists(p) else 0) for p in building_registry_imgs]}, "
                f"land={[(p, os.path.exists(p), os.path.getsize(p) if os.path.exists(p) else 0) for p in land_registry_imgs]}"
            )
        except Exception as e:
            diagnostic_reasons["tenant-status"] = _short_selenium_message(e, "임차인 현황 캡처 단계 오류")
            logger.warning(f"캡처 실패: {e}")

        # 예상명도비용은 캡처 이미지를 삽입하지 않고 템플릿 변수로 자동 입력한다.
        emit(3, "문서 캡처", "예상명도비용 계산값 준비 중...", percent=38)

        # 공시자료 팝업
        emit(3, "문서 캡처", "공시자료 팝업 열기...", percent=40)
        wait = WebDriverWait(driver, 20)
        try:
            bu_link = _find_public_data_link(driver, timeout=25)
            popup_handle = _open_public_data_page(driver, bu_link, timeout=20)
        except Exception as e:
            reason = _short_selenium_message(e, "공시자료 팝업 진입 실패")
            for diagnostic_key in ("building-register", "sale-spec", "registry-summary", "status-report"):
                diagnostic_reasons[diagnostic_key] = reason
            raise RuntimeError(_short_selenium_message(
                e,
                "공시자료 팝업 링크를 찾거나 클릭하지 못했습니다. 마이옥션 사건 상세 화면에서 공시자료/부동산표시 버튼 노출 여부를 확인해 주세요.",
            )) from e
        time.sleep(1)

        popup_handle = popup_handle or driver.current_window_handle
        wait = WebDriverWait(driver, 15)

        # 카카오맵
        emit(3, "문서 캡처", "전자지도/위성지도 캡처 중...", percent=45)
        try:
            kakao_map_img = capturer.open_kakao_and_capture(driver, popup_handle, "전자지도", KAKAO_MAP_PNG)
        except Exception as e:
            logger.warning(f"전자지도 캡처 실패: {e}")
        try:
            kakao_sat_img = capturer.open_kakao_and_capture(driver, popup_handle, "위성지도", KAKAO_SAT_PNG)
        except Exception as e:
            logger.warning(f"위성지도 캡처 실패: {e}")

        # [MODE] 건축물대장 탭 존재시 건축물 버전 강제 전환
        if LAND_MODE:
            try:
                _tab = click_tab_safe(wait, driver, ["건축물대장", "건축물"])
                if _tab:
                    logger.info("토지로 판별됐지만 '건축물대장' 탭 존재 → 건축물버전으로 전환")
                    LAND_MODE = False
                    data["LAND_MODE"] = False
                    data["BUILDING_MODE"] = True
            except Exception:
                pass

        # 건축물대장
        emit(3, "문서 캡처", "건축물대장 처리 중...", percent=50)
        if not LAND_MODE:
            try:
                clicked = click_tab_safe(wait, driver, ["건축물대장", "건축물"])
                if not clicked:
                    raise RuntimeError("건축물대장 탭을 찾지 못했습니다.")
                logger.info(f"팝업 내 '{clicked}' 탭 클릭 완료")
                time.sleep(1)
                iframe = wait.until(EC.presence_of_element_located((By.ID, "detail_target")))
                pdf_url = iframe.get_attribute("src")
                if not pdf_url:
                    raise RuntimeError("건축물대장 iframe src 없음")
                pdf_path = pdf_processor.download_pdf_with_cookies(driver, pdf_url, "building_register")
                total_building = pdf_processor.pdf_to_images(pdf_path, IMG_PATTERN, dpi=250)
            except Exception as e:
                diagnostic_reasons["building-register"] = _short_selenium_message(e, "건축물대장 캡처 실패")
                logger.warning(f"건축물대장 실패 → 생략(계속 진행): {e}")
        else:
            logger.info("토지버전 → 건축물대장 생략")

        # 매각물건명세서
        emit(3, "문서 캡처", "매각물건명세서 처리 중...", percent=55)
        try:
            clicked = click_tab_safe(wait, driver, ["매각물건명세서", "물건명세서"])
            time.sleep(1)
            iframe = wait.until(EC.presence_of_element_located((By.ID, "detail_target")))
            pdf_url = iframe.get_attribute("src")
            if pdf_url:
                pdf_path = pdf_processor.download_pdf_with_cookies(driver, pdf_url, "sale_spec")
                total_sale = pdf_processor.pdf_to_images(pdf_path, SALE_IMG_PATTERN, dpi=250)
        except Exception as e:
            diagnostic_reasons["sale-spec"] = _short_selenium_message(e, "매각물건명세서 캡처 실패")
            logger.warning(f"매각물건명세서 실패: {e}")

        # 등기부(건물)
        emit(3, "문서 캡처", "등기부 처리 중...", percent=60)
        if not LAND_MODE:
            try:
                clicked = click_tab_safe(wait, driver, ["등기부(건물)", "등기부", "건물"])
                if not clicked:
                    raise RuntimeError("등기부(건물) 탭을 찾지 못했습니다.")
                logger.info(f"팝업 내 '{clicked}' 탭 클릭 완료")
                time.sleep(1)
                iframe = wait.until(EC.presence_of_element_located((By.ID, "detail_target")))
                pdf_url = iframe.get_attribute("src")
                if not pdf_url:
                    raise RuntimeError("등기부(건물) iframe src 없음")

                reg_pdf = pdf_processor.download_pdf_with_cookies(driver, pdf_url, "registry_building")

                # "주요 등기사항 요약" 문구 페이지 찾기 → 가로 변환
                start_idx = pdf_processor.find_first_page_contains_text(reg_pdf, REGISTRY_NEEDLE)
                if start_idx == -1:
                    # 문구 못 찾음 → 마지막 페이지만 가로 변환
                    logger.warning(f"'{REGISTRY_NEEDLE}' 문구를 못 찾음 → 마지막 페이지만 사용")
                    reg_land_pdf = reg_pdf.replace(".pdf", "_last_landscape.pdf")
                    reg_pdf = pdf_processor.pdf_last_page_to_landscape(reg_pdf, reg_land_pdf, dpi=220)
                    track_file(reg_pdf)
                else:
                    # 문구 발견 → 해당 페이지부터 끝까지 가로 변환
                    logger.info(f"'{REGISTRY_NEEDLE}' 발견: {start_idx+1}페이지부터 가로 변환")
                    reg_land_pdf = reg_pdf.replace(".pdf", f"_from_{start_idx+1}_landscape.pdf")
                    reg_pdf = pdf_processor.pdf_pages_from_to_landscape(reg_pdf, reg_land_pdf, start_idx, dpi=220)
                    track_file(reg_pdf)

                total_registry = pdf_processor.pdf_to_images(reg_pdf, REGISTRY_IMG_PATTERN, dpi=250)
            except Exception as e:
                diagnostic_reasons["registry-summary"] = _short_selenium_message(e, "건물 등기사항 요약 캡처 실패")
                logger.warning(f"등기부(건물) 실패 → 생략(계속 진행): {e}")
        else:
            logger.info("토지버전 → 등기부(건물) 생략")

        # 등기부(토지) - 토지 모드일 때만
        if LAND_MODE:
            try:
                emit(3, "문서 캡처", "등기부(토지) 처리 중...", percent=62)
                clicked = click_tab_safe(wait, driver, ["등기부(토지)", "등기부", "토지"])
                if not clicked:
                    raise RuntimeError("등기부(토지) 탭을 찾지 못했습니다.")
                logger.info(f"팝업 내 '{clicked}' 탭 클릭 완료")
                time.sleep(1)
                iframe = wait.until(EC.presence_of_element_located((By.ID, "detail_target")))
                pdf_url = iframe.get_attribute("src")
                if not pdf_url:
                    raise RuntimeError("등기부(토지) iframe src 없음")
                reg_land_pdf = pdf_processor.download_pdf_with_cookies(driver, pdf_url, "registry_land")
                registry_land_img_pattern = str(CAPTURE_DIR / "registry_land_{page}.png")
                total_registry_land = pdf_processor.pdf_to_images(reg_land_pdf, registry_land_img_pattern, dpi=250)
            except Exception as e:
                diagnostic_reasons["registry-summary"] = _short_selenium_message(e, "토지 등기사항 요약 캡처 실패")
                logger.warning(f"등기부(토지) 실패 → 생략(계속 진행): {e}")

        # 감정평가서
        emit(3, "문서 캡처", "감정평가서 위치도 처리 중...", percent=65)
        try:
            clicked = click_tab_safe(wait, driver, ["감정평가서", "감정평가", "감정"])
            time.sleep(1)
            iframe = wait.until(EC.presence_of_element_located((By.ID, "detail_target")))
            pdf_url = iframe.get_attribute("src")
            if pdf_url:
                appr_pdf = pdf_processor.download_pdf_with_cookies(driver, pdf_url, "appraisal_report")

                # (A) 위치도 탐색 + 렌더링
                found = pdf_processor.find_appraisal_map_pages(appr_pdf)
                chosen = pdf_processor.choose_location_types(found)
                if chosen:
                    page_indices = [found[t] for t in chosen]
                    imgs = pdf_processor.render_pages_vector(appr_pdf, page_indices, APPRAISAL_PREFIX)
                    if len(imgs) >= 1:
                        loc_left_img = imgs[0]
                    if len(imgs) >= 2:
                        loc_right_img = imgs[1]

                # (B) 내부구조도 / 건물개황도 탐색 + 렌더링
                emit(3, "문서 캡처", "내부구조도 탐색 중...", percent=68)
                try:
                    overview_page = pdf_processor.find_building_overview_page(appr_pdf)
                    if overview_page >= 0:
                        overview_imgs = pdf_processor.render_pages_vector(
                            appr_pdf, [overview_page],
                            str(CAPTURE_DIR / "building_overview"),
                            dpi=260,
                        )
                        if overview_imgs:
                            building_overview_img = overview_imgs[0]
                            logger.info(f"도면/내부구조도 이미지 저장 경로: {building_overview_img}")
                except Exception as e2:
                    logger.warning(f"내부구조도 탐색 실패(계속 진행): {e2}")

        except Exception as e:
            logger.warning(f"감정평가서 실패: {e}")

        # 현황조사서
        emit(3, "문서 캡처", "현황조사서 처리 중...", percent=70)
        try:
            clicked = click_tab_safe(wait, driver, ["현황조사서"])
            time.sleep(1)
            status_pdf = pdf_processor.print_current_page_to_pdf(driver, "status_report", landscape=True)
            if pdf_processor.is_valid_pdf(status_pdf):
                total_status = pdf_processor.pdf_to_images(status_pdf, STATUS_IMG_PATTERN, dpi=300)
        except Exception as e:
            diagnostic_reasons["status-report"] = _short_selenium_message(e, "현황조사서 캡처 실패")
            logger.warning(f"현황조사서 실패: {e}")

        emit(3, "문서 캡처", "캡처 완료", percent=75)

    except Exception as e:
        message = _short_selenium_message(e, "문서 캡처 일부 단계에서 응답 대기 시간이 초과되었습니다.")
        logger.warning(f"generate_report capture exception caught, prs_ready={prs is not None}, raw_type={type(e).__name__}, safe_message={message}")
        if prs is not None:
            logger.warning(f"문서 캡처 일부 실패 후 저장 계속: {message}")
            emit(3, "문서 캡처", f"일부 문서 캡처 생략: {message}", percent=75)
        else:
            logger.error(f"파이프라인 실패: {message}")
            emit(0, "오류", message, status="error")
            return {"success": False, "message": message}
    finally:
        try:
            driver.quit()
        except Exception:
            pass

    # ===== STEP 4: PPT 이미지 삽입 =====
    emit(4, "PPT 이미지 삽입", "슬라이드에 이미지 삽입 중...", percent=75)

    if prs is None:
        return {"success": False, "message": "PPT 로드 실패"}

    # 관할법원 텍스트 치환
    try:
        slide40 = ppt_builder.find_slide_by_note_key(prs, "SLIDE_KEY=CALC_COST")
        if slide40:
            deposit_num = re.sub(r"[^0-9,]", "", data.get("deposit", "") or "")
            mapping = {
                "{deposit}": deposit_num,
                "{auction_date}": data.get("auction_date", "") or "",
                "{auction_start_time}": court_start_time or "",
                "Auction_start_time": court_start_time or "",
                "{auction_end_time}": court_end_time or "",
                "Auction_end_time": court_end_time or "",
            }
            ppt_builder.replace_placeholders_in_slide(slide40, mapping)
            ppt_builder.remove_won_unit_in_slide(slide40)
    except Exception as e:
        logger.warning(f"40번 슬라이드 실패: {e}")

    # 관할법원 캡처 삽입
    if court_capture_png and os.path.exists(court_capture_png):
        ppt_builder.insert_single_image(prs, "SLIDE_KEY=COURT_GUIDE", court_capture_png, use_note_key=True)

    # 문서 이미지 삽입
    cap_dir = str(CAPTURE_DIR)
    ppt_builder.insert_images_into_ppt(prs, total_building, "건축물대장", IMG_PATTERN, clone_from_next=True, forced_num="(5)")
    ppt_builder.insert_images_into_ppt(prs, total_sale, "매각물건명세서", SALE_IMG_PATTERN, clone_from_next=True)
    ppt_builder.insert_images_into_ppt(prs, total_registry, "등기사항 요약", REGISTRY_IMG_PATTERN, clone_from_next=True)
    if not total_registry and total_registry_land and registry_land_img_pattern:
        ppt_builder.insert_images_into_ppt(
            prs,
            total_registry_land,
            "등기사항 요약",
            registry_land_img_pattern,
            clone_from_next=True,
        )
    ppt_builder.insert_images_into_ppt(prs, total_status, "현황조사서", STATUS_IMG_PATTERN, clone_from_next=True)

    # 임차인/등기부현황
    if tenant_imgs:
        ppt_builder.insert_images_into_ppt(
            prs, len(tenant_imgs), "임차인 현황",
            os.path.join(cap_dir, "tenant_status_{page}.png"),
            clone_from_next=True, base_as_last=True,
        )

    if building_registry_imgs:
        ppt_builder.insert_images_into_ppt(
            prs, len(building_registry_imgs), "등기부 현황",
            os.path.join(cap_dir, "building_registry_status_{page}.png"),
            clone_from_next=True, base_as_last=True,
        )

    # 단일 이미지 삽입
    ppt_builder.insert_single_image(prs, "전자지도", KAKAO_MAP_PNG)
    ppt_builder.insert_single_image(prs, "위성지도", KAKAO_SAT_PNG)
    ppt_builder.insert_single_image(prs, "토지이용계획", land_use_plan_img)

    # 물건현황 (1): 좌측 위치도, 우측 내부구조도/호별배치도
    location_img = loc_left_img or loc_right_img
    if location_img or building_overview_img:
        inserted = ppt_builder.insert_location_and_structure_images(prs, location_img, building_overview_img)
        if inserted:
            logger.info("물건현황 (1) 위치도/내부구조도 좌우 삽입 완료")
        else:
            logger.warning("물건현황 (1) 위치도/내부구조도 좌우 삽입 실패")

    try:
        planner_inserted_calculators = _insert_planner_snapshots(prs, request.planner_snapshots)
        if planner_inserted_calculators:
            logger.info(f"옥션플래너 스냅샷 {len(planner_inserted_calculators)}건 PPT 삽입 완료")
    except Exception as e:
        diagnostic_reasons["planner"] = _short_selenium_message(e, "옥션플래너 PPT 삽입 실패")
        logger.warning(f"옥션플래너 스냅샷 PPT 삽입 실패: {e}")

    logger.info("예상명도비용 산출근거 캡처 이미지 삽입 생략: 템플릿 변수 자동입력 방식 사용")

    emit(4, "PPT 이미지 삽입", "삽입 완료", percent=85)

    # 명도 정액제/실비제 비용: 엑셀 토큰 적용 후에도 파싱 면적 기반 산식이 최종값이 되도록 저장 직전 반영
    try:
        eviction_values = forced_execution_estimator.build_eviction_cost_values(data)
        updated = ppt_builder.apply_eviction_cost_estimates(prs, eviction_values)
        if updated:
            logger.info(f"명도 정액제/실비제 비용 반영 완료: {updated}개 텍스트")
        else:
            logger.warning("명도 정액제/실비제 비용 반영 대상 텍스트를 찾지 못했습니다.")
    except Exception as e:
        logger.warning(f"명도 정액제/실비제 비용 반영 실패: {e}")

    required_planner = {
        "acquisition-tax": "취득세 및 대출",
        "loan-bid-estimator": "입찰가 산정표",
        "acquisition-cost-sheet": "취득비용 계산표",
    }
    received_snapshots = {
        str(item.get("calculator") or ""): item
        for item in (request.planner_snapshots or [])
        if isinstance(item, dict) and item.get("include") is not False
    }
    for calculator, label in required_planner.items():
        snapshot = received_snapshots.get(calculator)
        if calculator in planner_inserted_calculators:
            snapshot_message = snapshot.get("message") if isinstance(snapshot, dict) else {}
            snapshot_message = snapshot_message if isinstance(snapshot_message, dict) else {}
            image_data = (
                snapshot.get("image_data_url")
                or snapshot.get("imageDataUrl")
                or snapshot_message.get("image_data_url")
                or snapshot_message.get("imageDataUrl")
            ) if snapshot else ""
            has_image = isinstance(image_data, str) and image_data.startswith("data:image/")
            add_diagnostic(
                f"planner:{calculator}",
                label,
                "ok" if has_image else "warning",
                "이미지 수신 및 PPT 반영 완료" if has_image else "이미지 없이 계산값 표로 대체 반영",
            )
        elif snapshot:
            add_diagnostic(
                f"planner:{calculator}",
                label,
                "error",
                diagnostic_reasons.get("planner") or "자료는 수신했지만 PPT 삽입 위치를 찾지 못함",
            )
        else:
            add_diagnostic(f"planner:{calculator}", label, "error", "브라우저에서 이미지 자료가 전달되지 않음")

    if checklist_match_count and checklist_applied:
        add_diagnostic("checklist", "물건별 체크리스트", "ok", f"{checklist_match_count}개 항목 매칭 및 특이사항 반영 완료")
    elif checklist_match_count:
        add_diagnostic("checklist", "물건별 체크리스트", "error", f"{checklist_match_count}개 항목은 매칭됐지만 PPT 반영 실패")
    else:
        add_diagnostic(
            "checklist",
            "물건별 체크리스트",
            "error",
            diagnostic_reasons.get("checklist") or "공통·물건 카테고리 체크리스트가 매칭되지 않음",
        )

    document_checks = [
        ("building-register", "건축물대장", total_building, "skipped" if LAND_MODE else "error"),
        ("tenant-status", "전입세대·임차인 현황", len(tenant_imgs), "error"),
        ("registry-summary", "등기사항 요약", total_registry or total_registry_land, "error"),
        ("status-report", "현황조사서", total_status, "error"),
        ("sale-spec", "매각물건명세서", total_sale, "error"),
    ]
    for key, label, count, empty_status in document_checks:
        if count:
            add_diagnostic(key, label, "ok", f"{count}페이지 캡처 및 PPT 삽입 처리")
        elif empty_status == "skipped":
            add_diagnostic(key, label, "skipped", "토지 물건으로 판별되어 생략")
        else:
            add_diagnostic(
                key,
                label,
                empty_status,
                diagnostic_reasons.get(key) or "캡처 결과가 없어 PPT에 삽입되지 않음",
            )

    # ===== STEP 5: 저장 =====
    emit(5, "저장", "PPT 저장 중...", percent=90)
    output_file = _briefing_output_file(data, task_id)
    os.makedirs(os.path.dirname(output_file), exist_ok=True)
    ppt_builder.save_pptm_preserve_vba(settings.pptm_template, prs, output_file)

    cleanup_generated_files()

    diagnostic_failures = sum(1 for item in diagnostics if item.get("status") in ("warning", "error"))
    completion_message = "보고서 생성이 완료되었습니다."
    if diagnostic_failures:
        completion_message = f"보고서는 생성됐지만 진단 경고 {diagnostic_failures}건을 확인해 주세요."
    emit(5, "저장", completion_message, status="completed", percent=100)
    logger.info(f"보고서 생성 완료: {output_file}")

    return {
        "success": True,
        "output_file": output_file,
        "message": completion_message,
        "data": data,
        "diagnostics": diagnostics,
    }
