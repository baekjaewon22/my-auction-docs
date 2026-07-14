# -*- coding: utf-8 -*-
"""
권리분석 보증서 생성 서비스
- 마이옥션 상세 페이지를 기존 입력값으로 파싱
- HTML 템플릿을 렌더링
- Selenium/Chrome CDP printToPDF로 PDF 저장
"""

import base64
import html
import logging
import os
import re
import time
import asyncio
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Callable, Optional
from urllib.parse import urljoin

import fitz
from PIL import Image, ImageEnhance, ImageOps
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from selenium.webdriver.common.by import By
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.support.ui import WebDriverWait

from ..core.config import (
    CAPTURE_DIR,
    OUTPUT_DIR,
    SELENIUM_PROFILE_DIR,
    TESSERACT_PATH,
    ensure_dirs,
    load_config,
    settings,
)
from ..core.utils import normalize_myauction_detail_url
from ..models.schemas import ProgressUpdate, ReportRequest
from . import capturer, crawler, forced_execution_estimator, pdf_processor
from .special_situations import build_special_issue_lines
from .selenium_driver import (
    account_profile_dir,
    click_tab_safe,
    create_driver,
    log_detail_page_diagnostics,
    login_myauction,
    navigate_with_retry,
    safe_click,
    wait_document_ready,
)

try:
    import pytesseract
    if os.path.exists(TESSERACT_PATH):
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH
except ImportError:
    pytesseract = None

logger = logging.getLogger(__name__)

TOTAL_STEPS = 6
BODY_PLACEHOLDER_TOKENS = (
    "{{baseRightDescription}}",
    "{{tenantAnalysisText}}",
    "{{surplusDescription}}",
    "{{miscText}}",
    "{{reviewText}}",
    "{{specialSummaryText}}",
    "{{unpaidManagementFeeText}}",
    "{{saleSpecRemarksText}}",
    "{{statusSurveyEtcText}}",
    "{{특이사항요약}}",
    "{{미납관리비}}",
    "{{매각물건명세서비고}}",
    "{{현황조사서기타}}",
)
CASE_INFO_PLACEHOLDER_TOKENS = ("{{caseNumber}}", "{{caceNumber}}")
BODY_FONT_SIZE = Pt(11)
PRIORITY_REPAYMENT_FONT_SIZE = Pt(10)
CASE_INFO_FONT_SIZE = Pt(19)
SPECIAL_SUMMARY_SLIDE_CHAR_LIMIT = 850
SPECIAL_SUMMARY_WRAP_WIDTH = 48
SPECIAL_SUMMARY_MAX_WEIGHTED_LINES = 22
NO_TENANTS_TEXT = "조사된 임차인이 없으므로, 낙찰자에게 인수되는 임차권리는 없습니다."

BASE_RIGHT_TYPES = ("근저당권", "근저당", "저당권", "저당", "가압류", "압류", "강제경매", "임의경매")
RIGHT_TYPES = (
    "근저당권",
    "근저당",
    "저당권",
    "저당",
    "가압류",
    "압류",
    "강제경매",
    "임의경매",
    "임차권등기",
    "전세권",
    "지상권",
    "지역권",
    "가처분",
    "가등기",
    "소유권이전청구권가등기",
)


async def generate_rights_certificate(
    request: ReportRequest,
    progress_callback: Optional[Callable] = None,
    task_id: Optional[str] = None,
) -> dict:
    ensure_dirs()

    def emit(step: int, title: str, message: str, status: str = "running", percent: float = 0.0):
        if progress_callback:
            try:
                update = ProgressUpdate(
                    step=step,
                    total_steps=TOTAL_STEPS,
                    title=title,
                    message=message,
                    status=status,
                    percent=percent,
                )
                if asyncio.iscoroutinefunction(progress_callback):
                    try:
                        loop = asyncio.get_event_loop()
                        if loop.is_running():
                            loop.create_task(progress_callback(update))
                        else:
                            asyncio.run(progress_callback(update))
                    except RuntimeError:
                        pass
                else:
                    progress_callback(update)
            except Exception:
                pass
        logger.info(f"[권리분석 보증서 {step}/{TOTAL_STEPS}] {title}: {message}")

    input_url = (request.url or "").strip()
    url = normalize_myauction_detail_url(input_url, request.myauction_id)
    logger.info(f"[권리분석 보증서] 입력 URL: {input_url}")
    logger.info(f"[권리분석 보증서] URL 정규화 결과: {url}")
    logger.info(
        "[권리분석 보증서] URL 정규화 상세: "
        f"view_to_view3={'/view/' in input_url and '/view3/' not in input_url}, "
        f"myauction_id_appended={bool(request.myauction_id and request.myauction_id.strip() and request.myauction_id.strip() in url)}, "
        f"final_url={url}"
    )

    profile_dir = account_profile_dir(request.myauction_id) if request.remember_login else ""
    logger.info(
        f"[권리분석 보증서] Chrome profile 선택: remember_login={request.remember_login}, "
        f"profile_dir={profile_dir or '(임시/비저장 프로필)'}"
    )

    driver = None
    try:
        emit(0, "브라우저 준비", "Chrome 시작 중...", percent=5)
        driver = create_driver(profile_dir=profile_dir, headless=True)

        emit(1, "물건정보 확인", "마이옥션 로그인 중...", percent=15)
        login_myauction(driver, request.myauction_id, request.myauction_pw)

        emit(1, "물건정보 확인", "상세 페이지 접속 및 기본정보 확인 중...", percent=28)
        navigate_with_retry(driver, url, retries=3)
        wait_document_ready(driver, timeout=30)
        time.sleep(3)
        logger.info(
            f"[권리분석 보증서] 상세 페이지 로드: current_url={driver.current_url}, "
            f"title={driver.title or ''}, html_length={len(driver.page_source or '')}"
        )
        log_detail_page_diagnostics(driver, prefix="권리분석 보증서 상세 진입 직후")

        soup = crawler.fetch_soup_from_driver(driver)
        base_data = crawler.parse_myauction_detail(soup, url, driver=driver)
        emit(2, "매각물건명세서 확인", "등기부현황 및 매각물건명세서 확인 중...", percent=42)
        analysis_context = extract_rights_context(soup, driver=driver, task_id=task_id)
        data = {**base_data, **analysis_context}
        data["author_name"] = str(getattr(request, "author_name", "") or "").strip()
        data["author_title"] = str(getattr(request, "author_title", "") or "").strip()
        data["author_phone"] = str(getattr(request, "author_phone", "") or "").strip()

        emit(3, "권리분석 문구", "말소기준권리/임차인/특이사항 문구 구성 중...", percent=55)
        template_data = build_template_data(data)

        safe_task_id = re.sub(r"[^0-9A-Za-z._-]", "_", task_id or datetime.now().strftime("%Y%m%d_%H%M%S"))
        certificate_stem = _rights_certificate_filename_stem(
            template_data.get("caseNumber") or data.get("case_number") or safe_task_id
        )
        pptx_template = Path(getattr(settings, "rights_certificate_pptx_template", ""))

        if pptx_template.exists():
            emit(4, "보증서 템플릿", "PPT 보증서 템플릿에 변수 입력 중...", percent=70)
            pdf_output_path = OUTPUT_DIR / f"{certificate_stem}.pdf"
            pptx_output_path = pdf_output_path.with_suffix(".pptx")
            render_certificate_pptx_template(pptx_template, pptx_output_path, template_data)

            emit(5, "PDF/PPTX 변환", "보증서 PDF 저장 중...", percent=85)
            if export_pptx_to_pdf(pptx_output_path, pdf_output_path):
                output_path = pdf_output_path
            else:
                output_path = pptx_output_path
                emit(5, "PDF/PPTX 변환", "PDF 변환을 사용할 수 없어 PPTX로 저장했습니다.", percent=90)
        else:
            emit(4, "보증서 템플릿", "HTML 보증서 생성 중...", percent=70)
            html_text = render_certificate_template(settings.rights_certificate_template, template_data)

            html_path = OUTPUT_DIR / f"{certificate_stem}.html"
            output_path = OUTPUT_DIR / f"{certificate_stem}.pdf"

            with open(html_path, "w", encoding="utf-8") as f:
                f.write(html_text)

            emit(5, "PDF/PPTX 변환", "보증서 PDF 저장 중...", percent=85)
            render_html_to_pdf(driver, html_path, output_path)

        emit(6, "저장 완료", "권리분석 보증서 생성 완료", status="completed", percent=100)
        return {
            "success": True,
            "output_file": str(output_path),
            "message": "권리분석 보증서 생성이 완료되었습니다.",
            "data": data,
        }
    except Exception as e:
        logger.exception("권리분석 보증서 생성 실패")
        emit(0, "오류", str(e), status="error", percent=0)
        return {"success": False, "message": str(e)}
    finally:
        if driver:
            try:
                driver.quit()
            except Exception:
                pass


def _task_output_path(base_path: str, task_id: str) -> Path:
    root, ext = os.path.splitext(base_path)
    return Path(f"{root}_{task_id}{ext or '.pdf'}")


def _rights_certificate_filename_stem(case_number: str) -> str:
    safe_case_number = _safe_filename_part(case_number) or datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"권리분석_보증서_{safe_case_number}"


def _safe_filename_part(value: str) -> str:
    text = re.sub(r"\s+", "", str(value or "")).strip()
    text = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", text)
    text = text.strip(" ._-")
    if text in ("", "담당자확인필요", "사건번호미확인"):
        return ""
    return text[:80]


def extract_rights_context(soup, driver=None, task_id: Optional[str] = None) -> dict:
    selector_fields = _extract_selector_fields(soup)
    rights_ocr_context = extract_rights_context_by_ocr(driver, task_id=task_id) if driver else {}
    rights = merge_rights(_extract_rights(soup), rights_ocr_context.get("rights") or [])
    ocr_context = extract_tenant_context_by_ocr(driver, task_id=task_id) if driver else {}
    status_survey_context = extract_status_survey_context_by_ocr(driver, task_id=task_id) if driver else {}
    case_document_text = collect_case_document_text(driver) if driver else ""
    tenants = ocr_context.get("tenants") or _extract_tenants(soup)
    dividend_requests = _extract_dividend_requests(soup)
    related_cases = _extract_related_cases(soup)
    auction_applicant_creditors = _extract_auction_applicant_creditors(soup, selector_fields.get("case_number") or "")
    expected_dividend = _extract_expected_dividend(
        soup,
        selector_fields.get("case_number") or "",
        auction_applicant_creditors,
    )
    management_fee = _extract_management_fee(soup)
    market_data = _extract_market_data(soup)
    context = {
        "rights": rights,
        "rights_ocr_text": rights_ocr_context.get("rights_ocr_text", ""),
        "rights_ocr_images": rights_ocr_context.get("rights_ocr_images", []),
        "tenants": tenants,
        "tenant_source": ocr_context.get("tenant_source", ""),
        "tenant_ocr_text": ocr_context.get("tenant_ocr_text", ""),
        "tenant_ocr_images": ocr_context.get("tenant_ocr_images", []),
        "sale_spec_remarks": ocr_context.get("sale_spec_remarks", ""),
        "status_survey_etc": status_survey_context.get("status_survey_etc") or _extract_status_survey_etc_from_text(soup.get_text("\n", strip=True)),
        "status_survey_text": status_survey_context.get("status_survey_text", ""),
        "case_document_text": case_document_text,
        "dividend_requests": dividend_requests,
        "related_cases": related_cases,
        "auction_applicant_creditors": auction_applicant_creditors,
        "expected_dividend": expected_dividend,
        "management_fee": management_fee,
        "market_data": market_data,
    }
    context.update(selector_fields)
    return context


def _extract_selector_fields(soup) -> dict:
    case_number = _css_text(
        soup,
        "#header_detailz > h2 > strong > span",
        "#header_detail2 > h2 > strong > span",
        "#header_detail > h2 > strong > span",
    )
    court = _clean_court_label(_css_text(soup, "#dtl_table > table > tbody > tr:nth-child(2) > td > ul > li:nth-child(1)"))
    case_notice = _extract_case_notice(soup)
    selector_base_right = _extract_selector_base_right(soup)

    fields = {}
    if case_number:
        fields["case_number"] = case_number
    if court:
        fields["court"] = court
    if case_notice:
        fields["case_notice"] = case_notice
    if selector_base_right:
        fields["selector_base_right"] = selector_base_right
    return fields


def _extract_selector_base_right(soup) -> dict:
    row = None
    try:
        row = soup.select_one("#dtl_table > table > tbody > tr:nth-child(3)")
    except Exception:
        row = None
    if not row:
        return {}

    cells = [
        re.sub(r"\s+", " ", cell.get_text(" ", strip=True)).strip()
        for cell in row.find_all(["th", "td"], recursive=False)
    ]
    row_text = re.sub(r"\s+", " ", row.get_text(" ", strip=True)).strip()
    compact = re.sub(r"\s+", "", row_text)
    if not row_text or not any(token in compact for token in ("최선순위", "말소기준", "소멸기준", "근저당", "가압류", "압류", "강제경매", "임의경매")):
        return {}

    date = normalize_date(_first_any_date(row_text) or "")
    right_type = _first_match(row_text, RIGHT_TYPES) or _guess_right_type_after_date(row_text, date)
    creditor = ""

    for cell in cells:
        if not date and _first_any_date(cell):
            date = normalize_date(_first_any_date(cell))
        if not right_type and _first_match(cell, RIGHT_TYPES):
            right_type = _first_match(cell, RIGHT_TYPES)

    red_text = _css_text(soup, "#dtl_table > table > tbody > tr:nth-child(3) > td.dtn_red")
    if red_text:
        if not date and _first_any_date(red_text):
            date = normalize_date(_first_any_date(red_text))
        red_type = _first_match(red_text, RIGHT_TYPES) or _guess_right_type_after_date(red_text, date)
        if red_type:
            right_type = red_type

    for idx in (3, 2, 1):
        if idx < len(cells):
            candidate = _clean_creditor_candidate(cells[idx])
            if candidate and not _first_any_date(candidate) and not _first_match(candidate, RIGHT_TYPES):
                creditor = candidate
                break

    if not date:
        return {}
    return {
        "date": date,
        "type": right_type or "권리종류 확인 필요",
        "creditor": creditor,
        "isBaseRight": True,
        "rawText": row_text,
        "source": "dtl_table_selector",
    }


def _extract_case_notice(soup) -> str:
    text = _css_text(soup, "#dtl_table > table > tbody > tr:nth-child(1) > td")
    if not text:
        return ""
    text = re.sub(r"^(?:주의사항|특이사항|비고)\s*[:：\-]?\s*", "", text).strip()
    return _clean_document_note(text, limit=500)


def _clean_court_label(value: str) -> str:
    text = re.sub(r"\s+", " ", value or "").strip()
    if not text:
        return ""

    court_idx = text.find("법원")
    if court_idx >= 0:
        return text[: court_idx + len("법원")].strip()

    text = re.split(
        r"\s*/\s*|\(\s*\d{2,4}-\d{2,4}\s*\)|\b\d{2,4}-\d{2,4}\b|"
        r"서울특별시|부산광역시|대구광역시|인천광역시|광주광역시|대전광역시|울산광역시|세종특별자치시|"
        r"경기도|강원특별자치도|충청북도|충청남도|전북특별자치도|전라남도|경상북도|경상남도|제주특별자치도",
        text,
        maxsplit=1,
    )[0]
    return text.strip(" /,")


def build_template_data(data: dict) -> dict:
    try:
        cfg = load_config()
    except Exception:
        cfg = {}
    rights = data.get("rights") or []
    tenants = data.get("tenants") or []
    valid_tenants = [tenant for tenant in tenants if not _is_no_tenant_record(tenant)]
    dividend_requests = data.get("dividend_requests") or []
    related_cases = data.get("related_cases") or []
    management_fee = data.get("management_fee") or {}
    market_data = data.get("market_data") or {}

    registry_base_right = find_base_right(rights)
    sale_spec_base_right = data.get("sale_spec_base_right") or {}
    selector_base_right = data.get("selector_base_right") or {}
    if sale_spec_base_right.get("date"):
        base_right = sale_spec_base_right
    elif selector_base_right.get("date"):
        base_right = selector_base_right
    else:
        base_right = registry_base_right
    base_right = enrich_base_right_from_registry(base_right, rights)
    dividend_deadline = data.get("sale_spec_dividend_deadline") or ""
    tenant_texts = analyze_tenants(
        valid_tenants,
        base_right,
        dividend_requests,
        dividend_deadline,
        data.get("address") or "",
    )
    registered_takeover_texts = analyze_registered_takeover_rights(rights, base_right, dividend_requests)
    misc_items = build_misc_items(valid_tenants, management_fee, market_data)
    review_items = build_review_items(rights, valid_tenants, management_fee)
    unpaid_management_fee_text = build_unpaid_management_fee_text(management_fee)
    sale_spec_remarks_text = _polite_optional_note(
        data.get("sale_spec_remarks"),
        "매각물건명세서 비고란에 별도로 기재된 사항은 없습니다.",
    )
    status_survey_etc_text = _polite_optional_note(
        data.get("status_survey_etc"),
        "현황조사서 기타란에 별도로 기재된 사항은 없습니다.",
    )
    case_notice_text = _clean_document_note(data.get("case_notice"), limit=500)
    created_date = format_korean_date(datetime.now())
    case_number = data.get("case_number") or "담당자 확인 필요"
    base_right_date = base_right.get("date") if base_right else "담당자 확인 필요"
    base_right_type = base_right.get("type") if base_right else "등기부 원본 확인"
    base_right_description = data.get("base_right_description") or build_base_right_description(
        base_right,
        rights,
        registered_takeover_texts,
    )
    tenant_analysis_text = build_tenant_analysis_text(
        tenants,
        tenant_texts,
        data.get("tenant_ocr_text") or "",
        data.get("tenant_source") or "",
    )
    if registered_takeover_texts:
        tenant_analysis_text = combine_tenant_and_registered_takeover_texts(
            tenant_analysis_text,
            registered_takeover_texts,
        )
    surplus_description = analyze_surplus(data, rights, base_right, related_cases)
    special_summary_text = build_special_summary_text(
        data,
        rights,
        valid_tenants,
        base_right,
        tenant_texts,
        tenant_analysis_text,
        registered_takeover_texts,
        surplus_description,
        management_fee,
        market_data,
        sale_spec_remarks_text,
        status_survey_etc_text,
        case_notice_text,
    )
    no_tenants = tenant_analysis_text == NO_TENANTS_TEXT
    tenant_analyses = [] if no_tenants else [
        {"description": block.strip()}
        for block in tenant_analysis_text.split("\n\n")
        if block.strip()
    ]

    return {
        "caseNumber": case_number,
        "caceNumber": case_number,
        "court": _clean_court_label(data.get("court")) or "담당자 확인 필요",
        "propertyType": data.get("item_type") or "담당자 확인 필요",
        "propertyOverview": data.get("property_overview") or "",
        "물건개요": data.get("property_overview") or "",
        "appraisalValue": data.get("appraised_price") or "담당자 확인 필요",
        "minBidPrice": data.get("min_price") or "담당자 확인 필요",
        "bidDate": data.get("auction_date") or "담당자 확인 필요",
        "입찰기일": data.get("auction_date") or "담당자 확인 필요",
        "authorName": data.get("author_name") or cfg.get("author_name") or "담당자",
        "authorTitle": data.get("author_title") or cfg.get("author_title") or "",
        "authorPhone": data.get("author_phone") or cfg.get("author_phone") or "",
        "createdDate": created_date,
        "createDate": created_date,
        "baseRightDate": base_right_date,
        "baseRightType": base_right_type,
        "baseRightCreditor": base_right.get("creditor") if base_right else "",
        "baseRightDescription": base_right_description,
        "tenantAnalysisText": tenant_analysis_text,
        "tenantOcrText": data.get("tenant_ocr_text") or "",
        "tenantAnalyses": tenant_analyses,
        "noTenants": no_tenants,
        "surplusDescription": surplus_description,
        "specialSummaryText": special_summary_text,
        "특이사항요약": special_summary_text,
        "miscText": "\n".join(misc_items),
        "miscItems": misc_items,
        "unpaidManagementFeeText": unpaid_management_fee_text,
        "saleSpecRemarksText": sale_spec_remarks_text,
        "statusSurveyEtcText": status_survey_etc_text,
        "caseNoticeText": case_notice_text,
        "caseDocumentText": data.get("case_document_text") or "",
        "미납관리비": unpaid_management_fee_text,
        "매각물건명세서비고": sale_spec_remarks_text,
        "현황조사서기타": status_survey_etc_text,
        "주의사항": case_notice_text,
        "hasUnpaidFee": int(management_fee.get("unpaidAmount") or 0) > 0,
        "reviewText": "\n".join(review_items),
        "reviewItems": review_items,
    }


def find_base_right(rights: list[dict]) -> Optional[dict]:
    marked = [
        r for r in rights
        if r.get("date") and r.get("isBaseRight")
    ]
    marked.sort(key=lambda r: _date_sort_key(r.get("date")))
    if marked:
        return marked[0]

    candidates = [
        r for r in rights
        if r.get("date") and any(t in (r.get("type") or "") for t in BASE_RIGHT_TYPES)
    ]
    candidates.sort(key=lambda r: _date_sort_key(r.get("date")))
    return candidates[0] if candidates else None


def enrich_base_right_from_registry(base_right: Optional[dict], rights: list[dict]) -> Optional[dict]:
    if not base_right:
        return None
    enriched = dict(base_right)
    base_date = enriched.get("date") or ""
    base_type = enriched.get("type") or ""
    for right in rights:
        if not _same_date(right.get("date") or "", base_date):
            continue
        right_type = right.get("type") or ""
        if (
            base_type
            and right_type
            and base_type != "권리종류 확인 필요"
            and base_type not in right_type
            and right_type not in base_type
        ):
            continue
        for key in ("type", "creditor", "amount", "status", "note", "isBaseRight", "rawText"):
            if not enriched.get(key) and right.get(key):
                enriched[key] = right.get(key)
        break
    return enriched


def build_base_right_description(
    base_right: Optional[dict],
    rights: list[dict],
    registered_takeover_texts: Optional[list[str]] = None,
) -> str:
    if not base_right:
        return (
            "등기부현황에서 말소기준권리 확인이 필요합니다. "
            "등기부등본 원본과 매각물건명세서를 기준으로 담당자 최종 확인이 필요합니다."
        )
    creditor = base_right.get("creditor") or "담당자 확인 필요"
    extinguish_text = "말소기준권리 및 이후 모든 권리는 말소됩니다."
    if registered_takeover_texts:
        extinguish_text += " 다만 최선순위 설정일보다 앞선 전세권은 임차권리 인수사항에서 별도 검토합니다."
    else:
        extinguish_text += " 등기부상 낙찰자에게 인수되는 권리는 없습니다."
    return (
        f"최선순위 설정 {base_right.get('date')} 일자 {base_right.get('type')} [{creditor}]\n"
        f"{extinguish_text}"
    )


def analyze_registered_takeover_rights(
    rights: list[dict],
    base_right: Optional[dict],
    dividend_requests: list[dict],
) -> list[str]:
    base_date = (base_right or {}).get("date") or ""
    if not _has_valid_date(base_date):
        return []

    descriptions = []
    for right in rights:
        right_type = right.get("type") or ""
        if "전세권" not in right_type:
            continue
        right_date = right.get("date") or ""
        if not _date_before(right_date, base_date):
            continue
        creditor = right.get("creditor") or "전세권자"
        amount_text = fmt_money(right.get("amount")) if right.get("amount") else "담당자 확인 필요"
        right_label = f"{right_date} {right_type} [{creditor}]"
        if right.get("amount"):
            right_label += f" {amount_text}"

        if _right_has_dividend_request(right, dividend_requests):
            text = (
                f"최선순위 설정보다 앞선 전세권({right_label})이 확인됩니다. "
                "전세권자의 배당요구가 확인되므로 배당 후 소멸 여부를 원본 문서와 대조해 확인해 주시기 바랍니다."
            )
        else:
            text = (
                f"최선순위 설정보다 앞선 전세권({right_label})이 확인됩니다. "
                "전세권자의 배당요구가 확인되지 않으므로 해당 전세권은 낙찰자에게 인수됩니다."
            )
        if text not in descriptions:
            descriptions.append(text)
    return descriptions


def combine_tenant_and_registered_takeover_texts(tenant_text: str, registered_takeover_texts: list[str]) -> str:
    registered_text = "\n\n".join(registered_takeover_texts)
    if not registered_text:
        return tenant_text
    if not tenant_text or tenant_text == NO_TENANTS_TEXT:
        return registered_text
    return f"{tenant_text}\n\n{registered_text}"


def build_special_summary_text(
    data: dict,
    rights: list[dict],
    tenants: list[dict],
    base_right: Optional[dict],
    tenant_texts: list[str],
    tenant_analysis_text: str,
    registered_takeover_texts: list[str],
    surplus_description: str,
    management_fee: dict,
    market_data: dict,
    sale_spec_remarks_text: str,
    status_survey_etc_text: str,
    case_notice_text: str,
) -> str:
    lines: list[str] = []

    lines.append("1) 권리분석 핵심")
    lines.append(f"- 말소기준권리: {_base_right_summary(base_right)}")
    lines.append(f"- 임차내역: {_tenant_count_summary(tenants)}")
    lines.append(f"- 대항력 있는 임차인: {_opposing_tenant_summary(tenant_texts, tenant_analysis_text)}")
    lines.append(f"- 낙찰자 인수 권리: {_takeover_right_summary(tenant_analysis_text, registered_takeover_texts)}")
    lines.append(f"- 후순위 권리: {_junior_rights_summary(rights, base_right)}")

    lines.append("")
    lines.append("2) 가격 분석")
    lines.append(f"- 감정가: {data.get('appraised_price') or '담당자 확인 필요'}")
    lines.append(f"- 최저매각가: {data.get('min_price') or '담당자 확인 필요'}")
    lines.append(f"- 최저가율: {_min_price_rate_text(data)}")
    price_note = _price_analysis_note(data, market_data)
    if price_note:
        lines.append(f"- {price_note}")

    lines.append("")
    lines.append("3) 입찰 전략")
    lines.extend(_bid_strategy_lines(data, tenant_analysis_text, registered_takeover_texts))

    lines.append("")
    lines.append("4) 물건별 특이사항")
    lines.extend(_property_special_issue_lines(
        data,
        rights,
        tenant_texts,
        tenant_analysis_text,
        registered_takeover_texts,
        sale_spec_remarks_text,
        status_survey_etc_text,
        case_notice_text,
    ))
    lines.extend(_bid_check_lines(data, tenants, management_fee))

    return "\n".join(lines)


def _base_right_summary(base_right: Optional[dict]) -> str:
    if not base_right:
        return "등기부현황과 매각물건명세서 원본 확인이 필요합니다."
    date = base_right.get("date") or "일자 확인 필요"
    right_type = base_right.get("type") or "권리종류 확인 필요"
    creditor = base_right.get("creditor") or "권리자 확인 필요"
    return f"{date} 설정된 {right_type} [{creditor}]"


def _tenant_count_summary(tenants: list[dict]) -> str:
    valid_tenants = [tenant for tenant in tenants if not _is_no_tenant_record(tenant)]
    if not valid_tenants:
        return "없음"
    return f"{len(valid_tenants)}건 확인"


def _opposing_tenant_summary(tenant_texts: list[str], tenant_analysis_text: str) -> str:
    combined = "\n".join(tenant_texts + [tenant_analysis_text])
    if _text_has_takeover_tenant(combined):
        return "있음"
    if NO_TENANTS_TEXT in combined or "인수되는 임차권리는 없습니다" in combined:
        return "없음"
    return "원본 문서 확인 필요"


def _takeover_right_summary(tenant_analysis_text: str, registered_takeover_texts: list[str]) -> str:
    if registered_takeover_texts:
        return "선순위 전세권 등 인수 가능 권리 확인"
    if _text_has_takeover_tenant(tenant_analysis_text):
        return "대항력 있는 임차권리 인수 가능성 확인"
    if "인수되는 임차권리는 없습니다" in tenant_analysis_text:
        return "0원 또는 없음으로 분석"
    return "원본 문서 확인 필요"


def _junior_rights_summary(rights: list[dict], base_right: Optional[dict]) -> str:
    base_date = (base_right or {}).get("date") or ""
    if not _has_valid_date(base_date):
        return "말소기준권리 확정 후 확인 필요"
    juniors = [
        right for right in rights
        if _date_after(right.get("date") or "", base_date)
        and any(token in (right.get("type") or "") for token in RIGHT_TYPES)
    ]
    if juniors:
        return "말소기준권리 이후 권리는 낙찰 시 소멸되는 구조로 확인됩니다."
    return "말소기준권리 이후 별도 권리는 확인되지 않습니다."


def _min_price_rate_text(data: dict) -> str:
    if data.get("min_rate"):
        return data["min_rate"]
    appraised = parse_money(data.get("appraised_price"))
    min_price = parse_money(data.get("min_price"))
    if appraised > 0 and min_price > 0:
        return f"{round(min_price / appraised * 100)}%"
    return "담당자 확인 필요"


def _price_analysis_note(data: dict, market_data: dict) -> str:
    if market_data.get("recentDealPrice"):
        return f"최근 실거래가 {fmt_money(market_data.get('recentDealPrice'))} 기준으로 층수·면적·거래시점 차이를 비교해 주시기 바랍니다."
    rate = _min_price_rate_text(data)
    if rate != "담당자 확인 필요":
        return f"현재 최저가는 감정가 대비 {rate} 수준이므로 최근 실거래가와 매물 호가를 함께 확인해 주시기 바랍니다."
    return "감정가와 최저매각가를 기준으로 최근 실거래가와 매물 호가를 함께 확인해 주시기 바랍니다."


def _bid_check_lines(data: dict, tenants: list[dict], management_fee: dict) -> list[str]:
    lines = []
    if tenants:
        lines.append("- 점유 상태: 매각물건명세서 임차내역과 실제 점유자를 대조해 주시기 바랍니다.")
    else:
        lines.append("- 점유 상태: 임차내역이 없더라도 소유자 점유 여부를 현장에서 확인해 주시기 바랍니다.")
    unpaid = int(management_fee.get("unpaidAmount") or 0)
    if unpaid > 0:
        lines.append(f"- 관리비 체납: {build_unpaid_management_fee_text(management_fee)}")
    elif management_fee:
        note = management_fee.get("note") or "미납관리비 없음 또는 미확인"
        lines.append(
            f"- 관리비 체납: {_polite_confirmation(note)} "
            "관리사무소에서 최종 미납 금액 및 승계 범위를 재확인해 주시기 바랍니다."
        )
    else:
        lines.append("- 관리비 체납: 관리사무소에서 미납 금액 및 개월수를 확인해 주시기 바랍니다.")
    return lines


def _bid_strategy_lines(data: dict, tenant_analysis_text: str, registered_takeover_texts: list[str]) -> list[str]:
    lines = []
    min_price = parse_money(data.get("min_price"))
    appraised = parse_money(data.get("appraised_price"))
    rate = _min_price_rate_text(data)
    if min_price > 0 and appraised > 0:
        conservative = int(min_price * 1.03)
        competitive = int((min_price + appraised) / 2)
        ceiling = int(appraised * 0.9)
        if ceiling < competitive:
            ceiling = competitive
        lines.append(f"- 보수적 검토: {fmt_money(conservative)} 전후")
        lines.append(f"- 일반 경쟁 검토: {fmt_money(competitive)} 전후")
        lines.append(f"- 적극 입찰 상한: {fmt_money(ceiling)} 이내에서 관리")
    else:
        lines.append("- 입찰가: 감정가, 최저매각가, 최근 실거래가와 매물 호가를 기준으로 산정해 주시기 바랍니다.")
    if registered_takeover_texts or _text_has_takeover_tenant(tenant_analysis_text):
        lines.append("- 인수권리 가능성이 있으므로 인수금액을 입찰가에서 차감해 검토해 주시기 바랍니다.")
    else:
        lines.append(f"- 권리 리스크가 낮은 물건은 경쟁 가능성이 있으므로 {rate} 최저가율과 시세를 함께 비교해 주시기 바랍니다.")
    lines.append("- 경락잔금대출 한도와 금리는 입찰 전 사전 확인이 필요합니다.")
    return lines


def _property_special_issue_lines(
    data: dict,
    rights: list[dict],
    tenant_texts: list[str],
    tenant_analysis_text: str,
    registered_takeover_texts: list[str],
    sale_spec_remarks_text: str,
    status_survey_etc_text: str,
    case_notice_text: str,
) -> list[str]:
    issues: list[str] = []
    source_text = " ".join(
        str(value or "")
        for value in (
            data.get("item_type"),
            data.get("address"),
            data.get("appraisal_raw"),
            data.get("sale_spec_remarks"),
            data.get("status_survey_etc"),
            data.get("case_notice"),
            data.get("case_document_text"),
            sale_spec_remarks_text,
            status_survey_etc_text,
            case_notice_text,
            tenant_analysis_text,
        )
    )
    if case_notice_text:
        issues.append(f"- 주의사항: {case_notice_text}")
    for line in build_special_issue_lines(source_text):
        issues.append(f"- {line}")
    if any("임차권등기" in (right.get("type") or "") or "임차권등기" in (right.get("rawText") or "") for right in rights):
        issues.append("- 임차권등기: 임차권등기가 확인되므로 실제 점유관계와 배당·인수 여부를 원본 문서로 확인해 주시기 바랍니다.")
    if "대지권미등기" in source_text.replace(" ", ""):
        issues.append(_land_right_unregistered_issue_text())
    if _text_has_priority_repayment("\n".join(tenant_texts + [tenant_analysis_text])):
        issues.append(_small_tenant_priority_issue_text())
    if _text_has_takeover_tenant("\n".join(tenant_texts + [tenant_analysis_text])):
        issues.append("- 대항력 임차인: 최선순위 설정일보다 앞선 임차권리가 확인되므로 보증금 잔액 인수 가능성을 확인해 주시기 바랍니다.")
    if registered_takeover_texts:
        issues.append("- 선순위 전세권: 말소기준권리보다 앞선 전세권은 배당요구 여부에 따라 낙찰자 인수 가능성이 있습니다.")
    return issues


def _small_tenant_priority_issue_text() -> str:
    return (
        "- 소액임차인의 최우선변제: 소액임차인은 확정일자가 늦어 선순위 변제를 받지 못하더라도, "
        "선순위담보권자의 경매신청 등기 전에 대항력을 갖춘 경우 보증금 중 일정액을 다른 담보물권자보다 "
        "우선하여 변제받을 권리가 있습니다(주택임대차보호법 제3조①, 제8조①). "
        "요건은 소액임차인의 범위, 경매신청 등기 전 대항요건, 경매 또는 체납처분 매각, "
        "배당요구 또는 우선권행사 신고, 보증금 중 일정액 보호입니다. "
        "소액임차인의 우선변제 채권은 압류가 금지됩니다(민사집행법 제246조①제6호)."
    )


def _land_right_unregistered_issue_text() -> str:
    return (
        "- 대지권 미등기: 본 물건은 대지권 미등기 상태이므로, 대지사용권의 존재 및 내용, 대지지분, "
        "토지 등기부상 권리관계, 대지권 등기 가능 여부를 계약 전 반드시 확인하여야 합니다. "
        "대지권 미등기는 대지사용권이 없다는 의미로 단정할 수는 없으나, 등기부상 권리관계가 명확히 "
        "공시되지 않은 상태이므로 담보대출 제한, 등기 지연, 추가 비용, 제3자와의 권리분쟁 등 위험이 "
        "발생할 수 있습니다. 매수인은 관련 공부 및 전문가 검토 후 계약을 체결하여야 합니다."
    )


def _text_has_priority_repayment(text: str) -> bool:
    compact = re.sub(r"\s+", "", text or "")
    return "소액임차인" in compact or "최우선변제" in compact


def _text_has_takeover_tenant(text: str) -> bool:
    compact = re.sub(r"\s+", "", text or "")
    if "인수되는임차권리는없" in compact:
        return False
    return any(token in compact for token in ("대항력을갖춘임차인", "잔액은낙찰자에게인수", "임차권리는낙찰자에게인수"))


def build_tenant_analysis_text(
    tenants: list[dict],
    tenant_texts: list[str],
    tenant_ocr_text: str,
    tenant_source: str = "",
) -> str:
    valid_tenants = [tenant for tenant in tenants if not _is_no_tenant_record(tenant)]
    if valid_tenants:
        blocks = []
        for idx, tenant in enumerate(valid_tenants):
            detail = format_tenant_detail(tenant)
            takeover = tenant_texts[idx] if idx < len(tenant_texts) else "인수여부 확인이 필요합니다."
            blocks.append(f"{detail}\n인수여부: {takeover}")
        return "\n\n".join(blocks)
    if tenants or _tenant_ocr_text_indicates_no_tenants(tenant_ocr_text) or tenant_source == "sale_spec_ocr":
        return NO_TENANTS_TEXT
    if tenant_texts:
        return "\n".join(f"인수여부: {text}" for text in tenant_texts)
    if tenant_ocr_text:
        clipped = _clip_text(tenant_ocr_text, 1500)
        return (
            "임차인현황 확인 내용입니다.\n"
            f"{clipped}\n\n"
            "전입일, 확정일자, 배당요구일 및 보증금은 원본 문서와 대조해 담당자 최종 확인이 필요합니다."
        )
    return "임차인 현황은 현황조사서, 매각물건명세서, 전입세대 열람자료를 기준으로 담당자 확인이 필요합니다."


def format_tenant_detail(tenant: dict) -> str:
    return " / ".join(
        [
            f"점유자 성명: {tenant.get('name') or '미확인'}",
            f"점유구분: {tenant.get('occupancyType') or tenant.get('type') or '미확인'}",
            f"보증금: {fmt_money(tenant.get('deposit'))}",
            f"차임: {fmt_money_or_unknown(tenant.get('rent'))}",
            f"전입일: {tenant.get('moveInDate') or '미확인'}",
            f"확정일: {tenant.get('fixedDate') or '미확인'}",
            f"배당요구일: {tenant.get('depositClaimDate') or '미확인'}",
        ]
    )


def format_sale_spec_tenants(tenants: list[dict]) -> str:
    lines = []
    valid_tenants = [tenant for tenant in tenants if not _is_no_tenant_record(tenant)]
    if not valid_tenants:
        return NO_TENANTS_TEXT
    for tenant in valid_tenants:
        lines.append(format_tenant_detail(tenant))
    return "\n".join(lines)


def _is_no_tenant_record(tenant: dict) -> bool:
    name = re.sub(r"\s+", "", str(tenant.get("name") or ""))
    occupancy = re.sub(r"\s+", "", str(tenant.get("occupancyType") or tenant.get("type") or ""))
    has_dates = any(
        _has_valid_date(tenant.get(key) or "")
        for key in ("moveInDate", "fixedDate", "depositClaimDate")
    )
    has_money = parse_money(tenant.get("deposit")) > 0 or parse_money(tenant.get("rent")) > 0
    no_name = _tenant_name_indicates_no_tenants(name)
    no_occupancy = occupancy in ("", "없음", "해당없음", "해당사항없음", "공실", "미상", "미확인")
    return no_name and not has_dates and not has_money and (no_occupancy or "임차인" in occupancy or "점유자" in occupancy)


def _tenant_name_indicates_no_tenants(name: str) -> bool:
    compact = re.sub(r"\s+", "", str(name or ""))
    if compact in ("없음", "해당없음", "해당사항없음", "무", "없습니다"):
        return True
    if "없" not in compact:
        return False
    return any(keyword in compact for keyword in ("조사", "임차", "임대차", "점유", "내역", "관계"))


def _tenant_ocr_text_indicates_no_tenants(text: str) -> bool:
    no_tenant_phrases = (
        "점유자성명없음",
        "점유자없음",
        "임차인없음",
        "조사된임차인없음",
        "해당사항없음",
        "해당없음",
    )
    for raw_line in (text or "").splitlines():
        line = re.sub(r"\s+", "", raw_line)
        if line and any(phrase in line for phrase in no_tenant_phrases):
            return True
        if "없" in line and any(keyword in line for keyword in ("점유자성명", "점유자", "임차인", "임대차관계")):
            return True
    return False


def extract_rights_context_by_ocr(driver, task_id: Optional[str] = None, deadline: Optional[float] = None) -> dict:
    if not driver or not pytesseract:
        return {}

    safe_task_id = re.sub(r"[^0-9A-Za-z._-]", "_", task_id or datetime.now().strftime("%Y%m%d_%H%M%S"))
    image_paths = []
    texts = []
    timed_out = False
    for h3_text, suffix in (
        ("건물 등기부현황", "building_registry"),
        ("토지 등기부현황", "land_registry"),
        ("등기부현황", "registry"),
    ):
        prefix = os.path.join(str(CAPTURE_DIR), f"rights_{suffix}_ocr_{safe_task_id}")
        try:
            captured = capturer.capture_table_split_by_rows(driver, h3_text, prefix, rows_per_page=8, timeout=5)
        except Exception as e:
            logger.info(f"{h3_text} 문서 확인 생략: {e}")
            continue
        image_paths.extend(captured)
        for image_path in captured:
            remaining = (deadline - time.monotonic()) if deadline else 30
            if remaining <= 0:
                timed_out = True
                break
            text = ocr_image_to_text(image_path, timeout_seconds=min(30, max(1, int(remaining))))
            if text:
                texts.append(text)
        if timed_out:
            break

    raw_text = normalize_ocr_text("\n".join(texts))
    if not raw_text:
        return {"rights_ocr_images": image_paths, "_timed_out": timed_out}

    return {
        "rights": parse_rights_from_ocr(raw_text),
        "rights_ocr_text": raw_text,
        "rights_ocr_images": image_paths,
        "_timed_out": timed_out,
    }


def parse_rights_from_ocr(text: str) -> list[dict]:
    rights = []
    for raw_line in (text or "").splitlines():
        line = _strip_registry_sequence_text(raw_line.strip())
        if not line:
            continue
        right_type = _first_match(line, RIGHT_TYPES)
        date = normalize_date(_first_date(line) or "")
        if not right_type or not date:
            continue
        rights.append(
            {
                "seq": len(rights) + 1,
                "date": date,
                "type": right_type,
                "creditor": _guess_right_creditor_from_text(line, right_type),
                "amount": parse_money(line),
                "status": _extract_right_status(line),
                "note": _extract_right_note([], line),
                "isBaseRight": _is_base_right_row(line),
                "rawText": line,
            }
        )
    return _dedupe_by(rights, ("date", "type", "creditor", "amount"))


def _guess_right_creditor_from_text(line: str, right_type: str) -> str:
    cleaned = re.sub(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", " ", line)
    cleaned = cleaned.replace(right_type, " ")
    cleaned = re.sub(r"\d{1,3}(?:,\d{3})+|\d{4,}\s*원?", " ", cleaned)
    cleaned = re.sub(r"(순위|권리자|권리종류|접수|등기|채권최고액|청구금액|말소기준|소유권|이전|설정)", " ", cleaned)
    cleaned = re.sub(r"[:|()\[\],.]", " ", cleaned)
    tokens = [t.strip() for t in cleaned.split() if t.strip()]
    for token in tokens:
        if 1 < len(token) <= 40 and re.search(r"[가-힣A-Za-z]", token):
            return token
    return ""


def extract_tenant_context_by_ocr(driver, task_id: Optional[str] = None, deadline: Optional[float] = None) -> dict:
    if not driver or not pytesseract:
        return {}

    sale_spec_context = extract_sale_spec_tenant_context_by_ocr(driver, task_id=task_id, deadline=deadline)
    if sale_spec_context.get("tenants") or _ocr_text_has_tenant_signals(sale_spec_context.get("tenant_ocr_text", "")):
        return sale_spec_context

    safe_task_id = re.sub(r"[^0-9A-Za-z._-]", "_", task_id or datetime.now().strftime("%Y%m%d_%H%M%S"))
    prefix = os.path.join(str(CAPTURE_DIR), f"rights_tenant_ocr_{safe_task_id}")
    try:
        image_paths = capturer.capture_table_split_by_rows(driver, "임차인현황", prefix, rows_per_page=8, timeout=8)
    except Exception as e:
        logger.warning(f"임차인현황 문서 확인 실패: {e}")
        return {}

    texts = []
    timed_out = False
    for image_path in image_paths:
        remaining = (deadline - time.monotonic()) if deadline else 30
        if remaining <= 0:
            timed_out = True
            break
        text = ocr_image_to_text(image_path, timeout_seconds=min(30, max(1, int(remaining))))
        if text:
            texts.append(text)

    raw_text = normalize_ocr_text("\n".join(texts))
    if not raw_text:
        return {"tenant_ocr_images": image_paths, "_timed_out": timed_out}

    return {
        "tenants": parse_tenants_from_ocr(raw_text),
        "tenant_source": "tenant_status_ocr",
        "tenant_ocr_text": raw_text,
        "tenant_ocr_images": image_paths,
        "_timed_out": timed_out,
    }


def extract_sale_spec_tenant_context_by_ocr(driver, task_id: Optional[str] = None, deadline: Optional[float] = None) -> dict:
    safe_task_id = re.sub(r"[^0-9A-Za-z._-]", "_", task_id or datetime.now().strftime("%Y%m%d_%H%M%S"))
    pdf_text, image_paths = collect_sale_spec_text_and_images(driver, safe_task_id, deadline=deadline)
    if not image_paths and not pdf_text:
        return {}

    pdf_text = normalize_ocr_text(pdf_text)
    texts = []
    timed_out = False
    for image_path in image_paths:
        remaining = (deadline - time.monotonic()) if deadline else 30
        if remaining <= 0:
            timed_out = True
            break
        text = ocr_image_to_text(image_path, timeout_seconds=min(30, max(1, int(remaining))))
        if text:
            texts.append(text)

    raw_text = normalize_ocr_text("\n".join([pdf_text, *texts]))
    if not raw_text:
        return {"tenant_source": "sale_spec_ocr", "tenant_ocr_images": image_paths, "_timed_out": timed_out}

    tenants = parse_sale_spec_tenants_from_pdf_text(pdf_text) or parse_sale_spec_tenants_from_ocr(raw_text)
    sale_spec_context = parse_sale_spec_document_context(pdf_text or raw_text)
    dividend_deadline = sale_spec_context.get("dividendDeadline") or ""
    for tenant in tenants:
        if dividend_deadline and not tenant.get("depositDeadline"):
            tenant["depositDeadline"] = dividend_deadline

    return {
        "tenants": tenants,
        "tenant_source": "sale_spec_ocr",
        "tenant_ocr_text": raw_text,
        "tenant_ocr_images": image_paths,
        "sale_spec_base_right": sale_spec_context.get("baseRight") or {},
        "sale_spec_dividend_deadline": dividend_deadline,
        "sale_spec_remarks": sale_spec_context.get("remarks") or "",
        "_timed_out": timed_out,
    }


def extract_status_survey_context_by_ocr(driver, task_id: Optional[str] = None) -> dict:
    if not driver:
        return {}

    safe_task_id = re.sub(r"[^0-9A-Za-z._-]", "_", task_id or datetime.now().strftime("%Y%m%d_%H%M%S"))
    text = collect_status_survey_text(driver, safe_task_id)
    raw_text = normalize_ocr_text(text)
    if not raw_text:
        return {}
    return {
        "status_survey_text": raw_text,
        "status_survey_etc": _extract_status_survey_etc_from_text(raw_text),
    }


def collect_status_survey_text(driver, safe_task_id: str) -> str:
    base_handle = driver.current_window_handle
    base_url = driver.current_url
    before_handles = set(driver.window_handles)
    opened_handle = ""

    try:
        _open_status_survey_document(driver)
        end = time.time() + 8
        while time.time() < end:
            new_handles = list(set(driver.window_handles) - before_handles)
            if new_handles:
                opened_handle = new_handles[0]
                driver.switch_to.window(opened_handle)
                wait_document_ready(driver, timeout=15)
                break
            time.sleep(0.2)

        pdf_url = _find_current_pdf_url(driver)
        if pdf_url:
            pdf_path = pdf_processor.download_pdf_with_cookies(driver, pdf_url, f"rights_status_survey_{safe_task_id}")
            if pdf_processor.is_valid_pdf(pdf_path):
                return extract_pdf_text(pdf_path)
        try:
            return driver.find_element(By.TAG_NAME, "body").text
        except Exception:
            return ""
    except Exception as e:
        logger.info(f"현황조사서 기타 확인 생략: {e}")
        return ""
    finally:
        try:
            driver.switch_to.default_content()
        except Exception:
            pass
        try:
            if opened_handle and opened_handle in driver.window_handles:
                driver.close()
                driver.switch_to.window(base_handle)
            elif base_handle in driver.window_handles:
                driver.switch_to.window(base_handle)
                if driver.current_url != base_url:
                    driver.get(base_url)
                    wait_document_ready(driver, timeout=15)
        except Exception:
            pass


def _open_status_survey_document(driver) -> None:
    wait = WebDriverWait(driver, 8)
    candidates = [
        (By.PARTIAL_LINK_TEXT, "현황조사서"),
        (By.XPATH, "//div[@id='dtlw_link']//a[contains(normalize-space(.), '현황조사서')]"),
        (By.XPATH, "//a[contains(normalize-space(.), '현황조사서')]"),
    ]
    last_error = None
    for by, value in candidates:
        try:
            el = wait.until(EC.element_to_be_clickable((by, value)))
            safe_click(driver, el)
            time.sleep(1)
            return
        except Exception as e:
            last_error = e
    raise last_error if last_error else RuntimeError("현황조사서 링크를 찾지 못했습니다.")


def _extract_status_survey_etc_from_text(text: str) -> str:
    lines = [
        re.sub(r"\s+", " ", line).strip()
        for line in (text or "").splitlines()
        if re.sub(r"\s+", " ", line).strip()
    ]
    for idx, line in enumerate(lines):
        compact = line.replace(" ", "")
        if compact in ("기타", "기타사항") or "그밖의사항" in compact or "기타사항" in compact:
            same_line = re.sub(r"^\s*(?:기타사항|기타|그 밖의 사항|그밖의사항)\s*[:：]?\s*", "", line)
            collected = [same_line] if same_line and same_line != line else []
            for next_line in lines[idx + 1: idx + 8]:
                if any(next_line.startswith(word) for word in ("첨부", "작성", "사건", "부동산의 표시", "점유관계", "임대차관계")):
                    break
                collected.append(next_line)
            note = _clean_document_note(" ".join(collected))
            if note:
                return note
    return ""


def collect_case_document_text(driver) -> str:
    if not driver:
        return ""

    base_handle = driver.current_window_handle
    base_url = driver.current_url
    before_handles = set(driver.window_handles)
    opened_handle = ""

    try:
        _open_case_document_list(driver)
        end = time.time() + 8
        while time.time() < end:
            new_handles = list(set(driver.window_handles) - before_handles)
            if new_handles:
                opened_handle = new_handles[0]
                driver.switch_to.window(opened_handle)
                wait_document_ready(driver, timeout=15)
                break
            time.sleep(0.2)

        try:
            body_text = driver.find_element(By.TAG_NAME, "body").text
        except Exception:
            body_text = ""
        return _extract_case_document_signal_text(body_text)
    except Exception as e:
        logger.info(f"문건접수 내역 확인 생략: {e}")
        return ""
    finally:
        try:
            driver.switch_to.default_content()
        except Exception:
            pass
        try:
            if opened_handle and opened_handle in driver.window_handles:
                driver.close()
                driver.switch_to.window(base_handle)
            elif base_handle in driver.window_handles:
                driver.switch_to.window(base_handle)
                if driver.current_url != base_url:
                    driver.get(base_url)
                    wait_document_ready(driver, timeout=15)
        except Exception:
            pass


def _open_case_document_list(driver) -> None:
    wait = WebDriverWait(driver, 8)
    candidates = [
        (By.PARTIAL_LINK_TEXT, "문건접수"),
        (By.PARTIAL_LINK_TEXT, "문건"),
        (By.XPATH, "//div[@id='dtlw_link']//a[contains(normalize-space(.), '문건접수') or contains(normalize-space(.), '문건')]"),
        (By.XPATH, "//a[contains(normalize-space(.), '문건접수') or contains(normalize-space(.), '문건')]"),
    ]
    last_error = None
    for by, value in candidates:
        try:
            el = wait.until(EC.element_to_be_clickable((by, value)))
            safe_click(driver, el)
            time.sleep(1)
            return
        except Exception as e:
            last_error = e
    raise last_error if last_error else RuntimeError("문건접수 링크를 찾지 못했습니다.")


def _extract_case_document_signal_text(text: str) -> str:
    lines = [
        re.sub(r"\s+", " ", line).strip()
        for line in (text or "").splitlines()
        if re.sub(r"\s+", " ", line).strip()
    ]
    signal_keywords = (
        "유치권",
        "공사대금",
        "배제신청",
        "유치권배제",
        "유치권 배제",
        "권리신고",
        "권리 신고",
    )
    selected = []
    for idx, line in enumerate(lines):
        compact = re.sub(r"\s+", "", line)
        if not any(re.sub(r"\s+", "", keyword) in compact for keyword in signal_keywords):
            continue
        start = max(0, idx - 1)
        end = min(len(lines), idx + 2)
        selected.extend(lines[start:end])

    if selected:
        return _clean_document_note(" ".join(_dedupe_text_lines(selected)), limit=1200)
    return ""


def _dedupe_text_lines(lines: list[str]) -> list[str]:
    result = []
    seen = set()
    for line in lines:
        key = re.sub(r"\s+", "", line or "")
        if not key or key in seen:
            continue
        seen.add(key)
        result.append(line)
    return result


def collect_sale_spec_text_and_images(driver, safe_task_id: str, deadline: Optional[float] = None) -> tuple[str, list[str]]:
    base_handle = driver.current_window_handle
    base_url = driver.current_url
    before_handles = set(driver.window_handles)
    opened_handle = ""
    image_pattern = str(CAPTURE_DIR / f"rights_sale_spec_{safe_task_id}_{{page}}.png")
    pdf_text = ""

    try:
        _open_sale_spec_document(driver)
        end = time.time() + 8
        while time.time() < end:
            new_handles = list(set(driver.window_handles) - before_handles)
            if new_handles:
                opened_handle = new_handles[0]
                driver.switch_to.window(opened_handle)
                wait_document_ready(driver, timeout=15)
                break
            time.sleep(0.2)

        pdf_url = _find_current_pdf_url(driver)
        if not pdf_url:
            try:
                wait = WebDriverWait(driver, 8)
                click_tab_safe(wait, driver, ["매각물건명세서", "물건명세서"])
                time.sleep(1)
                pdf_url = _find_current_pdf_url(driver)
            except Exception:
                pass

        if pdf_url:
            pdf_path = pdf_processor.download_pdf_with_cookies(driver, pdf_url, f"rights_sale_spec_{safe_task_id}")
        else:
            pdf_path = pdf_processor.print_current_page_to_pdf(driver, f"rights_sale_spec_{safe_task_id}", landscape=True)

        if not pdf_processor.is_valid_pdf(pdf_path):
            return "", _capture_current_page_as_image(driver, safe_task_id)

        pdf_text = extract_pdf_text(pdf_path)
        remaining = (deadline - time.monotonic()) if deadline else 90
        if remaining <= 0:
            return pdf_text, []
        total = pdf_processor.pdf_to_images(
            pdf_path,
            image_pattern,
            dpi=300,
            timeout_seconds=min(90, max(1, int(remaining))),
        )
        return pdf_text, [image_pattern.format(page=i) for i in range(1, total + 1) if os.path.exists(image_pattern.format(page=i))]
    except Exception as e:
        logger.warning(f"매각물건명세서 이미지 생성 실패: {e}")
        return "", []
    finally:
        try:
            driver.switch_to.default_content()
        except Exception:
            pass
        try:
            if opened_handle and opened_handle in driver.window_handles:
                driver.close()
                driver.switch_to.window(base_handle)
            elif base_handle in driver.window_handles:
                driver.switch_to.window(base_handle)
                if driver.current_url != base_url:
                    driver.get(base_url)
                    wait_document_ready(driver, timeout=15)
        except Exception:
            pass


def capture_sale_spec_images(driver, safe_task_id: str) -> list[str]:
    _, image_paths = collect_sale_spec_text_and_images(driver, safe_task_id)
    return image_paths


def _open_sale_spec_document(driver) -> None:
    wait = WebDriverWait(driver, 8)
    candidates = [
        (By.PARTIAL_LINK_TEXT, "매각물건명세서"),
        (By.PARTIAL_LINK_TEXT, "물건명세서"),
        (By.XPATH, "//div[@id='dtlw_link']//a[contains(normalize-space(.), '매각물건명세서') or contains(normalize-space(.), '물건명세서')]"),
        (By.CSS_SELECTOR, "#dtlw_link > ul > li:nth-child(5) > a"),
    ]
    last_error = None
    for by, value in candidates:
        try:
            el = wait.until(EC.element_to_be_clickable((by, value)))
            safe_click(driver, el)
            time.sleep(1)
            return
        except Exception as e:
            last_error = e
    raise last_error if last_error else RuntimeError("매각물건명세서 링크를 찾지 못했습니다.")


def _find_current_pdf_url(driver) -> str:
    try:
        driver.switch_to.default_content()
    except Exception:
        pass

    current_url = driver.current_url or ""
    if ".pdf" in current_url.lower():
        return current_url

    selectors = [
        "iframe#detail_target",
        "iframe[src]",
        "embed[src]",
        "object[data]",
    ]
    for selector in selectors:
        try:
            for el in driver.find_elements(By.CSS_SELECTOR, selector):
                src = el.get_attribute("src") or el.get_attribute("data") or ""
                if src and (".pdf" in src.lower() or "pdf" in src.lower()):
                    return urljoin(current_url, src)
        except Exception:
            continue
    return ""


def _capture_current_page_as_image(driver, safe_task_id: str) -> list[str]:
    out_path = str(CAPTURE_DIR / f"rights_sale_spec_{safe_task_id}_page.png")
    try:
        image = capturer._capture_fullpage_png(driver)
        image.save(out_path, "PNG")
        return [out_path]
    except Exception as e:
        logger.warning(f"매각물건명세서 화면 캡처 실패: {e}")
        return []


def extract_pdf_text(pdf_path: str) -> str:
    try:
        doc = fitz.open(pdf_path)
        try:
            return "\n".join(doc.load_page(i).get_text("text") or "" for i in range(doc.page_count))
        finally:
            doc.close()
    except Exception as e:
        logger.warning(f"PDF 텍스트 추출 실패({pdf_path}): {e}")
        return ""


def ocr_image_to_text(image_path: str, timeout_seconds: int = 30) -> str:
    if not pytesseract:
        return ""
    try:
        img = Image.open(image_path).convert("RGB")
        img = ImageOps.grayscale(img)
        img = ImageEnhance.Contrast(img).enhance(1.8)
        img = ImageOps.autocontrast(img)
        try:
            return pytesseract.image_to_string(img, lang="kor+eng", config="--psm 6", timeout=timeout_seconds)
        except RuntimeError as exc:
            if "timeout" in str(exc).lower():
                logger.warning(f"이미지 OCR {timeout_seconds}초 제한시간 초과({image_path})")
                return ""
            return pytesseract.image_to_string(img, config="--psm 6", timeout=timeout_seconds)
    except Exception as e:
        logger.warning(f"이미지 텍스트 확인 실패({image_path}): {e}")
        return ""


def normalize_ocr_text(text: str) -> str:
    lines = []
    for line in (text or "").splitlines():
        cleaned = re.sub(r"[ \t]+", " ", line).strip()
        if cleaned:
            lines.append(cleaned)
    return "\n".join(lines)


def parse_sale_spec_tenants_from_pdf_text(text: str) -> list[dict]:
    lines = _sale_spec_occupancy_lines(text)
    if not lines:
        return []

    name = _guess_sale_spec_name_from_lines(lines)
    occupancy_type = _guess_sale_spec_occupancy_type_from_lines(lines)
    money_entries = [
        (idx, parse_money(line))
        for idx, line in enumerate(lines)
        if parse_money(line) > 0
    ]
    deposit_idx = -1
    deposit = 0
    rent = 0
    if money_entries:
        deposit_idx, deposit = money_entries[0]
        if len(money_entries) > 1:
            rent = money_entries[1][1]

    date_entries = [
        (idx, normalize_date(match.group(0)))
        for idx, line in enumerate(lines)
        for match in re.finditer(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", line)
    ]
    if deposit_idx >= 0:
        dates = [date for idx, date in date_entries if idx > deposit_idx]
    else:
        dates = [date for _, date in date_entries]

    tenant = {
        "name": name,
        "occupancyType": occupancy_type,
        "type": occupancy_type,
        "deposit": deposit,
        "rent": rent,
        "moveInDate": dates[0] if len(dates) > 0 else "",
        "fixedDate": dates[1] if len(dates) > 1 else "",
        "depositClaimDate": dates[2] if len(dates) > 2 else "",
        "depositDeadline": "",
        "isHUG": "주택도시보증공사" in text or "HUG" in text.upper(),
        "isVacant": "공실" in text,
    }
    tenant.update(_extract_increase_context(text))

    if not tenant["name"] and not tenant["deposit"] and not tenant["moveInDate"]:
        return []
    tenant["name"] = tenant["name"] or "미확인 점유자"
    tenant["occupancyType"] = tenant["occupancyType"] or "미확인"
    tenant["type"] = tenant["occupancyType"]
    return [tenant]


def parse_sale_spec_document_context(text: str) -> dict:
    lines = [
        re.sub(r"\s+", " ", line).strip()
        for line in (text or "").splitlines()
        if re.sub(r"\s+", " ", line).strip()
    ]
    context = {}
    base_right = {}
    dividend_deadline = ""

    for idx, line in enumerate(lines):
        compact = line.replace(" ", "")
        if "최선순위" in compact or (compact == "설정" and idx > 0 and "최선순위" in lines[idx - 1].replace(" ", "")):
            window = "\n".join(lines[idx: idx + 6])
            date = _first_any_date(window)
            if date:
                base_right = {
                    "date": date,
                    "type": _first_match(window, RIGHT_TYPES) or _guess_right_type_after_date(window, date),
                    "creditor": "",
                }
                break

    for idx, line in enumerate(lines):
        compact = line.replace(" ", "")
        if "배당요구종기" in compact:
            window = "\n".join(lines[idx: idx + 4])
            dividend_deadline = _first_any_date(window)
            if dividend_deadline:
                break

    if base_right:
        context["baseRight"] = base_right
    if dividend_deadline:
        context["dividendDeadline"] = dividend_deadline
    remarks = _extract_sale_spec_remarks_from_lines(lines)
    if remarks:
        context["remarks"] = remarks
    return context


def _extract_sale_spec_remarks_from_lines(lines: list[str]) -> str:
    start_idx = -1
    for idx, line in enumerate(lines):
        compact = line.replace(" ", "")
        if "비고란" in compact or compact in ("비고", "비고사항"):
            start_idx = idx
            break
        if compact.startswith("비고") and len(compact) > 2:
            return _clean_document_note(re.sub(r"^\s*비고\s*[:：]?\s*", "", line))

    if start_idx < 0:
        return ""

    stop_words = ("사건", "작성", "담임법관", "부동산의 표시", "최선순위", "배당요구종기")
    collected = []
    for line in lines[start_idx + 1: start_idx + 8]:
        if any(line.startswith(word) for word in stop_words):
            break
        if line in ("없음", "해당없음", "해당 사항 없음"):
            return ""
        if line and not line.startswith("※"):
            collected.append(line)
    return _clean_document_note(" ".join(collected))


def _clean_document_note(value: str, limit: int = 260) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip(" /,|")
    if text in ("", "없음", "해당없음", "해당 사항 없음", "미상"):
        return ""
    return _clip_text(text, limit)


def _guess_right_type_after_date(text: str, date: str) -> str:
    after = text
    if date:
        numbers = re.findall(r"\d+", date)
        if len(numbers) >= 3:
            pattern = r"\s*[.\-/년월일]*\s*".join(map(re.escape, numbers[:3]))
            after = re.split(pattern, text, maxsplit=1)[-1]
    match = re.search(r"[가-힣A-Za-z]+", after)
    return match.group(0) if match else ""


def _extract_increase_context(text: str) -> dict:
    if not re.search(r"증액|증가|추가", text or ""):
        return {}
    lines = [
        re.sub(r"\s+", " ", line).strip()
        for line in (text or "").splitlines()
        if re.sub(r"\s+", " ", line).strip()
    ]
    for idx, line in enumerate(lines):
        if not re.search(r"증액|증가|추가", line):
            continue
        window = "\n".join(lines[max(0, idx - 2): idx + 4])
        fixed_date = _extract_date_after_keywords(window, ("확정일자", "확정일", "확정")) or _first_any_date(window)
        return {
            "hasIncrease": True,
            "increaseFixedDate": fixed_date,
        }
    return {"hasIncrease": True, "increaseFixedDate": ""}


def _sale_spec_occupancy_lines(text: str) -> list[str]:
    lines = [
        re.sub(r"\s+", " ", line).strip()
        for line in (text or "").splitlines()
        if re.sub(r"\s+", " ", line).strip()
    ]
    if not lines:
        return []

    header_idx = -1
    for idx, line in enumerate(lines):
        if "배당요구일자" in line or "(배당요구일자)" in line:
            header_idx = idx
    if header_idx < 0:
        return []

    stop_words = (
        "등기된 부동산",
        "매각에 따라",
        "부동산의 표시",
        "※1:",
        "비고란",
        "사건",
        "작성",
        "담임법관",
    )
    result = []
    for line in lines[header_idx + 1:]:
        if any(line.startswith(word) for word in stop_words):
            break
        result.append(line)
    return result


def _guess_sale_spec_name_from_lines(lines: list[str]) -> str:
    excluded = {
        "현황조사",
        "권리신고",
        "주거",
        "상가",
        "임차인",
        "전부",
        "일부",
        "미상",
        "없음",
        "권원",
    }
    for line in lines:
        if line in excluded:
            continue
        if _first_date(line) or parse_money(line):
            continue
        if any(word in line for word in ("점유", "보증금", "차임", "전입", "확정", "배당")):
            continue
        if 1 < len(line) <= 30 and re.search(r"[가-힣A-Za-z]", line):
            return line
    return ""


def _guess_sale_spec_occupancy_type_from_lines(lines: list[str]) -> str:
    for idx, line in enumerate(lines):
        compact = line.replace(" ", "")
        if "주거임차인" in compact:
            return "주거 임차인"
        if "상가임차인" in compact:
            return "상가 임차인"
        if line in ("주거", "상가") and idx + 1 < len(lines) and "임차인" in lines[idx + 1]:
            return f"{line} 임차인"
    for line in lines:
        if "임차인" in line:
            return "임차인"
    return ""


def parse_sale_spec_tenants_from_ocr(text: str) -> list[dict]:
    tenants = []
    for raw_line in (text or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if not _looks_like_sale_spec_tenant_line(line):
            continue
        if _is_tenant_header_line(line):
            continue

        dates = [normalize_date(d) for d in re.findall(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", line)]
        tenant = {
            "name": _extract_labeled_value(line, ("점유자 성명", "점유자", "성명")) or _guess_sale_spec_tenant_name(line),
            "occupancyType": _extract_labeled_value(line, ("점유구분", "점유관계", "점유 부분", "점유")) or _guess_occupancy_type(line),
            "type": _guess_occupancy_type(line),
            "deposit": _extract_money_after_keywords(line, ("보증금", "임대차보증금", "전세금")),
            "rent": _extract_money_after_keywords(line, ("차임", "월차임", "월세")),
            "moveInDate": _extract_date_after_keywords(line, ("전입일", "전입일자", "전입")) or (dates[0] if len(dates) > 0 else ""),
            "fixedDate": _extract_date_after_keywords(line, ("확정일", "확정일자", "확정")) or (dates[1] if len(dates) > 1 else ""),
            "depositClaimDate": _extract_date_after_keywords(line, ("배당요구일", "배당요구일자", "배당요구")) or (dates[2] if len(dates) > 2 else ""),
            "depositDeadline": "",
            "isHUG": "주택도시보증공사" in line or "HUG" in line.upper(),
            "isVacant": "공실" in line,
        }
        tenant.update(_extract_increase_context(line))

        if not tenant["name"] and not tenant["moveInDate"] and not tenant["deposit"]:
            continue
        tenant["name"] = tenant["name"] or "미확인 점유자"
        tenant["occupancyType"] = tenant["occupancyType"] or "미확인"
        tenants.append(tenant)

    return _dedupe_by(tenants, ("name", "occupancyType", "moveInDate", "deposit", "rent"))


def _looks_like_sale_spec_tenant_line(line: str) -> bool:
    keywords = ("점유", "임차", "보증금", "차임", "월세", "전입", "확정", "배당요구")
    return any(k in line for k in keywords) or bool(re.search(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", line))


def _ocr_text_has_tenant_signals(text: str) -> bool:
    if not text:
        return False
    keywords = ("매각물건명세서", "점유자", "점유구분", "보증금", "차임", "전입일", "확정일", "배당요구")
    return sum(1 for keyword in keywords if keyword in text) >= 2


def _is_tenant_header_line(line: str) -> bool:
    header_keywords = ("점유자 성명", "점유구분", "보증금", "차임", "전입일", "확정일", "배당요구일")
    return sum(1 for keyword in header_keywords if keyword in line) >= 4 and not re.search(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", line)


def _extract_labeled_value(line: str, labels: tuple[str, ...]) -> str:
    for label in labels:
        pattern = re.compile(
            re.escape(label) + r"\s*[:：]?\s*([^\n/|,]+)",
        )
        match = pattern.search(line)
        if not match:
            continue
        value = match.group(1).strip()
        value = re.split(r"\s+(?:점유구분|보증금|차임|월세|전입|확정|배당요구)", value)[0].strip()
        value = re.sub(r"[:：|,]+$", "", value).strip()
        if value and value not in labels:
            return value
    return ""


def _extract_date_after_keywords(line: str, labels: tuple[str, ...]) -> str:
    date_pattern = r"(\d{4}\s*(?:[.\-/년])\s*\d{1,2}\s*(?:[.\-/월])\s*\d{1,2}\s*\.?\s*일?)"
    for label in labels:
        match = re.search(re.escape(label) + r"\s*[:：]?\s*" + date_pattern, line)
        if match:
            return normalize_date(match.group(1))
    return ""


def _extract_money_after_keywords(line: str, labels: tuple[str, ...]) -> int:
    for label in labels:
        match = re.search(
            re.escape(label) + r"\s*[:：]?\s*((?:\d{1,3}(?:,\d{3})+|\d{4,})(?:\s*원)?)",
            line,
        )
        if match:
            return parse_money(match.group(1))
    return 0


def _guess_occupancy_type(line: str) -> str:
    for keyword in ("주거임차인", "상가임차인", "임차인", "소유자", "채무자", "점유자", "공실", "미상"):
        if keyword in line:
            return keyword
    return ""


def _guess_sale_spec_tenant_name(line: str) -> str:
    cleaned = re.sub(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", " ", line)
    cleaned = re.sub(r"\d{1,3}(?:,\d{3})+|\d{4,}\s*원?", " ", cleaned)
    cleaned = re.sub(
        r"(점유자\s*성명|점유자|성명|점유구분|점유관계|점유|보증금|임대차보증금|전세금|차임|월차임|월세|전입일자?|전입|확정일자?|확정|배당요구일자?|배당요구|임차인|소유자|채무자|공실|없음|미상)",
        " ",
        cleaned,
    )
    cleaned = re.sub(r"[:：|()\[\],./]", " ", cleaned)
    tokens = [t.strip() for t in cleaned.split() if t.strip()]
    for token in tokens:
        if 1 < len(token) <= 20 and re.search(r"[가-힣A-Za-z]", token):
            return token
    return ""


def parse_tenants_from_ocr(text: str) -> list[dict]:
    tenants = []
    current_deadline = ""

    for raw_line in (text or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if "배당요구종기" in line or "종기" in line:
            current_deadline = normalize_date(_first_date(line) or "")
        if not ("임차" in line or "전입" in line or "확정" in line or "보증" in line or _first_date(line)):
            continue
        if any(header in line for header in ("성명", "점유", "전입", "확정")) and not re.search(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", line):
            continue

        dates = [normalize_date(d) for d in re.findall(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", line)]
        deposit = parse_money(line)
        name = _guess_ocr_tenant_name(line)

        if not dates and not deposit and not name:
            continue
        tenants.append(
            {
                "name": name or "미확인 임차인",
                "type": "주택도시보증공사" if "주택도시보증공사" in line or "HUG" in line.upper() else "",
                "moveInDate": dates[0] if len(dates) > 0 else "",
                "fixedDate": dates[1] if len(dates) > 1 else "",
                "depositClaimDate": dates[2] if len(dates) > 2 else "",
                "depositDeadline": current_deadline,
                "deposit": deposit,
                "rent": 0,
                "isHUG": "주택도시보증공사" in line or "HUG" in line.upper(),
                "isVacant": "공실" in line,
            }
        )

    return _dedupe_by(tenants, ("name", "moveInDate", "deposit"))


def _guess_ocr_tenant_name(line: str) -> str:
    cleaned = re.sub(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", " ", line)
    cleaned = re.sub(r"\d{1,3}(?:,\d{3})+|\d{4,}\s*원?", " ", cleaned)
    cleaned = re.sub(r"(임차인|전입일자?|전입|확정일자?|확정|배당요구일자?|배당요구|보증금|점유|월세|차임|없음|미상)", " ", cleaned)
    cleaned = re.sub(r"[:|()\[\],.]", " ", cleaned)
    tokens = [t.strip() for t in cleaned.split() if t.strip()]
    for token in tokens:
        if 1 < len(token) <= 20 and re.search(r"[가-힣A-Za-z]", token):
            return token
    return ""


def analyze_tenants(
    tenants: list[dict],
    base_right: Optional[dict],
    dividend_requests: list[dict],
    dividend_deadline: str = "",
    property_address: str = "",
) -> list[str]:
    descriptions = []
    for tenant in tenants:
        move_in = tenant.get("moveInDate") or ""
        fixed_date = tenant.get("fixedDate") or ""
        request_date = tenant.get("depositClaimDate") or ""
        deadline = dividend_deadline or tenant.get("depositDeadline") or ""
        if not request_date:
            req = _find_dividend_request(tenant.get("name") or "", dividend_requests)
            request_date = (req or {}).get("requestDate") or ""
            deadline = deadline or (req or {}).get("deadline") or ""

        base_date = (base_right or {}).get("date") or ""
        case = _tenant_takeover_case(move_in, fixed_date, request_date, base_date, deadline)
        text = case["text"]
        priority_text = _tenant_priority_repayment_text(
            tenant,
            base_date,
            property_address,
            bool(case.get("takeover")),
        )
        if priority_text:
            text = f"{text}\n{priority_text}"
        increase_text = _tenant_increase_takeover_text(tenant, base_date)
        if increase_text:
            text = f"{text}\n{increase_text}"
        if text not in descriptions:
            descriptions.append(text)
    return descriptions


def _tenant_takeover_case_text(move_in: str, fixed_date: str, request_date: str, base_date: str, deadline: str) -> str:
    return _tenant_takeover_case(move_in, fixed_date, request_date, base_date, deadline)["text"]


def _tenant_takeover_case(move_in: str, fixed_date: str, request_date: str, base_date: str, deadline: str) -> dict:
    no_takeover = "최선순위 설정보다 앞서 대항력을 갖춘 임차인이 없으므로,낙찰자에게 인수되는 임차권리는 없습니다."
    if not base_date or not _has_valid_date(move_in) or not _date_before(move_in, base_date):
        return {"text": no_takeover, "takeover": False}

    fixed_before_base = _date_before(fixed_date, base_date)
    request_on_time = _date_on_or_before(request_date, deadline)

    if fixed_before_base and request_on_time:
        return {
            "text": "최선순위 설정 보다 앞선 대항력을 갖춘 임차인이 있으므로,순위 배당 시 배당 받지 못하는 잔액이 있다면, 잔액은 낙찰자에게 인수됩니다.",
            "takeover": True,
        }

    if fixed_before_base and _date_after(request_date, deadline):
        return {
            "text": "최선순위 설정 보다 앞선 대항력을 갖춘 임차인이 있습니다. 배당 받지 못한 잔액은 낙찰자에게 인수 됩니다.",
            "takeover": True,
        }

    if not _has_valid_date(fixed_date) and request_on_time:
        return {
            "text": "최선순위 설정 보다 앞선 대항력을 갖춘 임차인이 있으므로,순위 배당 시 배당 받지 못하는 잔액이 있다면, 잔액은 낙찰자에게 인수됩니다.",
            "takeover": True,
        }

    return {
        "text": "최선순위 설정 보다 앞선 대항력을 갖춘 임차인이 있으므로, 임차권리는 낙찰자에게 인수 됩니다.",
        "takeover": True,
    }


def _tenant_priority_repayment_text(tenant: dict, base_date: str, property_address: str, residual_takeover: bool) -> str:
    occupancy = f"{tenant.get('occupancyType') or ''} {tenant.get('type') or ''}"
    if "상가" in occupancy or "사업자" in occupancy:
        return ""
    deposit = parse_money(tenant.get("deposit"))
    if deposit <= 0:
        return ""
    priority = _housing_lease_priority_repayment(base_date, property_address, deposit)
    if not priority:
        return ""

    priority_amount = min(deposit, priority["priority_amount"])
    remaining = max(deposit - priority_amount, 0)
    text = (
        f"주택임대차보호법상 소액임차인 최우선변제 대상입니다. "
        f"최우선변제금액 {fmt_money(priority_amount)}은 최우선 변제됩니다."
    )
    if remaining <= 0:
        return f"{text} 남은 잔존 금액은 없습니다."
    if residual_takeover:
        return f"{text} 잔존 금액 {fmt_money(remaining)}은 낙찰자에게 인수됩니다."
    return f"{text} 잔존 금액 {fmt_money(remaining)}은 낙찰자에게 인수되지 않습니다."


HOUSING_LEASE_PRIORITY_RULES = [
    ("1984.06.14", {
        "special": (3_000_000, 3_000_000),
        "other": (2_000_000, 2_000_000),
    }),
    ("1987.12.01", {
        "special": (5_000_000, 5_000_000),
        "other": (4_000_000, 4_000_000),
    }),
    ("1990.02.19", {
        "special": (20_000_000, 7_000_000),
        "other": (15_000_000, 5_000_000),
    }),
    ("1995.10.19", {
        "metropolitan": (30_000_000, 12_000_000),
        "other": (20_000_000, 8_000_000),
    }),
    ("2001.09.15", {
        "overcrowding": (40_000_000, 16_000_000),
        "metro_except": (35_000_000, 14_000_000),
        "other": (30_000_000, 12_000_000),
    }),
    ("2008.08.21", {
        "overcrowding": (60_000_000, 20_000_000),
        "metro_except": (50_000_000, 17_000_000),
        "other": (40_000_000, 14_000_000),
    }),
    ("2010.07.26", {
        "seoul": (75_000_000, 25_000_000),
        "overcrowding": (65_000_000, 22_000_000),
        "metro_plus": (55_000_000, 19_000_000),
        "other": (40_000_000, 14_000_000),
    }),
    ("2014.01.01", {
        "seoul": (95_000_000, 32_000_000),
        "overcrowding": (80_000_000, 27_000_000),
        "metro_plus": (60_000_000, 20_000_000),
        "other": (45_000_000, 15_000_000),
    }),
    ("2016.03.31", {
        "seoul": (100_000_000, 34_000_000),
        "overcrowding": (80_000_000, 27_000_000),
        "metro_plus": (60_000_000, 20_000_000),
        "other": (50_000_000, 17_000_000),
    }),
    ("2018.09.18", {
        "seoul": (110_000_000, 37_000_000),
        "overcrowding": (100_000_000, 34_000_000),
        "metro_plus": (60_000_000, 20_000_000),
        "other": (50_000_000, 17_000_000),
    }),
    ("2021.05.11", {
        "seoul": (150_000_000, 50_000_000),
        "overcrowding_plus": (130_000_000, 43_000_000),
        "metro_plus": (70_000_000, 23_000_000),
        "other": (60_000_000, 20_000_000),
    }),
    ("2023.02.21", {
        "seoul": (165_000_000, 55_000_000),
        "overcrowding_plus": (145_000_000, 48_000_000),
        "metro_plus": (85_000_000, 28_000_000),
        "other": (75_000_000, 25_000_000),
    }),
]


def _housing_lease_priority_repayment(base_date: str, address: str, deposit: int) -> dict:
    rule = _housing_lease_priority_rule(base_date)
    if not rule:
        return {}
    category = _housing_lease_region_category(address, rule["start"])
    deposit_limit, priority_amount = rule["limits"].get(category) or rule["limits"]["other"]
    if deposit > deposit_limit:
        return {}
    return {
        "category": category,
        "deposit_limit": deposit_limit,
        "priority_amount": priority_amount,
    }


def _housing_lease_priority_rule(base_date: str) -> dict:
    if not _has_valid_date(base_date):
        return {}
    selected = None
    for start, limits in HOUSING_LEASE_PRIORITY_RULES:
        if _date_sort_key(start) <= _date_sort_key(base_date):
            selected = {"start": start, "limits": limits}
        else:
            break
    return selected or {}


def _housing_lease_region_category(address: str, rule_start: str) -> str:
    text = re.sub(r"\s+", "", address or "")
    if _date_sort_key(rule_start) < _date_sort_key("1995.10.19"):
        return "special" if ("특별시" in text or "광역시" in text or "직할시" in text) else "other"
    if _date_sort_key(rule_start) < _date_sort_key("2001.09.15"):
        return "metropolitan" if _is_special_or_metropolitan_non_county(text) else "other"
    if _date_sort_key(rule_start) < _date_sort_key("2010.07.26"):
        if _is_overcrowding_area(text, include_seoul=True):
            return "overcrowding"
        if _is_metropolitan_except_incheon_and_county(text):
            return "metro_except"
        return "other"
    if _date_sort_key(rule_start) < _date_sort_key("2021.05.11"):
        if "서울특별시" in text or text.startswith("서울"):
            return "seoul"
        if _is_overcrowding_area(text, include_seoul=False):
            return "overcrowding"
        if _is_metropolitan_plus_city(text, include_sejong=_date_sort_key(rule_start) >= _date_sort_key("2016.03.31")):
            return "metro_plus"
        return "other"
    if "서울특별시" in text or text.startswith("서울"):
        return "seoul"
    if _is_overcrowding_area(text, include_seoul=False) or _contains_any(text, ("세종", "용인", "화성", "김포")):
        return "overcrowding_plus"
    if _is_metropolitan_city_non_county(text) or _contains_any(text, ("안산", "광주", "파주", "이천", "평택")):
        return "metro_plus"
    return "other"


def _is_special_or_metropolitan_non_county(address: str) -> bool:
    if "군" in address:
        return False
    return "특별시" in address or "광역시" in address or "직할시" in address


def _is_metropolitan_city_non_county(address: str) -> bool:
    return "광역시" in address and "군" not in address


def _is_metropolitan_except_incheon_and_county(address: str) -> bool:
    return _is_metropolitan_city_non_county(address) and "인천광역시" not in address


def _is_metropolitan_plus_city(address: str, include_sejong: bool = False) -> bool:
    if _is_metropolitan_city_non_county(address) and not _is_overcrowding_area(address, include_seoul=True):
        return True
    cities = ("안산", "용인", "김포", "광주")
    if _contains_any(address, cities):
        return True
    return include_sejong and "세종" in address


def _is_overcrowding_area(address: str, include_seoul: bool) -> bool:
    if include_seoul and ("서울특별시" in address or address.startswith("서울")):
        return True
    return _contains_any(
        address,
        (
            "인천광역시",
            "의정부",
            "구리",
            "남양주",
            "하남",
            "고양",
            "수원",
            "성남",
            "안양",
            "부천",
            "광명",
            "과천",
            "의왕",
            "군포",
            "시흥",
        ),
    )


def _contains_any(text: str, needles: tuple[str, ...]) -> bool:
    return any(needle in text for needle in needles)


def _tenant_increase_takeover_text(tenant: dict, base_date: str) -> str:
    if not tenant.get("hasIncrease"):
        return ""
    increase_fixed_date = tenant.get("increaseFixedDate") or ""
    if not _has_valid_date(base_date) or not _has_valid_date(increase_fixed_date):
        return ""
    if _date_before(increase_fixed_date, base_date):
        return "증액분도 인수됩니다."
    return "증액분은 인수 되지 않습니다."


def calculate_senior_debt_total(rights: list[dict], base_right: Optional[dict]) -> tuple[int, bool]:
    base_date = (base_right or {}).get("date") or ""
    if not _has_valid_date(base_date):
        return 0, False

    total = 0
    for right in rights or []:
        right_date = right.get("date") or ""
        amount = parse_money(right.get("amount"))
        if amount <= 0 or not _date_before(right_date, base_date):
            continue
        total += amount
    return total, True


def calculate_surplus_basis(data: dict, rights: list[dict], base_right: Optional[dict]) -> dict:
    base_date = (base_right or {}).get("date") or ""
    min_bid = parse_money(data.get("min_price"))
    execution_values = forced_execution_estimator.build_eviction_cost_values(data)
    execution_cost = int(execution_values.get("normal_execution_cost") or 0)
    senior_debt_total, senior_debt_available = calculate_senior_debt_total(rights, base_right)
    remainder = min_bid - execution_cost - senior_debt_total
    can_calculate = min_bid > 0 and execution_cost >= 0 and senior_debt_available
    return {
        "min_bid": min_bid,
        "execution_cost": execution_cost,
        "senior_debt_total": senior_debt_total,
        "remainder": remainder,
        "can_calculate": can_calculate,
        "base_date": base_date,
    }


def fmt_formula_money(value) -> str:
    try:
        amount = int(value)
    except (TypeError, ValueError):
        amount = parse_money(value)
    return f"{amount:,}원"


def fmt_signed_money(value) -> str:
    try:
        amount = int(value)
    except (TypeError, ValueError):
        amount = parse_money(value)
    return f"{amount:,}원"


def build_no_surplus_judgment_text(data: dict, rights: list[dict], base_right: Optional[dict]) -> str:
    basis = calculate_surplus_basis(data, rights, base_right)
    if not basis.get("can_calculate"):
        return (
            "최저매각가격, 집행비용 또는 선순위 채권총액 확인이 필요하여 "
            "무잉여 가능성을 확정하지 못했습니다."
        )

    if basis["remainder"] <= 0:
        return "경매신청채권자에게 배당될 금액이 남지 않는 것으로 판단되어 무잉여 가능성이 있습니다."
    return "경매신청채권자에게 배당될 금액이 남는 것으로 판단되어 무잉여 가능성은 낮습니다."


def analyze_surplus(data: dict, rights: list[dict], base_right: Optional[dict], related_cases: list[dict]) -> str:
    min_bid = parse_money(data.get("min_price"))
    appraised = parse_money(data.get("appraised_price"))
    total_debt = sum(int(r.get("amount") or 0) for r in rights if r.get("amount"))
    expected_dividend = data.get("expected_dividend") or {}
    applicant_creditors = data.get("auction_applicant_creditors") or []
    case_text = ", ".join(
        f"{c.get('type', '관련사건')} {c.get('caseNumber', '')}".strip()
        for c in related_cases
    )

    expected_amount = int(expected_dividend.get("auctionApplicantDividendAmount") or 0)
    if _has_duplicate_auction_case(related_cases):
        no_surplus_text = "중복경매 신청 사건이 확인되므로, 단순 무잉여를 이유로 한 절차 기각 가능성은 낮습니다."
    elif expected_dividend.get("auctionApplicantDividendFound") and expected_amount > 0:
        no_surplus_text = "경매신청채권자는 배당을 받을 수 있으므로 무잉여 가능성은 없습니다."
    elif _base_right_creditor_is_auction_applicant(base_right, applicant_creditors):
        no_surplus_text = "최선순위 설정권자와 경매신청채권자가 동일하여 우선 배당 가능성이 높으므로 무잉여 가능성은 없습니다."
    else:
        no_surplus_text = build_no_surplus_judgment_text(data, rights, base_right)

    if not appraised or not rights:
        withdrawal_text = (
            "감정가 또는 채권 내역 확인이 필요하여 경매취하 가능성을 확정하지 못했습니다. "
            "등기부 채권 총액과 감정가 대비 비율을 담당자가 확인해야 합니다."
        )
    else:
        debt_rate = total_debt / appraised
        debt_rate_text = _format_percent(debt_rate * 100)
        if debt_rate < 0.7:
            withdrawal_text = (
                f"확인된 채권 총액 ({fmt_money(total_debt)})은 감정가({fmt_money(appraised)}) 대비 "
                f"{debt_rate_text}로 70% 미만이므로 취하 가능성이 있습니다."
            )
        else:
            withdrawal_text = (
                f"확인된 채권 총액 ({fmt_money(total_debt)})은 감정가({fmt_money(appraised)}) 대비 "
                f"{debt_rate_text}로 70% 이상이므로 취하 가능성은 낮습니다."
            )

    parts = [
        f"무잉여 가능성: {no_surplus_text}",
        f"취하 가능성: {withdrawal_text}",
    ]
    if case_text:
        parts.append(f"관련 사건: {case_text}.")
    return "\n".join(parts)


def _format_percent(value: float) -> str:
    if abs(value - round(value)) < 0.05:
        return f"{round(value)}%"
    return f"{value:.1f}%"


def _has_duplicate_auction_case(related_cases: list[dict]) -> bool:
    for case in related_cases or []:
        case_type = str(case.get("type") or "")
        if "중복" in case_type:
            return True
    return False


def _base_right_creditor_is_auction_applicant(base_right: Optional[dict], applicant_creditors: list[str]) -> bool:
    creditor = (base_right or {}).get("creditor") or ""
    normalized_applicants = [
        _normalize_creditor_name(name)
        for name in applicant_creditors or []
        if _normalize_creditor_name(name)
    ]
    return _creditor_matches_any(creditor, normalized_applicants)


def build_misc_items(tenants: list[dict], management_fee: dict, market_data: dict) -> list[str]:
    items = []
    unpaid = int(management_fee.get("unpaidAmount") or 0)
    if management_fee:
        if unpaid > 0:
            items.append(build_unpaid_management_fee_text(management_fee))
        else:
            note = management_fee.get("note") or "미납관리비 없음 또는 미확인"
            items.append(f"관리비 현황은 {_polite_confirmation(note)}")
    else:
        items.append("관리비 미납 금액 및 개월수는 관리사무소에 직접 확인해 주시기 바랍니다.")

    if any(t.get("isVacant") for t in tenants):
        items.append("임차권등기 또는 현황자료상 공실 가능성이 있으므로 점유 현황을 현장에서 확인해 주시기 바랍니다.")

    if market_data.get("recentDealPrice"):
        items.append(
            f"최근 실거래가 {fmt_money(market_data.get('recentDealPrice'))}이(가) 확인됩니다. "
            "층수·면적·거래시점 차이를 고려해 비교해 주시기 바랍니다."
        )

    items.append("본 보증서 발급 전 등기부등본, 매각물건명세서, 현황조사서 최신본을 반드시 재확인해 주시기 바랍니다.")
    return items


def build_unpaid_management_fee_text(management_fee: dict) -> str:
    unpaid = int(management_fee.get("unpaidAmount") or 0)
    if unpaid > 0:
        due_text = management_fee.get("dueThroughText") or _extract_management_fee_due_text(management_fee.get("note") or "")
        if due_text:
            return (
                f"미납관리비 {due_text} 약 {fmt_money(unpaid)}이 존재합니다. "
                "낙찰시 명도시점까지 미납관리비가 추가로 더 발생 될 수 있으며, "
                "미납관리비는 낙찰자에게 인수됩니다. 자세한 내용은 뒷장에 참조 바랍니다."
            )
        return (
            f"미납관리비 약 {fmt_money(unpaid)}이 존재합니다. "
            "낙찰시 명도시점까지 미납관리비가 추가로 더 발생 될 수 있으며, "
            "미납관리비는 낙찰자에게 인수됩니다. 자세한 내용은 뒷장에 참조 바랍니다."
        )
    if management_fee:
        note = management_fee.get("note") or "미납관리비 없음 또는 미확인"
        return f"{_polite_confirmation(note)} 최종 입찰 전 관리사무소에 재확인해 주시기 바랍니다."
    return "미납관리비 금액 및 개월수는 관리사무소에 직접 확인해 주시기 바랍니다."


def _polite_optional_note(value: str, fallback: str) -> str:
    return _polite_confirmation(value) or fallback


def _polite_confirmation(value: str) -> str:
    note = _clean_document_note(value)
    if not note:
        return ""
    if note.endswith(("습니다.", "입니다.", "바랍니다.", "됩니다.", "합니다.", "없습니다.")):
        return note
    note = note.rstrip(".。")
    return f"{note}{_euro_josa(note)} 확인됩니다."


def _euro_josa(value: str) -> str:
    text = str(value or "").strip()
    if not text:
        return "로"
    last = text[-1]
    code = ord(last) - 0xAC00
    if 0 <= code <= 11171 and code % 28:
        return "으로"
    return "로"


def build_review_items(rights: list[dict], tenants: list[dict], management_fee: dict) -> list[str]:
    items = []
    if not rights:
        items.append("등기부현황 또는 권리분석 테이블에서 권리 목록을 확인해 주시기 바랍니다.")
    if not tenants:
        items.append("임차인현황 테이블에서 임차인 목록을 확인해 주시기 바랍니다.")
    if not management_fee:
        items.append("관리비 현황은 관리사무소 또는 현장에서 확인해 주시기 바랍니다.")
    if not items:
        items.append("원본 문서와 대조해 최종 검토해 주시기 바랍니다.")
    return items


def render_certificate_template(template_path: str, data: dict) -> str:
    path = Path(template_path)
    if not path.exists():
        raise FileNotFoundError(f"권리분석 보증서 템플릿을 찾을 수 없습니다: {path}")

    template = path.read_text(encoding="utf-8")
    rendered = template

    rendered = _replace_each(rendered, "tenantAnalyses", data.get("tenantAnalyses") or [])
    rendered = _replace_each(rendered, "miscItems", data.get("miscItems") or [])
    rendered = _replace_each(rendered, "reviewItems", data.get("reviewItems") or [])

    rendered = _replace_if(rendered, "noTenants", bool(data.get("noTenants")))
    rendered = _replace_if(rendered, "hasUnpaidFee", bool(data.get("hasUnpaidFee")))

    for key, value in data.items():
        if isinstance(value, (list, dict)):
            continue
        rendered = rendered.replace("{{" + key + "}}", html.escape("" if value is None else str(value)))
    rendered = rendered.replace("{{address}}", "").replace("{{adress}}", "")
    rendered = re.sub(r"[ \t]*(?:/|,|\|)[ \t]*(?=(?:<br>|</p>|\n|$))", "", rendered)
    return rendered


def render_certificate_pptx_template(template_path: Path, output_path: Path, data: dict) -> None:
    if not template_path.exists():
        raise FileNotFoundError(f"권리분석 보증서 PPT 템플릿을 찾을 수 없습니다: {template_path}")

    prs = Presentation(str(template_path))
    special_summary_chunks = _prepare_special_summary_slides(prs, data)
    mapping = {
        "{{" + key + "}}": _ppt_value(value)
        for key, value in data.items()
        if not isinstance(value, (list, dict))
    }
    mapping["{{address}}"] = ""
    mapping["{{adress}}"] = ""

    special_slide_index = 0
    for slide in prs.slides:
        slide_mapping = mapping
        if special_summary_chunks and _slide_contains_text(slide, "{{specialSummaryText}}"):
            chunk = special_summary_chunks[min(special_slide_index, len(special_summary_chunks) - 1)]
            slide_mapping = dict(mapping)
            slide_mapping["{{specialSummaryText}}"] = chunk
            slide_mapping["{{특이사항요약}}"] = chunk
            special_slide_index += 1
        _replace_placeholders_in_shapes(slide.shapes, slide_mapping)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def export_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    try:
        import win32com.client
    except Exception as e:
        logger.warning(f"PowerPoint PDF 변환 모듈을 사용할 수 없습니다: {e}")
        return False

    app = None
    presentation = None
    try:
        pdf_path.parent.mkdir(parents=True, exist_ok=True)
        app = win32com.client.DispatchEx("PowerPoint.Application")
        try:
            app.DisplayAlerts = 0
        except Exception:
            pass
        presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        presentation.SaveAs(str(pdf_path.resolve()), 32)
        return pdf_path.exists() and pdf_path.stat().st_size > 0
    except Exception as e:
        logger.warning(f"PowerPoint PDF 변환 실패: {e}")
        return False
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass


def _prepare_special_summary_slides(prs: Presentation, data: dict) -> list[str]:
    text = _ppt_value(data.get("specialSummaryText") or data.get("특이사항요약") or "")
    chunks = _split_special_summary_text(text)
    if len(chunks) <= 1:
        return chunks

    base_index = _find_slide_index_with_text(prs, "{{specialSummaryText}}")
    if base_index is None:
        return chunks

    base_slide = prs.slides[base_index]
    for offset in range(1, len(chunks)):
        new_slide = _duplicate_slide(prs, base_slide)
        _move_slide(prs, prs.slides.index(new_slide), base_index + offset)
    return chunks


def _split_special_summary_text(text: str) -> list[str]:
    lines = str(text or "").splitlines()
    if not lines:
        return [""]

    chunks: list[str] = []
    current: list[str] = []
    current_chars = 0
    current_weight = 0

    for line in lines:
        line_chars = len(line)
        line_weight = _special_summary_line_weight(line)
        should_split = (
            bool(current)
            and (
                current_chars + line_chars > SPECIAL_SUMMARY_SLIDE_CHAR_LIMIT
                or current_weight + line_weight > SPECIAL_SUMMARY_MAX_WEIGHTED_LINES
            )
        )
        if should_split:
            chunks.append("\n".join(current).strip())
            current = []
            current_chars = 0
            current_weight = 0
            if line.strip() and not re.match(r"^\d+\)", line.strip()):
                current.append("(계속)")
                current_chars += len("(계속)")
                current_weight += 1

        current.append(line)
        current_chars += line_chars + 1
        current_weight += line_weight

    if current:
        chunks.append("\n".join(current).strip())
    return chunks or [""]


def _special_summary_line_weight(line: str) -> int:
    text = line.strip()
    if not text:
        return 1
    return max(1, (len(text) + SPECIAL_SUMMARY_WRAP_WIDTH - 1) // SPECIAL_SUMMARY_WRAP_WIDTH)


def _find_slide_index_with_text(prs: Presentation, needle: str) -> Optional[int]:
    for idx, slide in enumerate(prs.slides):
        if _slide_contains_text(slide, needle):
            return idx
    return None


def _slide_contains_text(slide, needle: str) -> bool:
    return any(needle in text for text in _iter_shape_texts(slide.shapes))


def _iter_shape_texts(shapes):
    for shape in shapes:
        if hasattr(shape, "shapes"):
            yield from _iter_shape_texts(shape.shapes)
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame.text
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text


def _duplicate_slide(prs: Presentation, slide):
    new_slide = prs.slides.add_slide(slide.slide_layout)
    rel_map = _copy_slide_relationships(slide, new_slide)
    _copy_slide_background(slide, new_slide, rel_map)
    for shape in slide.shapes:
        new_el = deepcopy(shape._element)
        _remap_relationship_ids(new_el, rel_map)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def _copy_slide_background(source_slide, target_slide, rel_map: dict[str, str]) -> None:
    source_c_sld = source_slide._element.find(qn("p:cSld"))
    target_c_sld = target_slide._element.find(qn("p:cSld"))
    if source_c_sld is None or target_c_sld is None:
        return

    source_bg = source_c_sld.find(qn("p:bg"))
    if source_bg is None:
        return

    target_bg = target_c_sld.find(qn("p:bg"))
    if target_bg is not None:
        target_c_sld.remove(target_bg)

    new_bg = deepcopy(source_bg)
    _remap_relationship_ids(new_bg, rel_map)
    target_c_sld.insert(0, new_bg)


def _copy_slide_relationships(source_slide, target_slide) -> dict[str, str]:
    rel_map: dict[str, str] = {}
    relationships = sorted(
        source_slide.part.rels.values(),
        key=lambda rel: int(re.sub(r"\D+", "", rel.rId) or 0),
    )
    for rel in relationships:
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            new_rid = target_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = target_slide.part.relate_to(rel.target_part, rel.reltype)
        rel_map[rel.rId] = new_rid
    return rel_map


def _remap_relationship_ids(element, rel_map: dict[str, str]) -> None:
    if not rel_map:
        return
    for node in element.iter():
        for attr_name, attr_value in list(node.attrib.items()):
            if attr_value in rel_map:
                node.set(attr_name, rel_map[attr_value])


def _move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slide_id = xml_slides[old_index]
    xml_slides.remove(slide_id)
    if new_index > old_index:
        new_index -= 1
    xml_slides.insert(new_index, slide_id)


def _replace_placeholders_in_shapes(shapes, mapping: dict[str, str]) -> None:
    for shape in shapes:
        if hasattr(shape, "shapes"):
            _replace_placeholders_in_shapes(shape.shapes, mapping)

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    _replace_text_frame_placeholders(cell.text_frame, mapping)

        if getattr(shape, "has_text_frame", False):
            _replace_text_frame_placeholders(shape.text_frame, mapping)


def _replace_text_frame_placeholders(text_frame, mapping: dict[str, str]) -> None:
    old_text = "\n".join(p.text for p in text_frame.paragraphs)
    new_text = old_text
    body_font_size = BODY_FONT_SIZE if any(token in old_text for token in BODY_PLACEHOLDER_TOKENS) else None
    case_info_font_size = CASE_INFO_FONT_SIZE if any(token in old_text for token in CASE_INFO_PLACEHOLDER_TOKENS) else None
    for token, value in mapping.items():
        new_text = new_text.replace(token, value)
    new_text = re.sub(r"[ \t]*(?:/|,|\|)[ \t]*(?=\n|$)", "", new_text)

    if new_text == old_text:
        return

    font_template = _first_run_font(text_frame)
    text_frame.clear()
    for index, line in enumerate(new_text.split("\n")):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.clear()
        run = paragraph.add_run()
        run.text = line
        _copy_font(font_template, run.font)
        if case_info_font_size is not None:
            run.font.size = case_info_font_size
        elif body_font_size is not None:
            run.font.size = _font_size_for_body_line(line)


def _font_size_for_body_line(line: str):
    if "최우선변제" in line or "잔존 금액" in line:
        return PRIORITY_REPAYMENT_FONT_SIZE
    return BODY_FONT_SIZE


def _first_run_font(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            return run.font
    return None


def _copy_font(src, dst) -> None:
    if src is None:
        return
    try:
        dst.name = src.name
        dst.size = src.size
        dst.bold = src.bold
        dst.italic = src.italic
        dst.underline = src.underline
    except Exception:
        pass
    try:
        dst.color.rgb = src.color.rgb
    except Exception:
        pass


def _ppt_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "예" if value else "아니오"
    return str(value)


def render_html_to_pdf(driver, html_path: Path, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    driver.get(html_path.resolve().as_uri())
    wait_document_ready(driver, timeout=20)
    result = driver.execute_cdp_cmd(
        "Page.printToPDF",
        {
            "landscape": False,
            "printBackground": True,
            "preferCSSPageSize": True,
            "marginTop": 0,
            "marginRight": 0,
            "marginBottom": 0,
            "marginLeft": 0,
        },
    )
    output_path.write_bytes(base64.b64decode(result["data"]))


def _replace_each(template: str, key: str, items: list) -> str:
    pattern = re.compile(r"{{#each " + re.escape(key) + r"}}([\s\S]*?){{/each}}")

    def repl(match):
        block = match.group(1)
        parts = []
        for item in items:
            piece = block
            if isinstance(item, dict):
                for item_key, value in item.items():
                    piece = piece.replace("{{this." + item_key + "}}", html.escape(str(value or "")))
            else:
                piece = piece.replace("{{this}}", html.escape(str(item or "")))
            parts.append(piece)
        return "\n".join(parts)

    return pattern.sub(repl, template)


def _replace_if(template: str, key: str, enabled: bool) -> str:
    pattern = re.compile(r"{{#if " + re.escape(key) + r"}}([\s\S]*?){{/if}}")
    return pattern.sub(lambda m: m.group(1) if enabled else "", template)


def _extract_rights(soup) -> list[dict]:
    rights = []
    for table in soup.find_all("table"):
        table_text = table.get_text(" ", strip=True)
        if not any(t in table_text for t in RIGHT_TYPES):
            continue
        if not re.search(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", table_text):
            continue
        for row in _iter_table_rows(table):
            row = _strip_registry_sequence_cells(row)
            row_text = " ".join(row)
            right_type = _first_match(row_text, RIGHT_TYPES)
            date = _first_date(row_text)
            if not right_type or not date:
                continue
            rights.append(
                {
                    "seq": len(rights) + 1,
                    "date": normalize_date(date),
                    "type": right_type,
                    "creditor": _guess_creditor(row, right_type),
                    "amount": parse_money(row_text),
                    "status": _extract_right_status(row_text),
                    "note": _extract_right_note(row, row_text),
                    "isBaseRight": _is_base_right_row(row_text),
                    "rawText": row_text,
                }
            )
    return _dedupe_by(rights, ("date", "type", "creditor", "amount"))


def _strip_registry_sequence_text(text: str) -> str:
    return re.sub(r"^\s*\d{1,3}\s+(?=\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2})", "", text or "").strip()


def _strip_registry_sequence_cells(cells: list[str]) -> list[str]:
    if len(cells) >= 2 and re.fullmatch(r"\d{1,3}", cells[0] or "") and _first_date(cells[1]):
        return cells[1:]
    return cells


def merge_rights(*groups: list[dict]) -> list[dict]:
    merged: list[dict] = []
    for group in groups:
        for right in group or []:
            match = _find_matching_right(merged, right)
            if not match:
                item = dict(right)
                item["seq"] = len(merged) + 1
                merged.append(item)
                continue
            for key, value in right.items():
                if key == "isBaseRight":
                    match[key] = bool(match.get(key)) or bool(value)
                elif key == "amount" and not match.get(key) and value:
                    match[key] = value
                elif key not in ("seq",) and not match.get(key) and value:
                    match[key] = value
    return merged


def _find_matching_right(rights: list[dict], target: dict) -> Optional[dict]:
    target_date = target.get("date") or ""
    target_type = target.get("type") or ""
    target_creditor = target.get("creditor") or ""
    for right in rights:
        if not _same_date(right.get("date") or "", target_date):
            continue
        right_type = right.get("type") or ""
        if target_type and right_type and target_type not in right_type and right_type not in target_type:
            continue
        right_creditor = right.get("creditor") or ""
        if target_creditor and right_creditor and not _creditor_names_match(target_creditor, right_creditor):
            continue
        return right
    return None


def _is_base_right_row(text: str) -> bool:
    compact = re.sub(r"\s+", "", text or "")
    return any(token in compact for token in ("말소기준권리", "말소기준", "소멸기준"))


def _extract_right_status(text: str) -> str:
    compact = re.sub(r"\s+", "", text or "")
    if "소멸기준" in compact:
        return "소멸기준"
    if "인수" in compact:
        return "인수"
    if "소멸" in compact:
        return "소멸"
    return ""


def _extract_right_note(cells: list[str], row_text: str) -> str:
    if cells:
        for cell in reversed(cells):
            compact = re.sub(r"\s+", "", cell or "")
            if any(token in compact for token in ("말소기준", "존속기간", "배당요구", "카단", "타경")):
                return cell[:200]
    if _is_base_right_row(row_text) or "배당요구" in (row_text or ""):
        return (row_text or "")[:200]
    return ""


def _right_has_dividend_request(right: dict, dividend_requests: list[dict]) -> bool:
    raw_text = f"{right.get('rawText') or ''} {right.get('note') or ''}"
    compact = re.sub(r"\s+", "", raw_text)
    if any(token in compact for token in ("배당요구없", "배당요구하지않", "배당요구안")):
        return False
    if "배당요구" in compact:
        return True
    return _find_dividend_request(right.get("creditor") or "", dividend_requests) is not None


def _extract_tenants(soup) -> list[dict]:
    tenants = []
    for table in soup.find_all("table"):
        table_text = table.get_text(" ", strip=True)
        if "임차" not in table_text and "보증금" not in table_text:
            continue
        for row in _iter_table_rows(table):
            row_text = " ".join(row)
            if "보증금" not in row_text and not re.search(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", row_text):
                continue
            dates = [normalize_date(d) for d in re.findall(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", row_text)]
            name = row[0] if row else ""
            if not name or any(word in name for word in ("임차", "점유", "성명")):
                name = _guess_name(row)
            tenants.append(
                {
                    "name": name or "임차인",
                    "type": "주택도시보증공사" if "주택도시보증공사" in row_text or "HUG" in row_text.upper() else "",
                    "moveInDate": dates[0] if len(dates) > 0 else "",
                    "fixedDate": dates[1] if len(dates) > 1 else "",
                    "depositClaimDate": dates[2] if len(dates) > 2 else "",
                    "depositDeadline": "",
                    "deposit": parse_money(row_text),
                    "rent": 0,
                    "isHUG": "주택도시보증공사" in row_text or "HUG" in row_text.upper(),
                    "isVacant": "공실" in row_text,
                }
            )
    return _dedupe_by(tenants, ("name", "moveInDate", "deposit"))


def _extract_dividend_requests(soup) -> list[dict]:
    requests = []
    for table in soup.find_all("table"):
        table_text = table.get_text(" ", strip=True)
        if "배당" not in table_text:
            continue
        for row in _iter_table_rows(table):
            row_text = " ".join(row)
            dates = [normalize_date(d) for d in re.findall(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", row_text)]
            if not dates:
                continue
            requests.append(
                {
                    "creditor": _guess_name(row) or (row[0] if row else ""),
                    "requestDate": dates[0],
                    "deadline": dates[1] if len(dates) > 1 else "",
                    "amount": parse_money(row_text),
                }
            )
    return _dedupe_by(requests, ("creditor", "requestDate", "deadline", "amount"))


def _extract_expected_dividend(soup, case_number: str = "", applicant_creditors: Optional[list[str]] = None) -> dict:
    applicant_creditors = applicant_creditors if applicant_creditors is not None else _extract_auction_applicant_creditors(soup, case_number)
    table_result = _extract_expected_dividend_from_tables(soup, applicant_creditors, case_number)
    if table_result:
        return table_result

    direct = _parse_money_cell(
        _css_text(soup, "#dtl_table > table > tbody > tr:nth-child(8) > td:nth-child(4)")
    )
    if direct.get("found"):
        return {
            "auctionApplicantDividendAmount": direct["amount"],
            "auctionApplicantDividendFound": True,
            "source": "dtl_table_selector",
        }

    return {}


def _extract_expected_dividend_from_tables(soup, applicant_creditors: list[str], case_number: str = "") -> dict:
    applicant_names = [_normalize_creditor_name(name) for name in applicant_creditors]
    applicant_names = [name for name in applicant_names if name]
    normalized_case_number = _normalize_case_number(case_number)

    stock = soup.select_one("#dtl_stock") or soup
    for table in stock.find_all("table"):
        table_text = table.get_text(" ", strip=True)
        if not _is_expected_dividend_table_text(table_text):
            continue
        rows = _iter_table_rows(table)
        if not rows:
            continue

        header_idx = _find_header_row_index(rows, ("채권배당금", "배당금", "배당액", "순위배당"))
        header = rows[header_idx] if header_idx is not None else rows[0]
        data_rows = rows[header_idx + 1:] if header_idx is not None else rows[1:]
        creditor_idx = _find_header_index(header, ("권리자", "채권자", "성명"))
        amount_idx = _find_header_index(header, ("채권배당금", "배당금", "배당액", "순위배당"))
        for row in data_rows:
            row_text = " ".join(row)
            creditor = _expected_dividend_row_creditor(row, creditor_idx)
            is_applicant_row = (
                _is_auction_applicant_row(row_text)
                or _row_has_case_number(row_text, normalized_case_number)
                or _creditor_matches_any(creditor, applicant_names)
            )
            if not is_applicant_row:
                continue

            candidates = []
            if amount_idx is not None and amount_idx < len(row):
                candidates.append(row[amount_idx])
            if len(row) >= 4:
                candidates.append(row[3])
            candidates.append(row_text)

            for candidate in candidates:
                parsed = _parse_money_cell(candidate)
                if parsed.get("found"):
                    return {
                        "auctionApplicantDividendAmount": parsed["amount"],
                        "auctionApplicantDividendFound": True,
                        "source": "expected_dividend_table",
                        "auctionApplicantCreditor": creditor,
                        "rowText": row_text,
                    }
    return {}


def _is_expected_dividend_table_text(table_text: str) -> bool:
    compact = re.sub(r"\s+", "", table_text or "")
    return (
        "예상배당표" in compact
        or "채권배당금" in compact
        or "입찰예상가" in compact
        or "최저경매가기준" in compact
    )


def _extract_auction_applicant_creditors(soup, case_number: str = "") -> list[str]:
    creditors = []
    case_matched_creditors = []
    normalized_case_number = _normalize_case_number(case_number)
    for table in soup.find_all("table"):
        table_text = table.get_text(" ", strip=True)
        if "임의경매" not in table_text and "강제경매" not in table_text:
            continue
        for row in _iter_table_rows(table):
            row_text = " ".join(row)
            if "임의경매" not in row_text and "강제경매" not in row_text:
                continue
            creditor = _guess_auction_applicant_creditor(row)
            if creditor:
                creditors.append(creditor)
                if _row_has_case_number(row_text, normalized_case_number):
                    case_matched_creditors.append(creditor)
    return _dedupe_strings(case_matched_creditors or creditors)


def _guess_auction_applicant_creditor(row: list[str]) -> str:
    for idx, cell in enumerate(row):
        if "임의경매" not in cell and "강제경매" not in cell:
            continue
        for candidate in row[idx + 1:]:
            creditor = _clean_creditor_candidate(candidate)
            if creditor:
                return creditor
        row_text = " ".join(row)
        tail = re.split(r"임의경매|강제경매", row_text, maxsplit=1)[-1]
        return _clean_creditor_candidate(tail)
    return ""


def _clean_creditor_candidate(value: str) -> str:
    text = re.sub(r"\[[^\]]*\]", " ", str(value or ""))
    text = re.sub(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", " ", text)
    text = re.sub(r"\d{4}\s*타경\s*\d+", " ", text)
    text = re.sub(r"\d{1,3}(?:,\d{3})+|\d{4,}\s*원?", " ", text)
    text = re.sub(r"\b\d+\b", " ", text)
    text = re.sub(r"(소멸|인수|청구금액|채권액|채권최고액|비고|권리자|권리|임의경매|강제경매)", " ", text)
    text = re.sub(r"\s+", " ", text).strip(" /,|")
    if not re.search(r"[가-힣A-Za-z]", text):
        return ""
    return text


def _expected_dividend_row_creditor(row: list[str], creditor_idx: Optional[int]) -> str:
    if creditor_idx is not None and creditor_idx < len(row):
        creditor = _clean_creditor_candidate(row[creditor_idx])
        if creditor:
            return creditor
    if len(row) > 1:
        creditor = _clean_creditor_candidate(row[1])
        if creditor:
            return creditor
    return _clean_creditor_candidate(" ".join(row))


def _creditor_matches_any(creditor: str, applicant_names: list[str]) -> bool:
    normalized = _normalize_creditor_name(creditor)
    if not normalized:
        return False
    return any(_creditor_names_match(normalized, applicant) for applicant in applicant_names)


def _normalize_case_number(value: str) -> str:
    text = re.sub(r"\s+", "", str(value or ""))
    m = re.search(r"(\d{4})타경(\d+)", text)
    return f"{m.group(1)}타경{m.group(2)}" if m else text


def _row_has_case_number(row_text: str, normalized_case_number: str) -> bool:
    compact = re.sub(r"\s+", "", str(row_text or ""))
    return bool(normalized_case_number and normalized_case_number in compact)


def _normalize_creditor_name(value: str) -> str:
    text = str(value or "")
    for token in ("주식회사", "(주)", "㈜", "유한회사", "합자회사", "합명회사", "재단법인", "사단법인"):
        text = text.replace(token, "")
    return re.sub(r"[^0-9A-Za-z가-힣]+", "", text)


def _creditor_names_match(left: str, right: str) -> bool:
    left = _normalize_creditor_name(left)
    right = _normalize_creditor_name(right)
    if not left or not right:
        return False
    if left == right:
        return True
    return min(len(left), len(right)) >= 3 and (left in right or right in left)


def _dedupe_strings(values: list[str]) -> list[str]:
    result = []
    seen = set()
    for value in values:
        key = _normalize_creditor_name(value)
        if key and key not in seen:
            seen.add(key)
            result.append(value)
    return result


def _is_auction_applicant_row(text: str) -> bool:
    compact = re.sub(r"\s+", "", text or "")
    return (
        "경매신청채권자" in compact
        or ("경매" in compact and "신청" in compact and "채권자" in compact)
        or "임의경매" in compact
        or "강제경매" in compact
    )


def _find_header_row_index(rows: list[list[str]], keywords: tuple[str, ...]) -> Optional[int]:
    for idx, row in enumerate(rows):
        if _find_header_index(row, keywords) is not None:
            return idx
    return None


def _find_header_index(row: list[str], keywords: tuple[str, ...]) -> Optional[int]:
    for idx, cell in enumerate(row):
        compact = re.sub(r"\s+", "", cell or "")
        if any(keyword in compact for keyword in keywords):
            return idx
    return None


def _parse_money_cell(value) -> dict:
    text = re.sub(r"\s+", "", str(value or ""))
    if not text:
        return {"found": False, "amount": 0}
    if re.search(r"(?:^|[^\d])0원", text) or text in ("0", "0원", "-0원"):
        return {"found": True, "amount": 0}
    amount = parse_money(text)
    return {"found": amount > 0, "amount": amount}


def _extract_related_cases(soup) -> list[dict]:
    cases = []
    text = soup.get_text(" ", strip=True)
    case_pattern = r"(?:\d{4})\s*타경\s*\d+"
    for m in re.finditer(r"중복경매|병합경매|중복|병합", text):
        case_type = m.group(0)
        if case_type in ("중복", "병합"):
            case_type = f"{case_type}경매"

        after = text[m.end(): m.end() + 35]
        after_match = re.search(case_pattern, after)
        if after_match:
            case_number = _normalize_case_number(after_match.group(0))
            cases.append({"caseNumber": case_number, "creditor": "", "type": case_type, "filingDate": ""})
            continue

        before = text[max(0, m.start() - 35): m.start()]
        before_matches = list(re.finditer(case_pattern, before))
        if before_matches:
            case_number = _normalize_case_number(before_matches[-1].group(0))
            cases.append({"caseNumber": case_number, "creditor": "", "type": case_type, "filingDate": ""})
    return _dedupe_by(cases, ("caseNumber", "type"))


def _extract_management_fee(soup) -> dict:
    text = re.sub(r"\s+", " ", soup.get_text(" ", strip=True))
    if "관리비" not in text:
        return {}

    candidates: list[tuple[int, int, str]] = []
    for match in re.finditer(r"(?:미납\s*)?관리비|체납\s*관리비", text):
        nearby = text[max(0, match.start() - 120): min(len(text), match.end() + 260)]
        if "관리비" not in nearby:
            continue
        unpaid_hint = any(word in nearby for word in ("미납", "체납"))
        amount = _extract_management_fee_amount(nearby) if unpaid_hint else 0
        keyword_score = 2 if re.search(r"미납\s*관리비|체납\s*관리비", match.group(0)) else 1
        amount_score = 10 if amount > 0 else 0
        candidates.append((amount_score + keyword_score, amount, nearby))

    if not candidates:
        idx = text.find("관리비")
        nearby = text[max(0, idx - 120): min(len(text), idx + 240)]
        candidates.append((0, 0, nearby))

    _, amount, nearby = max(candidates, key=lambda item: item[0])
    return {
        "unpaidAmount": amount,
        "unpaidMonths": 0,
        "dueThroughText": _extract_management_fee_due_text(nearby),
        "checkDate": normalize_date(_first_date(nearby) or ""),
        "note": _clip_text(nearby, 160),
    }


def _extract_management_fee_due_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "")
    match = re.search(r"(?<!\d)(\d{2,4})\s*(?:년|[.\-/])\s*(\d{1,2})\s*월?\s*까지", text)
    if match:
        year = int(match.group(1))
        month = int(match.group(2))
        return f"{year % 100:02d}년 {month}월까지"
    match = re.search(r"(?<!\d)(\d{1,2})\s*월\s*까지", text)
    if match:
        return f"{int(match.group(1))}월까지"
    return ""


def _extract_management_fee_amount(text: str) -> int:
    patterns = (
        r"(?:미납|체납)\s*관리비[^\d]{0,40}(\d{1,3}(?:,\d{3})+|\d{4,})\s*원?",
        r"관리비[^\d]{0,40}(?:미납|체납)?[^\d]{0,40}(\d{1,3}(?:,\d{3})+|\d{4,})\s*원?",
        r"(\d{1,3}(?:,\d{3})+|\d{4,})\s*원[^\d]{0,50}(?:미납|체납)\s*관리비",
    )
    for pattern in patterns:
        amounts = [parse_money(match.group(1)) for match in re.finditer(pattern, text or "")]
        amounts = [amount for amount in amounts if amount > 0]
        if amounts:
            return max(amounts)
    return parse_money(text)


def _extract_market_data(soup) -> dict:
    text = soup.get_text(" ", strip=True)
    if "실거래" not in text and "매각사례" not in text:
        return {}
    idx = max(text.find("실거래"), text.find("매각사례"))
    nearby = text[max(0, idx - 80): idx + 240]
    return {
        "recentDealPrice": parse_money(nearby),
        "recentDealDate": normalize_date(_first_date(nearby) or ""),
        "trend": "",
    }


def _iter_table_rows(table) -> list[list[str]]:
    rows = []
    for tr in table.find_all("tr"):
        cells = [c.get_text(" ", strip=True) for c in tr.find_all(["th", "td"])]
        cells = [re.sub(r"\s+", " ", c).strip() for c in cells if c and c.strip()]
        if cells:
            rows.append(cells)
    return rows


def _css_text(soup, *selectors: str) -> str:
    for selector in selectors:
        try:
            el = soup.select_one(selector)
        except Exception:
            el = None
        if not el:
            continue
        text = re.sub(r"\s+", " ", el.get_text(" ", strip=True)).strip()
        if text:
            return text
    return ""


def _first_match(text: str, candidates: tuple[str, ...]) -> str:
    for candidate in candidates:
        if candidate in text:
            return candidate
    return ""


def _first_date(text: str) -> str:
    m = re.search(r"\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2}", text or "")
    return m.group(0) if m else ""


def _first_any_date(text: str) -> str:
    pattern = r"\d{4}\s*(?:[.\-/년])\s*\d{1,2}\s*(?:[.\-/월])\s*\d{1,2}\s*\.?\s*일?"
    match = re.search(pattern, text or "")
    return normalize_date(match.group(0)) if match else ""


def normalize_date(date_str: str) -> str:
    if not date_str:
        return ""
    nums = re.findall(r"\d+", date_str)
    if len(nums) < 3:
        return date_str
    return f"{int(nums[0]):04d}.{int(nums[1]):02d}.{int(nums[2]):02d}"


def _date_sort_key(date_str: Optional[str]) -> tuple[int, int, int]:
    nums = re.findall(r"\d+", date_str or "")
    if len(nums) < 3:
        return (9999, 99, 99)
    return (int(nums[0]), int(nums[1]), int(nums[2]))


def _has_valid_date(date_str: str) -> bool:
    return _date_sort_key(date_str) != (9999, 99, 99)


def _date_before(left: str, right: str) -> bool:
    return _has_valid_date(left) and _has_valid_date(right) and _date_sort_key(left) < _date_sort_key(right)


def _same_date(left: str, right: str) -> bool:
    return _has_valid_date(left) and _has_valid_date(right) and _date_sort_key(left) == _date_sort_key(right)


def _date_on_or_before(left: str, right: str) -> bool:
    return _has_valid_date(left) and _has_valid_date(right) and _date_sort_key(left) <= _date_sort_key(right)


def _date_after(left: str, right: str) -> bool:
    return _has_valid_date(left) and _has_valid_date(right) and _date_sort_key(left) > _date_sort_key(right)


def parse_money(value) -> int:
    if value is None:
        return 0
    if isinstance(value, (int, float)):
        return int(value)
    text = str(value)
    comma_nums = re.findall(r"\d{1,3}(?:,\d{3})+", text)
    if comma_nums:
        return max(int(n.replace(",", "")) for n in comma_nums)
    won_nums = re.findall(r"(\d{4,})\s*원", text)
    if won_nums:
        return max(int(n) for n in won_nums)
    if re.fullmatch(r"\d{4,}", text.strip()):
        return int(text.strip())
    return 0


def fmt_money(value) -> str:
    amount = parse_money(value)
    if amount <= 0:
        return "담당자 확인 필요"
    return f"{amount:,}원"


def fmt_money_or_unknown(value) -> str:
    amount = parse_money(value)
    if amount <= 0:
        return "없음 또는 미확인"
    return f"{amount:,}원"


def format_korean_date(date: datetime) -> str:
    return f"{date.year}년 {date.month}월 {date.day}일"


def _clip_text(text: str, limit: int) -> str:
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "..."


def _guess_creditor(cells: list[str], right_type: str) -> str:
    for cell in cells:
        if right_type in cell or _first_date(cell) or parse_money(cell):
            continue
        if len(cell) <= 2:
            continue
        return cell[:80]
    return ""


def _guess_name(cells: list[str]) -> str:
    for cell in cells:
        if _first_date(cell) or parse_money(cell):
            continue
        if any(word in cell for word in ("임차", "배당", "보증금", "전입", "확정")):
            continue
        if 1 < len(cell) <= 40:
            return cell
    return ""


def _find_dividend_request(name: str, requests: list[dict]) -> Optional[dict]:
    if not name:
        return None
    for request in requests:
        creditor = request.get("creditor") or ""
        if _creditor_names_match(name, creditor):
            return request
    return None


def _dedupe_by(items: list[dict], keys: tuple[str, ...]) -> list[dict]:
    seen = set()
    unique = []
    for item in items:
        marker = tuple(item.get(k) for k in keys)
        if marker in seen:
            continue
        seen.add(marker)
        unique.append(item)
    return unique
