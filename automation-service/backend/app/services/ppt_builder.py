# -*- coding: utf-8 -*-
"""
PPT 생성 서비스
- PPTM 템플릿 로드
- 데이터 채우기 (1번/3번 슬라이드)
- 사진 삽입
- 슬라이드 복제/이동
- 노란 박스 찾기 → 이미지 교체
- VBA 보존 저장
"""

import os
import re
import logging
import shutil
import tempfile
import zipfile
from io import BytesIO
from copy import deepcopy
from typing import Optional, Tuple, Dict

import requests
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from ..core.config import settings, CAPTURE_DIR
from ..core.utils import track_file
from .capturer import trim_white_margin

logger = logging.getLogger(__name__)

TARGET_SLIDE_INDEX = 2  # 0-base → 3번 슬라이드
DEFAULT_PPT_FONT_NAME = "G마켓 산스 TTF Medium"
OPINION_NAVY = RGBColor(20, 48, 92)
OPINION_BLUE = RGBColor(0, 150, 214)
RIGHTS_OPINION_HEADING_PT = 18
RIGHTS_OPINION_BODY_PT = 12


# ============================================================
# AltText / 토큰
# ============================================================
def get_alt_text(shape) -> str:
    try:
        cNvPr = shape._element.xpath(".//p:cNvPr")[0]
        return cNvPr.get("descr") or cNvPr.get("title") or ""
    except Exception:
        return ""


_NUM_RE = re.compile(r"(\d[\d,]*)")


def parse_token(alt: str) -> Tuple[str, Dict[str, str]]:
    alt = (alt or "").strip()
    if not alt.startswith("EXCEL_"):
        return "", {}
    parts = alt.split("|")
    head = parts[0]
    opts: Dict[str, str] = {}
    for p in parts[1:]:
        if "=" in p:
            k, v = p.split("=", 1)
            opts[k.strip().upper()] = v
    if ":" not in head:
        return "", {}
    kind, payload = head.split(":", 1)
    opts["PAYLOAD"] = payload.strip()
    return kind.strip().upper(), opts


# ============================================================
# 텍스트 유틸
# ============================================================
def replace_first_number_preserve_runs(text_frame, new_number: str) -> bool:
    runs, full = [], ""
    for p in text_frame.paragraphs:
        for r in p.runs:
            runs.append(r)
            full += r.text

    m = _NUM_RE.search(full)
    if not m:
        return False

    start, end = m.span(1)
    idx, replaced = 0, False
    for r in runs:
        _force_run_font(r)
        t = r.text or ""
        rs, re_ = idx, idx + len(t)
        if re_ <= start or rs >= end:
            idx = re_
            continue
        a = max(start - rs, 0)
        b = min(end - rs, len(t))
        if not replaced:
            r.text = t[:a] + new_number + t[b:]
            replaced = True
        else:
            r.text = t[:a] + t[b:]
        idx = re_
    return True


def set_text_keep_style(text_frame, new_text: str):
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    font_tpl = p.runs[0].font if p.runs else None
    p.clear()
    r = p.add_run()
    r.text = new_text
    _copy_font_style(font_tpl, r.font)
    _force_run_font(r)


def _copy_font_style(font_tpl, dst_font) -> None:
    if font_tpl:
        dst_font.size = font_tpl.size
        dst_font.bold = font_tpl.bold
        dst_font.italic = font_tpl.italic
        dst_font.underline = font_tpl.underline
        try:
            dst_font.color.rgb = font_tpl.color.rgb
        except Exception:
            pass


def _force_run_font(run, font_name: str = DEFAULT_PPT_FONT_NAME) -> None:
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r_pr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            node = r_pr.find(qn(tag))
            if node is None:
                node = OxmlElement(tag)
                r_pr.append(node)
            node.set("typeface", font_name)
    except Exception:
        pass


def set_multiline_text_keep_style(text_frame, new_text: str):
    first_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    font_tpl = first_paragraph.runs[0].font if first_paragraph.runs else None
    align_tpl = first_paragraph.alignment

    text_frame.clear()
    lines = str(new_text or "").splitlines() or [""]
    for idx, line in enumerate(lines):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.alignment = align_tpl
        r = p.add_run()
        r.text = line
        _copy_font_style(font_tpl, r.font)
        _force_run_font(r)


_WON_RE = re.compile(r"(\d[\d,]*)\s*원")


def remove_won_unit_in_slide(slide):
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        tf = shp.text_frame
        old = "\n".join([p.text for p in tf.paragraphs])
        new = _WON_RE.sub(r"\1", old)
        if new != old:
            set_multiline_text_keep_style(tf, new)


# ============================================================
# 슬라이드 검색
# ============================================================
def find_slide_by_keyword(prs: Presentation, keyword: str):
    hangul_re = re.compile(r"[가-힣]")
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if not text or keyword not in text:
                continue
            rest = text.replace(keyword, "")
            if hangul_re.search(rest):
                continue
            return slide
    for slide in prs.slides:
        try:
            notes = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes = ""
        if keyword in notes:
            return slide
    return None


def find_slide_by_note_key(prs: Presentation, key: str):
    for s in prs.slides:
        try:
            notes = s.notes_slide.notes_text_frame.text or ""
            if key in notes:
                return s
        except Exception:
            continue
    return None


def find_slide_by_note_keywords(prs: Presentation, keywords: list[str]):
    for s in prs.slides:
        try:
            notes = s.notes_slide.notes_text_frame.text or ""
        except Exception:
            continue
        if any(keyword in notes for keyword in keywords):
            return s
    return None


def find_slide_index_by_note_key(prs: Presentation, key: str) -> int:
    for idx, s in enumerate(prs.slides):
        try:
            notes = s.notes_slide.notes_text_frame.text or ""
            if key in notes:
                return idx
        except Exception:
            continue
    return -1


# ============================================================
# 노란 박스 찾기
# ============================================================
def find_yellow_box(slide):
    candidates = []
    for shape in slide.shapes:
        text = ""
        if hasattr(shape, "text"):
            try:
                text = shape.text or ""
            except Exception:
                text = ""
        if text.strip():
            continue
        try:
            fill = shape.fill
        except Exception:
            continue
        if not fill or fill.type is None:
            continue
        try:
            area = shape.width * shape.height
        except Exception:
            continue
        candidates.append((area, shape))
    if candidates:
        candidates.sort(key=lambda x: x[0], reverse=True)
        return candidates[0][1]
    return None


def find_yellow_boxes_left_to_right(slide, limit=2):
    candidates = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and (shape.text or "").strip():
                continue
        except Exception:
            pass
        try:
            fill = shape.fill
            if not fill or fill.type is None:
                continue
        except Exception:
            continue
        try:
            area = shape.width * shape.height
        except Exception:
            continue
        candidates.append((shape.left, -area, shape))
    candidates.sort(key=lambda x: (x[0], x[1]))
    return [c[2] for c in candidates[:limit]]


def find_image_boxes_left_to_right(slide, limit=2):
    candidates = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and (shape.text or "").strip():
                continue
        except Exception:
            pass
        try:
            area = shape.width * shape.height
        except Exception:
            continue
        if area < 500000 * 500000:
            continue
        shape_type = getattr(shape, "shape_type", None)
        name = getattr(shape, "name", "") or ""
        has_fill = False
        try:
            has_fill = bool(shape.fill and shape.fill.type is not None)
        except Exception:
            has_fill = False
        if shape_type == 13 or "그림" in name or "Picture" in name or has_fill:
            candidates.append((shape.left, -area, shape))
    candidates.sort(key=lambda x: (x[0], x[1]))
    return [c[2] for c in candidates[:limit]]


# ============================================================
# 슬라이드 복제/이동
# ============================================================
def duplicate_slide(prs: Presentation, slide):
    new_slide = prs.slides.add_slide(slide.slide_layout)
    for shape in slide.shapes:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slide_id = xml_slides[old_index]
    xml_slides.remove(slide_id)
    if new_index > old_index:
        new_index -= 1
    xml_slides.insert(new_index, slide_id)


# ============================================================
# 데이터 채우기
# ============================================================
def fill_slide_with_data(prs: Presentation, data: dict):
    for idx, slide in enumerate(prs.slides):
        values = dict(data or {})
        main_addr = values.get("address", "") or ""
        old_addr = values.get("address_old", "") or ""
        if idx == 0 and old_addr:
            values["address"] = f"{main_addr}\n{old_addr}"
        _replace_template_tokens_in_shapes(slide.shapes, values)


def _replace_template_tokens_in_shapes(shapes, data: dict) -> int:
    updated = 0
    for shape in shapes:
        if hasattr(shape, "shapes"):
            updated += _replace_template_tokens_in_shapes(shape.shapes, data)

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    updated += _replace_template_tokens_in_text_frame(cell.text_frame, data)

        if getattr(shape, "has_text_frame", False):
            updated += _replace_template_tokens_in_text_frame(shape.text_frame, data)
    return updated


def _replace_template_tokens_in_text_frame(text_frame, data: dict) -> int:
    old = "\n".join([p.text for p in text_frame.paragraphs])
    if not old or ("{" not in old and "[" not in old):
        return 0

    new = old
    token_values = []
    for key, value in (data or {}).items():
        if value is None:
            value = ""
        value = str(value)
        key = str(key)
        token_values.extend([
            (f"{{{{{key}}}}}", value),
            (f"{{{key}}}", value),
            ("{" + key + "}}", value),
            ("{{" + key + "}", value),
            (f"[{key}]", value),
        ])

    for token, value in sorted(token_values, key=lambda item: len(item[0]), reverse=True):
        if token in new:
            new = new.replace(token, value)

    if new == old:
        return 0

    set_multiline_text_keep_style(text_frame, new)
    return 1


def insert_main_photo(prs: Presentation, photo_url: str):
    if not photo_url:
        return
    try:
        resp = requests.get(photo_url, timeout=15)
        resp.raise_for_status()
    except Exception as e:
        logger.warning(f"사진 다운로드 실패: {e}")
        return
    img_bytes = BytesIO(resp.content)
    slide = prs.slides[TARGET_SLIDE_INDEX]
    try:
        picture_ph = slide.placeholders[13]
        picture_ph.insert_picture(img_bytes)
        logger.info("3번 슬라이드에 대표사진 삽입 완료")
    except Exception as e:
        logger.warning(f"사진 placeholder 삽입 실패: {e}")
        for shape in list(slide.shapes):
            if get_alt_text(shape) != "MAIN_PHOTO_BOX":
                continue
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception:
                pass
            try:
                slide.shapes.add_picture(img_bytes, left, top, width=width, height=height)
                logger.info("3번 슬라이드 MAIN_PHOTO_BOX에 대표사진 삽입 완료")
            except Exception as e2:
                logger.warning(f"MAIN_PHOTO_BOX 사진 삽입 실패: {e2}")
            break


def replace_placeholders_in_slide(slide, mapping: dict):
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        tf = shp.text_frame
        old = "\n".join([p.text for p in tf.paragraphs])
        new = old
        for k, v in mapping.items():
            new = new.replace(k, "" if v is None else str(v))
        if new != old:
            set_multiline_text_keep_style(tf, new)


def apply_property_status_opinion(prs: Presentation, opinion_text: str) -> bool:
    slide = _find_property_status_opinion_slide(prs)
    if slide is None:
        logger.warning("담당자 종합의견 (1) 물건현황 슬라이드를 찾지 못했습니다.")
        return False

    target = _find_main_body_text_shape(slide)
    if target is None or not getattr(target, "has_text_frame", False):
        logger.warning("담당자 종합의견 (1) 물건현황 본문 텍스트 박스를 찾지 못했습니다.")
        return False

    target.text_frame.word_wrap = True
    _set_property_status_text(target.text_frame, opinion_text or "")
    return True


def _set_property_status_text(text_frame, opinion_text: str, body_size_pt: float = 12) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    try:
        text_frame.margin_left = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.margin_top = Pt(1)
        text_frame.margin_bottom = Pt(1)
    except Exception:
        pass

    first = True
    for raw_line in str(opinion_text or "").splitlines():
        line = re.sub(r"\s+", " ", raw_line or "").strip()
        if not line:
            continue
        line = "- " + line.lstrip("-ㆍ• ").strip()
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        paragraph.level = 0
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(6)
        paragraph.line_spacing = 1.12
        _disable_paragraph_numbering(paragraph)
        _add_styled_run(paragraph, line, size_pt=body_size_pt, bold=False, color=OPINION_NAVY)

    if first:
        paragraph = text_frame.paragraphs[0]
        _disable_paragraph_numbering(paragraph)
        _add_styled_run(paragraph, "", size_pt=body_size_pt, bold=False, color=OPINION_NAVY)


def _set_rights_analysis_rich_text(
    text_frame,
    opinion_text: str,
    heading_size_pt: float = 9.8,
    body_size_pt: float = 9.3,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    try:
        text_frame.margin_left = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.margin_top = Pt(1)
        text_frame.margin_bottom = Pt(1)
    except Exception:
        pass

    first = True
    paragraph_count = 0
    for raw_line in str(opinion_text or "").splitlines():
        line = re.sub(r"\s+", " ", raw_line or "").strip()
        if not line:
            continue

        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        paragraph_count += 1
        paragraph.alignment = None
        paragraph.level = 0
        _disable_paragraph_numbering(paragraph)

        if re.match(r"^\d+\)", line):
            paragraph.space_before = Pt(7) if paragraph_count > 1 else Pt(0)
            paragraph.space_after = Pt(3)
            _add_styled_run(paragraph, line, size_pt=heading_size_pt, bold=True, color=OPINION_NAVY)
            continue

        paragraph.space_before = Pt(2)
        paragraph.space_after = Pt(2)
        paragraph.line_spacing = 1.08
        _add_highlighted_opinion_line(paragraph, line, body_size_pt=body_size_pt)

    if first:
        paragraph = text_frame.paragraphs[0]
        _disable_paragraph_numbering(paragraph)
        _add_styled_run(paragraph, "", size_pt=body_size_pt, bold=False, color=OPINION_NAVY)


def _add_highlighted_opinion_line(paragraph, line: str, body_size_pt: float = 9.3) -> None:
    text = _normalize_opinion_line(line)
    segments = _split_opinion_highlight_segments(text)
    for segment, emphasized in segments:
        _add_styled_run(
            paragraph,
            segment,
            size_pt=body_size_pt,
            bold=emphasized,
            color=OPINION_BLUE if emphasized else OPINION_NAVY,
        )


def _disable_paragraph_numbering(paragraph) -> None:
    try:
        p_pr = paragraph._p.get_or_add_pPr()
        for child in list(p_pr):
            if child.tag in {
                qn("a:buAutoNum"),
                qn("a:buChar"),
                qn("a:buBlip"),
                qn("a:buNone"),
            }:
                p_pr.remove(child)
        p_pr.append(OxmlElement("a:buNone"))
        for attr in ("marL", "indent"):
            if attr in p_pr.attrib:
                del p_pr.attrib[attr]
    except Exception:
        pass


def _normalize_opinion_line(line: str) -> str:
    line = str(line or "").strip()
    if line.startswith("-"):
        return "- " + line.lstrip("- ").strip()
    return line


def _split_opinion_highlight_segments(text: str) -> list[tuple[str, bool]]:
    highlight_phrases = (
        "등기부 상 낙찰자가 인수해야 하는 권리는 없습니다",
        "등기부상 낙찰자가 인수해야 하는 권리는 없습니다",
        "낙찰자가 인수해야 하는 권리는 없습니다",
        "낙찰자에게 인수되는 임차권리는 없습니다",
        "인수되는 임차권리는 없습니다",
        "인수되는 권리는 없습니다",
        "취하 가능성은 낮습니다",
        "취하가능성은 낮습니다",
        "취하 가능성 존재 합니다",
        "취하 가능성은 존재합니다",
        "무잉여 가능성은 없습니다",
        "무잉여가능성은 없습니다",
        "무잉여 가능성이 존재합니다",
    )
    matches: list[tuple[int, int]] = []
    for phrase in highlight_phrases:
        start = 0
        while True:
            idx = text.find(phrase, start)
            if idx < 0:
                break
            matches.append((idx, idx + len(phrase)))
            start = idx + len(phrase)

    if not matches:
        return [(text, False)]

    matches.sort(key=lambda span: (span[0], -(span[1] - span[0])))
    merged: list[tuple[int, int]] = []
    for start, end in matches:
        if merged and start <= merged[-1][1]:
            merged[-1] = (merged[-1][0], max(merged[-1][1], end))
        else:
            merged.append((start, end))

    segments: list[tuple[str, bool]] = []
    cursor = 0
    for start, end in merged:
        if cursor < start:
            segments.append((text[cursor:start], False))
        segments.append((text[start:end], True))
        cursor = end
    if cursor < len(text):
        segments.append((text[cursor:], False))
    return [(segment, emphasized) for segment, emphasized in segments if segment]


def _add_styled_run(paragraph, text: str, size_pt: float, bold: bool, color: RGBColor):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    try:
        run.font.color.rgb = color
    except Exception:
        pass
    _force_run_font(run)
    return run


def apply_rights_analysis_opinion(prs: Presentation, opinion_text: str) -> bool:
    slide = find_slide_by_note_key(prs, "SLIDE_KEY=OPINION_RIGHTS_ANALYSIS")
    if slide is None:
        slide = _find_opinion_slide_after_toc(prs, offset=2)
    if slide is None:
        slide = _find_slide_by_body_keywords(prs, ("말소기준", "임차권리", "경매취하"))
    if slide is None:
        logger.warning("담당자 종합의견 (2) 권리분석 슬라이드를 찾지 못했습니다.")
        return False

    target = _find_main_body_text_shape(slide)
    if target is None or not getattr(target, "has_text_frame", False):
        logger.warning("담당자 종합의견 (2) 권리분석 본문 텍스트 박스를 새로 생성합니다.")
        target = slide.shapes.add_textbox(735013, 1486429, 9534525, 5478251)

    target.text_frame.word_wrap = True
    _set_rights_analysis_rich_text(
        target.text_frame,
        opinion_text or "",
        heading_size_pt=RIGHTS_OPINION_HEADING_PT,
        body_size_pt=RIGHTS_OPINION_BODY_PT,
    )
    return True


def apply_special_opinion(prs: Presentation, opinion_text: str) -> bool:
    if not str(opinion_text or "").strip():
        return False

    slide = find_slide_by_note_key(prs, "SLIDE_KEY=OPINION_SPECIAL")
    if slide is None:
        slide = _find_opinion_slide_after_toc(prs, offset=3)
    if slide is None:
        logger.warning("담당자 종합의견 (3) 특이사항 슬라이드를 찾지 못했습니다.")
        return False

    target = _find_main_body_text_shape(slide)
    if target is None or not getattr(target, "has_text_frame", False):
        target = slide.shapes.add_textbox(735013, 1486429, 9534525, 5478251)

    target.text_frame.word_wrap = True
    visible_line_count = len([line for line in str(opinion_text or "").splitlines() if line.strip()])
    heading_size = 11 if visible_line_count > 16 else 14
    body_size = 8.5 if visible_line_count > 20 else (9.5 if visible_line_count > 14 else 14)
    _set_rights_analysis_rich_text(
        target.text_frame,
        opinion_text or "",
        heading_size_pt=heading_size,
        body_size_pt=body_size,
    )
    return True


def _find_property_status_opinion_slide(prs: Presentation):
    slide = find_slide_by_note_key(prs, "SLIDE_KEY=OPINION_PROPERTY_STATUS")
    if slide is None:
        slide = _find_opinion_slide_after_toc(prs, offset=1)
    if slide is not None:
        return slide

    slides = list(prs.slides)
    for slide in slides:
        try:
            notes = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes = ""
        if "텍스트 작성하세요" not in notes:
            continue
        for shape in slide.shapes:
            text = (getattr(shape, "text", "") or "").strip()
            if text.startswith("본건은") and "주위" in text:
                return slide
    return None


def _find_opinion_slide_after_toc(prs: Presentation, offset: int):
    slides = list(prs.slides)
    for idx, slide in enumerate(slides):
        try:
            notes = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes = ""
        if "SLIDE_KEY=OPINION_TOC" in notes or "물건현황|권리분석|특이사항" in notes:
            target_idx = idx + offset
            return slides[target_idx] if 0 <= target_idx < len(slides) else None
        for shape in slide.shapes:
            text = (getattr(shape, "text", "") or "").replace("\n", " ")
            compact = re.sub(r"\s+", "", text)
            if (
                "물건현황|권리분석|특이사항" in compact
                or all(token in compact for token in ("물건현황", "권리분석", "특이사항"))
            ):
                target_idx = idx + offset
                return slides[target_idx] if 0 <= target_idx < len(slides) else None
    return None


def _find_slide_by_body_keywords(prs: Presentation, keywords: tuple[str, ...]):
    for slide in prs.slides:
        combined = "\n".join(
            (getattr(shape, "text", "") or "")
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        if all(keyword in combined for keyword in keywords):
            return slide
    return None


def _find_main_body_text_shape(slide):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").strip()
        if re.fullmatch(r"\d+", text):
            continue
        area = int(shape.width) * int(shape.height)
        candidates.append((area, shape))
    if not candidates:
        return None
    candidates.sort(key=lambda item: item[0], reverse=True)
    return candidates[0][1]


def apply_eviction_cost_estimates(prs: Presentation, values: dict) -> int:
    """명도 정액제/실비제 비용 슬라이드에 강제집행 계산 결과를 반영한다."""
    attorney_fee = int(values.get("attorney_fee") or 0)
    normal_cost = int(values.get("normal_execution_cost") or 0)
    myungsung_cost = int(values.get("myungsung_execution_cost") or 0)
    flat_total = int(values.get("flat_total") or myungsung_cost)
    cost_plus_total = int(values.get("cost_plus_total") or (attorney_fee + normal_cost))
    updated = 0

    # 새 템플릿의 {{명도_...}} 변수는 Notes/AltText와 무관하게 저장 직전 한 번 더 보장 치환한다.
    for slide in prs.slides:
        updated += _replace_template_tokens_in_shapes(slide.shapes, values)

    for slide in prs.slides:
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            pass
        if not any(key in notes for key in ("EXCEL_TEXT_AUTO", "명도정액제", "변호사수임료")):
            continue

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            alt = get_alt_text(shape)
            current = (shape.text or "").strip()
            if not alt and not current:
                continue

            is_left_area = shape.left < (prs.slide_width / 2)
            new_text = ""

            if "All 총괄시트!F9" in alt:
                new_text = _format_flat_total(flat_total, current)
            elif "EXCEL_CALC:ACTUAL_COST_EST" in alt:
                new_text = _format_cost_plus_total(cost_plus_total, current)
            elif "All 총괄시트!F10" in alt:
                new_text = _format_won_with_space(attorney_fee)
            elif "강제집행 비용계산표!D11" in alt:
                new_text = _format_won_with_space(normal_cost)

            if new_text:
                set_text_keep_style(shape.text_frame, new_text)
                updated += 1

    return updated


def _format_won_with_space(value: int) -> str:
    return f"{int(value):,} 원"


def _format_manwon_number(value: int) -> str:
    return f"{round(int(value) / 10000):,}"


def _format_flat_total(value: int, current: str) -> str:
    text = f"총 {_format_manwon_number(value)}만원"
    return f"({text})" if current.startswith("(") else text


def _format_cost_plus_total(value: int, current: str) -> str:
    text = f"약 {_format_manwon_number(value)}만원 + @"
    if current.startswith("("):
        return f"({text}‥.)"
    return text


# ============================================================
# 이미지 삽입 (복제 + 노란박스 교체)
# ============================================================
def insert_images_into_ppt(prs, total_pages, keyword, img_pattern,
                           clone_from_next=False, force_title=False,
                           forced_num="(5)", base_as_last=False):
    if total_pages <= 0:
        return

    base_slide = find_slide_by_keyword(prs, keyword)
    if base_slide is None:
        logger.error(f"'{keyword}' 기준 슬라이드를 찾지 못했습니다.")
        return

    # 제목 텍스트 보존
    base_title_text = ""
    base_forced_num = forced_num
    for shp in base_slide.shapes:
        if not hasattr(shp, "text"):
            continue
        txt = (shp.text or "").strip()
        if txt and keyword in txt:
            base_title_text = re.sub(r"-\d+\s*$", "", txt).strip()
            m = re.search(r"\(\d+\)", base_title_text)
            if m:
                base_forced_num = m.group(0)
            break
    if not base_title_text:
        base_title_text = f"{base_forced_num} {keyword}"

    def _apply_title(dst):
        for shp in dst.shapes:
            if hasattr(shp, "text") and keyword in (shp.text or ""):
                if getattr(shp, "has_text_frame", False):
                    set_text_keep_style(shp.text_frame, base_title_text)
                else:
                    shp.text = base_title_text
                return
        for shp in dst.shapes:
            if get_alt_text(shp) != "SUBTITLE":
                continue
            if getattr(shp, "has_text_frame", False):
                set_text_keep_style(shp.text_frame, base_title_text)
            elif hasattr(shp, "text"):
                shp.text = base_title_text
            return

    # 템플릿 슬라이드 결정
    template_slide = base_slide
    base_index = prs.slides.index(base_slide)

    if clone_from_next:
        if keyword in ["건축물대장", "토지이용계획", "전자지도", "위성지도", "위치도"]:
            tpl_key = "SLIDE_KEY=TEMPLATE_OBJ_STATUS_PAGE"
        elif keyword in ["매각물건명세서", "등기사항 요약", "현황조사서", "임차인 현황", "등기부 현황"] or "등기부 현황" in keyword:
            tpl_key = "SLIDE_KEY=TEMPLATE_RIGHT_ANALYSIS_PAGE"
        else:
            tpl_key = ""

        if tpl_key:
            tpl_idx = find_slide_index_by_note_key(prs, tpl_key)
            if tpl_idx >= 0:
                template_slide = prs.slides[tpl_idx]

    # 슬라이드-페이지 매핑
    slide_for_page = {}
    if not base_as_last:
        slide_for_page[1] = base_slide
        for page in range(2, total_pages + 1):
            new_s = duplicate_slide(prs, template_slide)
            slide_for_page[page] = new_s
            _apply_title(new_s)
        for page in range(2, total_pages + 1):
            old_idx = prs.slides.index(slide_for_page[page])
            move_slide(prs, old_idx, base_index + (page - 1))
    else:
        for page in range(1, total_pages):
            new_s = duplicate_slide(prs, template_slide)
            slide_for_page[page] = new_s
            _apply_title(new_s)
        slide_for_page[total_pages] = base_slide
        for page in range(1, total_pages):
            old_idx = prs.slides.index(slide_for_page[page])
            move_slide(prs, old_idx, base_index + (page - 1))

    # 이미지 삽입
    for page in range(1, total_pages + 1):
        slide = slide_for_page.get(page)
        if slide is None:
            continue
        img_path = img_pattern.format(page=page)
        if not os.path.exists(img_path):
            continue
        yellow = find_yellow_box(slide)
        if yellow is None:
            continue
        left, top, width, height = yellow.left, yellow.top, yellow.width, yellow.height
        slide.shapes._spTree.remove(yellow._element)

        trimmed = img_path.replace(".png", "_trim.png")
        try:
            trim_white_margin(img_path, trimmed)
            track_file(trimmed)
            use_path = trimmed
        except Exception:
            use_path = img_path

        slide.shapes.add_picture(use_path, left, top, width=width, height=height)

        # 제목 번호
        if total_pages > 1:
            for shp in slide.shapes:
                if hasattr(shp, "text") and keyword in (shp.text or ""):
                    base = re.sub(r"-\d+\s*$", "", shp.text).strip()
                    if getattr(shp, "has_text_frame", False):
                        set_text_keep_style(shp.text_frame, f"{base}-{page}")
                    else:
                        shp.text = f"{base}-{page}"
                    break


def insert_single_image(prs, keyword_or_key, image_path, use_note_key=False):
    if not image_path or not os.path.exists(image_path):
        return
    slide = (find_slide_by_note_key(prs, keyword_or_key) if use_note_key
             else find_slide_by_keyword(prs, keyword_or_key))
    if slide is None:
        return
    yellow = find_yellow_box(slide)
    if yellow is None:
        return
    l, t, w, h = yellow.left, yellow.top, yellow.width, yellow.height
    slide.shapes._spTree.remove(yellow._element)

    trimmed = image_path.replace(".png", "_trim.png")
    try:
        trim_white_margin(image_path, trimmed)
        track_file(trimmed)
        use_path = trimmed
    except Exception:
        use_path = image_path
    slide.shapes.add_picture(use_path, l, t, width=w, height=h)


def insert_single_image_by_note_keywords(prs, keywords: list[str], image_path: str):
    if not image_path or not os.path.exists(image_path):
        return False
    slide = find_slide_by_note_keywords(prs, keywords)
    if slide is None:
        logger.warning(f"노트 키워드 슬라이드를 찾지 못했습니다: {keywords}")
        return False
    yellow = find_yellow_box(slide)
    if yellow is None:
        l, t = Inches(0.65), Inches(1.15)
        w, h = prs.slide_width - Inches(1.3), prs.slide_height - Inches(1.65)
    else:
        l, t, w, h = yellow.left, yellow.top, yellow.width, yellow.height
        slide.shapes._spTree.remove(yellow._element)

    trimmed = image_path.replace(".png", "_trim.png")
    try:
        trim_white_margin(image_path, trimmed)
        track_file(trimmed)
        use_path = trimmed
    except Exception:
        use_path = image_path
    try:
        with PILImage.open(use_path) as img:
            img_w, img_h = img.size
        img_ratio = img_w / img_h
        box_ratio = w / h
        if img_ratio > box_ratio:
            new_w = w
            new_h = int(w / img_ratio)
            new_l = l
            new_t = t + int((h - new_h) / 2)
        else:
            new_h = h
            new_w = int(h * img_ratio)
            new_l = l + int((w - new_w) / 2)
            new_t = t
        slide.shapes.add_picture(use_path, new_l, new_t, width=new_w, height=new_h)
    except Exception:
        slide.shapes.add_picture(use_path, l, t, width=w, height=h)
    return True


def insert_key_value_table_by_note_keywords(
    prs,
    keywords: list[str],
    title: str,
    rows: list[tuple[str, str]],
) -> bool:
    """옥션플래너 값만 전달된 경우 지정 슬라이드의 이미지 박스를 표로 대체한다."""
    if not rows:
        return False
    slide = find_slide_by_note_keywords(prs, keywords)
    if slide is None:
        logger.warning(f"노트 키워드 표 삽입 슬라이드를 찾지 못했습니다: {keywords}")
        return False
    yellow = find_yellow_box(slide)
    if yellow is None:
        left, top = Inches(0.65), Inches(1.15)
        width, height = prs.slide_width - Inches(1.3), prs.slide_height - Inches(1.65)
    else:
        left, top, width, height = yellow.left, yellow.top, yellow.width, yellow.height
        slide.shapes._spTree.remove(yellow._element)
    visible_rows = rows[:16]
    table = slide.shapes.add_table(len(visible_rows) + 1, 2, left, top, width, height).table
    table.columns[0].width = int(width * 0.38)
    table.columns[1].width = width - table.columns[0].width

    header = table.cell(0, 0)
    header.merge(table.cell(0, 1))
    header.text = title
    header.fill.solid()
    header.fill.fore_color.rgb = OPINION_NAVY
    for paragraph in header.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            _force_run_font(run)

    for row_index, (label, value) in enumerate(visible_rows, start=1):
        for column_index, text in enumerate((label, value)):
            cell = table.cell(row_index, column_index)
            cell.text = str(text or "-")
            cell.margin_left = 70000
            cell.margin_right = 70000
            if column_index == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 240, 248)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.RIGHT
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.bold = column_index == 0
                    run.font.color.rgb = OPINION_NAVY
                    _force_run_font(run)
    return True


def insert_internal_structure_image(prs, image_path: str) -> bool:
    if not image_path or not os.path.exists(image_path):
        return False
    target_slide = find_slide_by_note_key(prs, "SLIDE_KEY=OBJ_LOCATION")
    target_keywords = ("내부구조도/호별배치도", "내부구조도", "호별배치도")
    if target_slide is None:
        for slide in prs.slides:
            slide_text = " ".join(
                (getattr(shape, "text", "") or "").replace("\n", " ").strip()
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            compact = re.sub(r"\s+", "", slide_text)
            if any(keyword in compact for keyword in target_keywords) and find_image_boxes_left_to_right(slide, limit=1):
                target_slide = slide
                break
    if target_slide is None:
        logger.warning("내부구조도/호별배치도 대상 슬라이드를 찾지 못했습니다.")
        return False

    yellow = find_yellow_box(target_slide)
    if yellow is not None:
        l, t, w, h = yellow.left, yellow.top, yellow.width, yellow.height
        target_slide.shapes._spTree.remove(yellow._element)
    else:
        boxes = find_image_boxes_left_to_right(target_slide, limit=2)
        if not boxes:
            logger.warning("내부구조도/호별배치도 슬라이드에서 이미지 삽입 박스를 찾지 못했습니다.")
            return False
        box = boxes[0]
        l, t, w, h = box.left, box.top, box.width, box.height
        target_slide.shapes._spTree.remove(box._element)

    trimmed = image_path.replace(".png", "_trim.png")
    try:
        trim_white_margin(image_path, trimmed)
        track_file(trimmed)
        use_path = trimmed
    except Exception:
        use_path = image_path

    try:
        with PILImage.open(use_path) as img:
            img_w, img_h = img.size
        img_ratio = img_w / img_h
        box_ratio = w / h
        if img_ratio > box_ratio:
            new_w = w
            new_h = int(w / img_ratio)
            new_l = l
            new_t = t + int((h - new_h) / 2)
        else:
            new_h = h
            new_w = int(h * img_ratio)
            new_l = l + int((w - new_w) / 2)
            new_t = t
        target_slide.shapes.add_picture(use_path, new_l, new_t, width=new_w, height=new_h)
    except Exception:
        target_slide.shapes.add_picture(use_path, l, t, width=w, height=h)
    return True


def insert_location_and_structure_images(prs, location_img: str = "", structure_img: str = "") -> bool:
    """물건현황 (1) 슬라이드에 좌측 위치도, 우측 내부구조도/호별배치도를 삽입한다."""
    if (not location_img or not os.path.exists(location_img)) and (not structure_img or not os.path.exists(structure_img)):
        return False

    target_slide = find_slide_by_note_key(prs, "SLIDE_KEY=OBJ_LOCATION")
    if target_slide is None:
        target_keywords = ("내부구조도/호별배치도", "내부구조도", "호별배치도", "위치도")
        for slide in prs.slides:
            slide_text = " ".join(
                (getattr(shape, "text", "") or "").replace("\n", " ").strip()
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            compact = re.sub(r"\s+", "", slide_text)
            if all(keyword in compact for keyword in ("위치도", "내부구조도")) or any(keyword in compact for keyword in target_keywords):
                target_slide = slide
                break
    if target_slide is None:
        logger.warning("위치도/내부구조도 공용 대상 슬라이드를 찾지 못했습니다.")
        return False

    boxes = find_yellow_boxes_left_to_right(target_slide, limit=2)
    if len(boxes) < 2:
        boxes = find_image_boxes_left_to_right(target_slide, limit=2)
    if len(boxes) < 2:
        logger.warning("위치도/내부구조도 공용 슬라이드에서 좌우 이미지 박스 2개를 찾지 못했습니다.")
        return False

    def _put(box, img_path, label):
        if not img_path or not os.path.exists(img_path):
            logger.warning(f"{label} 이미지가 없어 삽입을 생략합니다.")
            return False
        l, t, w, h = box.left, box.top, box.width, box.height
        try:
            target_slide.shapes._spTree.remove(box._element)
        except Exception:
            pass
        trimmed = img_path.replace(".png", "_trim.png")
        try:
            trim_white_margin(img_path, trimmed)
            track_file(trimmed)
            use_path = trimmed
        except Exception:
            use_path = img_path
        target_slide.shapes.add_picture(use_path, l, t, width=w, height=h)
        logger.info(f"{label} 이미지 삽입 완료: {use_path}")
        return True

    inserted = False
    inserted = _put(boxes[0], location_img, "좌측 위치도") or inserted
    inserted = _put(boxes[1], structure_img, "우측 내부구조도/호별배치도") or inserted
    return inserted


def insert_two_images_location(prs, keyword, left_img, right_img=""):
    slide = find_slide_by_keyword(prs, keyword)
    if slide is None:
        logger.warning(f"'{keyword}' 위치도 삽입 대상 슬라이드를 찾지 못했습니다.")
        return
    boxes = find_yellow_boxes_left_to_right(slide, limit=2)
    if not boxes:
        boxes = find_image_boxes_left_to_right(slide, limit=2)
    if not boxes:
        logger.warning(f"'{keyword}' 위치도 삽입 박스를 찾지 못했습니다.")
        return

    slide_text = " ".join(
        (getattr(shape, "text", "") or "").replace("\n", " ").strip()
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    compact = re.sub(r"\s+", "", slide_text)
    if "내부구조도" in compact and len(boxes) > 1:
        chosen_img = right_img if right_img and os.path.exists(right_img) else left_img
        logger.info("위치도 슬라이드가 내부구조도와 공용이므로 오른쪽 박스에 위치도 1장을 삽입합니다.")
        boxes = [boxes[1]]
        left_img = chosen_img
        right_img = ""

    def _put(box, img_path):
        if not img_path or not os.path.exists(img_path):
            return
        l, t, w, h = box.left, box.top, box.width, box.height
        slide.shapes._spTree.remove(box._element)
        trimmed = img_path.replace(".png", "_trim.png")
        try:
            trim_white_margin(img_path, trimmed)
            track_file(trimmed)
            use_path = trimmed
        except Exception:
            use_path = img_path
        slide.shapes.add_picture(use_path, l, t, width=w, height=h)

    _put(boxes[0], left_img)
    if len(boxes) > 1:
        _put(boxes[1], right_img)


# ============================================================
# VBA 보존 저장
# ============================================================
def save_pptm_preserve_vba(template_pptm: str, prs: Presentation, out_pptm: str):
    tmp_pptx = tempfile.mktemp(suffix=".pptx")
    prs.save(tmp_pptx)

    vba_bin = None
    try:
        with zipfile.ZipFile(template_pptm, "r") as zt:
            try:
                vba_bin = zt.read("ppt/vbaProject.bin")
            except KeyError:
                vba_bin = None
    except Exception:
        vba_bin = None

    if not vba_bin:
        try:
            shutil.move(tmp_pptx, out_pptm)
        finally:
            if os.path.exists(tmp_pptx):
                try:
                    os.remove(tmp_pptx)
                except Exception:
                    pass
        return

    with zipfile.ZipFile(tmp_pptx, "r") as zi, \
         zipfile.ZipFile(out_pptm, "w", compression=zipfile.ZIP_DEFLATED) as zo:
        for item in zi.infolist():
            if item.filename == "ppt/vbaProject.bin":
                continue
            zo.writestr(item, zi.read(item.filename))
        zo.writestr("ppt/vbaProject.bin", vba_bin)

    try:
        os.remove(tmp_pptx)
    except Exception:
        pass
